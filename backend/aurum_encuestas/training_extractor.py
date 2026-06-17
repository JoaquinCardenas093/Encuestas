"""Training extractor: parse training PPTs, learn layouts + chart visual styles."""
from collections import defaultdict
from datetime import UTC, datetime
from pathlib import Path

from lxml import etree
from pptx import Presentation

from .models import LayoutBank, LayoutElement, LearnedLayout


NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"


def signature_for_slide(n_charts: int, chart_types: list[str], n_chart_an: int, n_q_an: int, has_slide_an: bool) -> str:
    types = ",".join(sorted(chart_types))
    return f"{n_charts}|{types}|{n_chart_an}|{n_q_an}|{1 if has_slide_an else 0}"


def signature_loose_count(n_charts: int) -> str:
    """Looser signature: just chart count. Used for fallback matching."""
    return f"count:{n_charts}"


def extract_layouts_from_pptx(pptx_path: str) -> list[LearnedLayout]:
    prs = Presentation(pptx_path)
    layouts: list[LearnedLayout] = []

    for slide_idx, slide in enumerate(prs.slides):
        charts = [sh for sh in slide.shapes if getattr(sh, "has_chart", False)]
        text_boxes = [sh for sh in slide.shapes if sh.has_text_frame and not getattr(sh, "has_chart", False)]
        if not charts:
            continue

        chart_types = []
        chart_els: list[LayoutElement] = []
        per_chart_style = {}
        for i, ch in enumerate(charts):
            ct = _xl_to_app(ch.chart.chart_type) if hasattr(ch.chart, "chart_type") else "BAR"
            chart_types.append(ct)
            chart_els.append(LayoutElement(
                role=f"chart_{i}",
                x=ch.left or 0, y=ch.top or 0, cx=ch.width or 0, cy=ch.height or 0,
                chart_type=ct,
            ))
            per_chart_style[f"chart_{i}"] = _extract_chart_style(ch.chart)

        text_els: list[LayoutElement] = []
        text_styles: dict[str, dict] = {}
        for ti, tb in enumerate(text_boxes):
            text = tb.text_frame.text or ""
            if "@Titulo" in text or "@Notas" in text or "@" in text:
                continue
            role = "slide_analysis" if _is_bottom(tb, prs.slide_height) else f"chart_analysis_{ti}"
            text_els.append(LayoutElement(
                role=role,
                x=tb.left or 0, y=tb.top or 0, cx=tb.width or 0, cy=tb.height or 0,
            ))
            text_styles[role] = _extract_textbox_style(tb)

        n_chart_an = sum(1 for e in text_els if e.role.startswith("chart_analysis"))
        n_slide_an = sum(1 for e in text_els if e.role == "slide_analysis")
        signature = signature_for_slide(len(charts), chart_types, n_chart_an, 0, n_slide_an > 0)

        free_area = {
            "x": min(e.x for e in chart_els),
            "y": min(e.y for e in chart_els),
            "cx": prs.slide_width,
            "cy": prs.slide_height,
        }

        layouts.append(LearnedLayout(
            id=f"lay_{Path(pptx_path).stem}_{slide_idx:03d}",
            signature=signature,
            source=f"{Path(pptx_path).name}#slide{slide_idx + 1}",
            free_area=free_area,
            elements=chart_els + text_els,
            chart_style=per_chart_style,
            text_style=text_styles,
        ))

    return layouts


def _is_bottom(shape, slide_height: int) -> bool:
    top = shape.top or 0
    return top > slide_height * 0.7


_SCHEME_FALLBACK = {
    "accent1": "5B9BD5", "accent2": "ED7D31", "accent3": "A5A5A5",
    "accent4": "FFC000", "accent5": "4472C4", "accent6": "70AD47",
    "dk1": "000000", "lt1": "FFFFFF", "dk2": "44546A", "lt2": "E7E6E6",
}


def _resolve_color(color_el) -> str | None:
    """Resolve any color element (srgbClr, sysClr, schemeClr) to a hex RGB string.
    Applies lumMod/lumOff modifiers from children."""
    tag = color_el.tag.split("}")[-1]
    base = None
    if tag == "srgbClr":
        base = color_el.get("val")
    elif tag == "sysClr":
        base = color_el.get("lastClr") or color_el.get("val")
        # sysClr "window" → white, "windowText" → black, etc.
        sys_map = {"window": "FFFFFF", "windowText": "000000"}
        if base is None and color_el.get("val") in sys_map:
            base = sys_map[color_el.get("val")]
    elif tag == "schemeClr":
        scheme = color_el.get("val")
        base = _SCHEME_FALLBACK.get(scheme)
    if not base:
        return None
    base = base.upper().lstrip("#")
    if len(base) != 6:
        return None

    # Apply lumMod / lumOff
    try:
        rgb = [int(base[i:i+2], 16) for i in (0, 2, 4)]
        for child in color_el:
            ctag = child.tag.split("}")[-1]
            v = child.get("val")
            if v is None:
                continue
            try:
                pct = int(v) / 100000.0
            except ValueError:
                continue
            if ctag == "lumMod":
                rgb = [int(c * pct) for c in rgb]
            elif ctag == "lumOff":
                rgb = [int(c + (255 - c) * pct) for c in rgb]
            elif ctag == "shade":
                rgb = [int(c * pct) for c in rgb]
            elif ctag == "tint":
                rgb = [int(c + (255 - c) * pct) for c in rgb]
        rgb = [max(0, min(255, c)) for c in rgb]
        return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    except Exception:
        return base


def _extract_chart_style(chart) -> dict:
    """Extract visual properties from a chart: series/dPt colors only, data labels, legend, font.

    Returns: {colors: [hex...], has_data_labels: bool, show_percent: bool, show_value: bool,
              show_cat_name: bool, legend: "r"|"b"|None, font: str|None, font_size: int|None}
    """
    style = {
        "colors": [],
        "has_data_labels": False,
        "show_percent": False,
        "show_value": False,
        "show_cat_name": False,
        "legend": None,
        "font": None,
        "font_size": None,
    }
    try:
        chart_xml = chart._chartSpace

        # Color extraction from series/dPt fills, supports srgbClr, sysClr (with lastClr/lumMod),
        # schemeClr (with fallback palette + lumMod/tint/shade modifiers).
        color_tags = {f"{{{NS_A}}}srgbClr", f"{{{NS_A}}}sysClr", f"{{{NS_A}}}schemeClr"}
        for ser in chart_xml.iter(f"{{{NS_C}}}ser"):
            # series-level fill (used for BAR, COLUMN, LINE, etc.)
            for spPr in ser.findall(f"{{{NS_C}}}spPr"):
                for solidFill in spPr.iter(f"{{{NS_A}}}solidFill"):
                    for child in solidFill:
                        if child.tag in color_tags:
                            hex_color = _resolve_color(child)
                            if hex_color and hex_color not in style["colors"]:
                                style["colors"].append(hex_color)
                            break
                    break
            # data point fills (used for PIE/DONUT — one dPt per slice)
            for dpt in ser.findall(f"{{{NS_C}}}dPt"):
                spPr = dpt.find(f"{{{NS_C}}}spPr")
                if spPr is None:
                    continue
                for solidFill in spPr.iter(f"{{{NS_A}}}solidFill"):
                    for child in solidFill:
                        if child.tag in color_tags:
                            hex_color = _resolve_color(child)
                            if hex_color and hex_color not in style["colors"]:
                                style["colors"].append(hex_color)
                            break
                    break

        # Legend position
        legend = chart_xml.find(f".//{{{NS_C}}}legend")
        if legend is not None:
            pos_el = legend.find(f"{{{NS_C}}}legendPos")
            style["legend"] = pos_el.get("val") if pos_el is not None else "r"

        # Data labels — chart-level dLbls
        dlbls = chart_xml.find(f".//{{{NS_C}}}plotArea//{{{NS_C}}}dLbls")
        if dlbls is not None:
            def _flag(name):
                el = dlbls.find(f"{{{NS_C}}}{name}")
                return el is not None and el.get("val", "1") != "0"
            style["show_value"] = _flag("showVal")
            style["show_percent"] = _flag("showPercent")
            style["show_cat_name"] = _flag("showCatName")
            style["has_data_labels"] = style["show_value"] or style["show_percent"] or style["show_cat_name"]

        # Font from chart's txPr
        latin = chart_xml.find(f".//{{{NS_C}}}txPr//{{{NS_A}}}latin")
        if latin is not None:
            style["font"] = latin.get("typeface")
        rpr = chart_xml.find(f".//{{{NS_C}}}txPr//{{{NS_A}}}defRPr")
        if rpr is not None and rpr.get("sz"):
            style["font_size"] = int(rpr.get("sz")) // 100
    except Exception:
        pass
    return style


def _extract_textbox_style(tb) -> dict:
    """Extract font, size, color, fill from a textbox shape."""
    style = {"font": None, "font_size": None, "color": None, "fill": None}
    try:
        for para in tb.text_frame.paragraphs:
            for run in para.runs:
                if run.font:
                    if run.font.name and not style["font"]:
                        style["font"] = run.font.name
                    if run.font.size and not style["font_size"]:
                        style["font_size"] = int(run.font.size.pt)
                    try:
                        if run.font.color and run.font.color.rgb and not style["color"]:
                            style["color"] = str(run.font.color.rgb)
                    except Exception:
                        pass
                if style["font"] and style["font_size"] and style["color"]:
                    break
    except Exception:
        pass
    return style


_REVERSE_CHART_MAP = {
    5: "PIE",
    -4120: "DONUT",
    57: "BAR",
    51: "COLUMN",
    58: "BAR_STACKED",
    52: "COLUMN_STACKED",
    4: "LINE",
    1: "AREA",
    -4151: "RADAR",
}


def _xl_to_app(xl_type) -> str:
    try:
        v = int(xl_type)
    except (TypeError, ValueError):
        return "BAR"
    return _REVERSE_CHART_MAP.get(v, "BAR")


def build_bank_from_pptxs(pptx_paths: list[str]) -> LayoutBank:
    all_layouts: list[LearnedLayout] = []
    for p in pptx_paths:
        try:
            all_layouts.extend(extract_layouts_from_pptx(p))
        except Exception:
            continue
    return LayoutBank(
        extracted_at=datetime.now(UTC).isoformat(),
        source_pptxs=[Path(p).name for p in pptx_paths],
        layouts=all_layouts,
    )


def aggregate_chart_style_by_type(bank: LayoutBank) -> dict[str, dict]:
    """Merge chart_style across all layouts per chart_type. Returns: {chart_type: merged_style}.
    Uses most common values: union colors, mode for booleans, first non-null for fonts."""
    by_type: dict[str, list[dict]] = defaultdict(list)
    for lay in bank.layouts:
        for role, style in (lay.chart_style or {}).items():
            ct = None
            for el in lay.elements:
                if el.role == role and el.chart_type:
                    ct = el.chart_type
                    break
            if ct:
                by_type[ct].append(style)

    out: dict[str, dict] = {}
    for ct, styles in by_type.items():
        colors_ordered: list[str] = []
        for s in styles:
            for c in s.get("colors", []):
                if c not in colors_ordered:
                    colors_ordered.append(c)
        font = next((s.get("font") for s in styles if s.get("font")), None)
        font_size = next((s.get("font_size") for s in styles if s.get("font_size")), None)
        legends = [s.get("legend") for s in styles if s.get("legend")]
        legend = max(set(legends), key=legends.count) if legends else None
        labels_votes = sum(1 for s in styles if s.get("has_data_labels"))
        has_data_labels = labels_votes >= (len(styles) / 2)
        out[ct] = {
            "colors": colors_ordered[:12],
            "font": font,
            "font_size": font_size,
            "legend": legend,
            "has_data_labels": has_data_labels,
        }
    return out
