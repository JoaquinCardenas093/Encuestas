"""Chart element renderer — dispatches all 9 chart types via python-pptx."""
from __future__ import annotations

import logging
from typing import TYPE_CHECKING

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.util import Emu, Pt

if TYPE_CHECKING:
    from .render_context import RenderContext

log = logging.getLogger(__name__)

# Mapping from spec chart_type string to python-pptx XL_CHART_TYPE
_CHART_TYPE_MAP: dict[str, int] = {
    "PIE":               XL_CHART_TYPE.PIE,
    "DONUT":             XL_CHART_TYPE.DOUGHNUT,
    "BAR_HORIZONTAL":    XL_CHART_TYPE.BAR_CLUSTERED,
    "BAR_CLUSTERED":     XL_CHART_TYPE.BAR_CLUSTERED,
    "BAR_STACKED":       XL_CHART_TYPE.BAR_STACKED,
    "COLUMN_CLUSTERED":  XL_CHART_TYPE.COLUMN_CLUSTERED,
    "COLUMN_STACKED":    XL_CHART_TYPE.COLUMN_STACKED,
    "LINE":              XL_CHART_TYPE.LINE,
    "AREA":              XL_CHART_TYPE.AREA,
    # Fase A defensive: grouped render is Fase B. UI doesn't expose these
    # types yet, but a hand-crafted JSON could carry them. Fall back to
    # single-mode shape so the slide doesn't crash.
    "PIE_GROUPED":             XL_CHART_TYPE.PIE,
    "BAR_HORIZONTAL_GROUPED":  XL_CHART_TYPE.BAR_CLUSTERED,
}

_LABEL_POSITION_MAP = {
    "inside":       XL_LABEL_POSITION.INSIDE_END,
    "outside_end":  XL_LABEL_POSITION.OUTSIDE_END,
    "center":       XL_LABEL_POSITION.CENTER,
    "best_fit":     XL_LABEL_POSITION.BEST_FIT,
}

_LEGEND_POSITION_MAP = {
    "right":  XL_LEGEND_POSITION.RIGHT,
    "bottom": XL_LEGEND_POSITION.BOTTOM,
    "top":    XL_LEGEND_POSITION.TOP,
    "left":   XL_LEGEND_POSITION.LEFT,
}


def _compute_grid_dims(n: int, grid_cols: int | None) -> tuple[int, int]:
    """Return (rows, cols) for PIE_GROUPED grid.

    User-set grid_cols overrides auto rule (rows=1 if N<=3 else 2 if N<=6 else 3).
    """
    if grid_cols and grid_cols >= 1:
        cols = grid_cols
        rows = (n + cols - 1) // cols
        return rows, cols
    rows = 1 if n <= 3 else (2 if n <= 6 else 3)
    cols = (n + rows - 1) // rows
    return rows, cols


def _add_title_textbox(slide, x: int, y: int, w: int, h: int, text: str, ctx) -> None:
    """Centered bold title text-box above a PIE_GROUPED grid."""
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(ctx.typography.get("title_size", 16))
    run.font.bold = True
    run.font.name = ctx.typography.get("font_family", "Calibri")


def _render_pie_grouped(slide, element: dict, source_chart, ctx) -> None:
    """Render PIE_GROUPED as N pie chart shapes inside a grid bbox.

    Each cat of source_chart.breakdown_ids[0] becomes one pie. Per-pie title
    uses cat_titles[cat] or cat label. Optional chart.title renders as a
    text-box above the grid.
    """
    x, y, cx, cy = _resolve_position(element.get("position", {}), ctx)

    bds = list(getattr(source_chart, "breakdown_ids", []) or [])
    primary_bd = bds[0] if bds else None
    all_bds = getattr(source_chart, "all_breakdowns_data", {}) or {}
    bd_data = all_bds.get(primary_bd, {}) if primary_bd else {}
    categories = bd_data.get("categories", {}) or {}

    if not categories:
        log.warning("PIE_GROUPED with empty breakdown — skipping render")
        return

    cat_list = list(categories.keys())
    n = len(cat_list)
    rows, cols = _compute_grid_dims(n, getattr(source_chart, "grid_cols", None))
    cat_titles = getattr(source_chart, "cat_titles", None) or {}
    question = getattr(source_chart, "question", None)
    options = list(question.options) if question else []

    sc_title = (getattr(source_chart, "title", None) or "").strip()
    el_title = (element.get("title") or "").strip()
    title_str = sc_title or el_title
    title_band_h = 400_000 if title_str else 0
    if title_str:
        _add_title_textbox(slide, x, y, cx, title_band_h, title_str, ctx)

    grid_y = y + title_band_h
    grid_h = cy - title_band_h
    gap_x = int(0.02 * cx)
    gap_y = int(0.06 * grid_h)
    cell_w = (cx - gap_x * (cols - 1)) // cols
    cell_h = (grid_h - gap_y * (rows - 1)) // rows

    fallback_colors = list(ctx.chart_colors or [])
    per_chart_colors = list(getattr(source_chart, "colors", []) or [])

    for i, cat in enumerate(cat_list):
        r, c = divmod(i, cols)
        cell_x = x + c * (cell_w + gap_x)
        cell_y = grid_y + r * (cell_h + gap_y)

        cd = CategoryChartData()
        cd.categories = options
        values = [float((categories[cat].get(opt) or {}).get("pct", 0) or 0) for opt in options]
        cd.add_series("", values)

        try:
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                Emu(cell_x), Emu(cell_y),
                Emu(cell_w), Emu(cell_h),
                cd,
            )
        except Exception as exc:
            log.error("PIE_GROUPED: add_chart failed for cat %r: %s", cat, exc)
            continue
        sub_chart = chart_shape.chart
        sub_chart.has_legend = False

        # Per-slice colors
        n_pts = len(options)
        effective_colors: list[str] = []
        for j in range(max(n_pts, len(per_chart_colors), len(fallback_colors))):
            if j < len(per_chart_colors) and per_chart_colors[j]:
                effective_colors.append(per_chart_colors[j])
            elif fallback_colors:
                effective_colors.append(fallback_colors[j % len(fallback_colors)])
            else:
                effective_colors.append("#7F7F7F")
        _apply_series_colors(sub_chart, effective_colors)

        # Labels: category + percentage inside slice
        _apply_labels(sub_chart, {
            "show_category_name": True, "show_percentage": True,
            "show_value": False, "format": "0.0%",
            "position": "outside_end",
        }, ctx)

        # Per-pie title via cat_titles override
        sub_title = cat_titles.get(cat) or cat
        sub_chart.has_title = True
        sub_chart.chart_title.text_frame.text = sub_title

        # Pie rotation (existing rule)
        sorted_vals = sorted(values, reverse=True)
        if sorted_vals:
            total = sum(v for v in sorted_vals if v) or 1.0
            dom = (sorted_vals[0] or 0) / total
            angle = 180 if abs(dom - 0.5) < 0.05 else int(round(-90 - dom * 180)) % 360
            _set_pie_first_slice_angle(sub_chart, angle)


def render(slide, element: dict, ctx: RenderContext) -> None:
    """Render a chart element onto slide in-place."""
    # Resolve data source first so we can read the source chart's user-selected
    # chart_type and override the pattern's hardcoded chart_type. This lets
    # users change a slide's main chart from PIE → COLUMN_CLUSTERED via the UI
    # without editing the pattern.
    data_source = element.get("data_source", {})
    chart_ref_index = data_source.get("chart_ref_index", 0)
    value_field = data_source.get("value_field", "pct")

    charts_list = getattr(ctx.slide_config, "charts", []) or []
    if chart_ref_index >= len(charts_list):
        log.warning("chart_ref_index %d out of range (have %d charts) — skipping", chart_ref_index, len(charts_list))
        return

    source_chart = charts_list[chart_ref_index]
    # UI selection wins: source_chart.chart_type comes from AddChartModal and
    # is the authoritative type. The pattern's chart_type is now layout-only
    # advice — fall back to it only if the UI didn't pick anything (legacy charts).
    ui_chart_type = (getattr(source_chart, "chart_type", None) or "").strip()
    pattern_chart_type = (element.get("chart_type") or "").strip()
    chart_type_str = ui_chart_type or pattern_chart_type or "BAR_HORIZONTAL"

    if chart_type_str == "PIE_GROUPED":
        _render_pie_grouped(slide, element, source_chart, ctx)
        return

    xl_chart_type = _CHART_TYPE_MAP.get(chart_type_str)
    if xl_chart_type is None:
        log.warning("Unknown chart_type %r — falling back to BAR_CLUSTERED", chart_type_str)
        xl_chart_type = XL_CHART_TYPE.BAR_CLUSTERED
    is_pie = chart_type_str in ("PIE", "DONUT")

    # Build CategoryChartData (sorted in-place if requested)
    sort_order = element.get("sort", "none")
    chart_data, sorted_values = _build_chart_data(source_chart, value_field, sort_order)

    # Resolve position
    x, y, cx, cy = _resolve_position(element.get("position", {}), ctx)

    # Add chart shape
    try:
        chart_shape = slide.shapes.add_chart(xl_chart_type, Emu(x), Emu(y), Emu(cx), Emu(cy), chart_data)
    except Exception as exc:
        log.error("Failed to add chart shape: %s", exc)
        return

    chart = chart_shape.chart

    # Compose per-slice colors: source_chart.colors[i] wins over ctx global cascade.
    per_chart_colors = list(getattr(source_chart, "colors", []) or [])
    n_pts = len(sorted_values)
    effective_colors: list[str] = []
    fallback = ctx.chart_colors or []
    for i in range(max(n_pts, len(per_chart_colors), len(fallback))):
        if i < len(per_chart_colors) and per_chart_colors[i]:
            effective_colors.append(per_chart_colors[i])
        elif fallback:
            effective_colors.append(fallback[i % len(fallback)])
        else:
            effective_colors.append("#7F7F7F")
    _apply_series_colors(chart, effective_colors)

    # Apply labels — pie/donut force show_category_name + show_percentage, NO value.
    labels_cfg = dict(element.get("labels", {}))
    if is_pie:
        labels_cfg.setdefault("show_category_name", True)
        labels_cfg.setdefault("show_percentage", True)
        labels_cfg["show_value"] = False  # never raw fraction
        labels_cfg.setdefault("format", "0.0%")
    _apply_labels(chart, labels_cfg, ctx)

    # Pie/donut rotation rules:
    #   ~50/50  (|dom-0.5|<0.05) → firstSliceAng=180  vertical split, first
    #                              option on the LEFT half (matches Aurora ref).
    #   else                     → small slice centered ~45° (top-right). The
    #                              dominant slice covers the rest, leading the
    #                              eye to the right side of the disc.
    if is_pie and sorted_values:
        total = sum(v for v in sorted_values if v) or 1.0
        dom_frac = (sorted_values[0] or 0) / total
        if abs(dom_frac - 0.5) < 0.05:
            angle = 180
        else:
            angle = int(round(-90 - dom_frac * 180)) % 360
        _set_pie_first_slice_angle(chart, angle)

    # Chart title (all single-shape types)
    sc_title = (getattr(source_chart, "title", None) or "").strip()
    el_title = (element.get("title") or "").strip()
    title_str = sc_title or el_title
    if title_str:
        chart.has_title = True
        chart.chart_title.text_frame.text = title_str
    else:
        chart.has_title = False

    # Legend (Fase B):
    show_legend = bool(getattr(source_chart, "show_legend", False))
    if is_pie:
        chart.has_legend = False
    elif chart_type_str == "BAR_HORIZONTAL_GROUPED" and show_legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False


def render_chart_at(
    slide,
    source_chart,
    x: int, y: int, cx: int, cy: int,
    ctx: RenderContext,
    chart_type: str = "BAR_HORIZONTAL",
    value_field: str = "pct",
    sort: str = "desc_by_value",
    show_legend: bool = False,
    show_category_name: bool = True,
    show_percentage: bool = True,
    number_format: str = "0.0%",
    label_position: str = "outside_end",
) -> None:
    """Render a chart at absolute EMU coords, bypassing pattern position logic.

    Used by table_renderer to embed a mini-chart inside a breakdown panel
    without needing a synthetic element/position dict.
    """
    xl_chart_type = _CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.BAR_CLUSTERED)
    is_pie = chart_type in ("PIE", "DONUT")
    chart_data, sorted_values = _build_chart_data(source_chart, value_field, sort)

    try:
        chart_shape = slide.shapes.add_chart(
            xl_chart_type, Emu(x), Emu(y), Emu(cx), Emu(cy), chart_data,
        )
    except Exception as exc:
        log.error("render_chart_at: add_chart failed: %s", exc)
        return
    chart = chart_shape.chart

    per_chart_colors = list(getattr(source_chart, "colors", []) or [])
    fallback = ctx.chart_colors or []
    n_pts = len(sorted_values)
    effective_colors: list[str] = []
    for i in range(max(n_pts, len(per_chart_colors), len(fallback))):
        if i < len(per_chart_colors) and per_chart_colors[i]:
            effective_colors.append(per_chart_colors[i])
        elif fallback:
            effective_colors.append(fallback[i % len(fallback)])
        else:
            effective_colors.append("#7F7F7F")
    _apply_series_colors(chart, effective_colors)

    labels_cfg = {
        "show_category_name": show_category_name and not is_pie,  # for pies category appears in slice label
        "show_percentage": show_percentage,
        "show_value": False,
        "format": number_format,
        "position": label_position,
        "font_size": 8,
    }
    if is_pie:
        labels_cfg["show_category_name"] = True
    _apply_labels(chart, labels_cfg, ctx)

    if is_pie and sorted_values:
        total = sum(v for v in sorted_values if v) or 1.0
        dom_frac = (sorted_values[0] or 0) / total
        if abs(dom_frac - 0.5) < 0.05:
            angle = 180
        else:
            angle = int(round(-90 - dom_frac * 180)) % 360
        _set_pie_first_slice_angle(chart, angle)

    if is_pie or not show_legend:
        chart.has_legend = False
    else:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    chart.has_title = False


def _set_pie_first_slice_angle(chart, angle_deg: int) -> None:
    """Set <c:firstSliceAng val="..."/> on the pie/doughnut plot.

    python-pptx does not expose firstSliceAng directly, so we manipulate the
    underlying XML. Angle 0 = 12 o'clock; values increase clockwise.
    """
    try:
        from pptx.oxml.ns import qn

        plot = chart.plots[0]
        plot_el = plot._element  # <c:pieChart> or <c:doughnutChart>
        ang_el = plot_el.find(qn("c:firstSliceAng"))
        if ang_el is None:
            ang_el = plot_el.makeelement(qn("c:firstSliceAng"), {"val": str(angle_deg)})
            plot_el.append(ang_el)
        else:
            ang_el.set("val", str(angle_deg))
    except Exception as exc:
        log.debug("Could not set firstSliceAng=%d: %s", angle_deg, exc)


def _build_chart_data(source_chart, value_field: str, sort: str):
    """Extract CategoryChartData from a chart using breakdown_ids list.

    Single-series mode (breakdown_ids=[] or first element == "general"): plot
    the General/Total row.

    Multi-series mode (any other first breakdown_ids element): one series per
    breakdown category (e.g. Hombre/Mujer for sexo). General row is excluded
    so we never double-plot it alongside the breakdown categories.
    """
    cd = CategoryChartData()
    question = getattr(source_chart, "question", None)
    options = list(question.options) if question else []
    data = getattr(source_chart, "data", {}) or {}
    bds = list(getattr(source_chart, "breakdown_ids", []) or [])
    primary = bds[0] if bds else ""
    is_general = (not primary) or primary.lower() == "general"

    if is_general:
        primary = data.get("General") or data.get("Total") or (next(iter(data.values())) if data else {})
        if not options and primary:
            options = list(primary.keys())
        if sort in ("desc_by_value", "asc_by_value") and primary:
            reverse = sort == "desc_by_value"
            options = sorted(options, key=lambda o: (primary.get(o) or {}).get(value_field, 0), reverse=reverse)
        cd.categories = options
        values = [float((primary.get(o) or {}).get(value_field, 0) or 0) for o in options]
        cd.add_series("", values)
        return cd, values

    # Breakdown chart: every row except General becomes a series.
    cats = [k for k in data.keys() if k.lower() not in ("general", "total")]
    if not cats:
        cats = list(data.keys())

    # Backfill options from data keys when question.options is empty (parity with general branch).
    if not options and cats:
        first_row = data.get(cats[0]) or {}
        options = list(first_row.keys())

    # Sort options by the FIRST category's value so all series share an axis order.
    if options and sort in ("desc_by_value", "asc_by_value") and cats:
        reverse = sort == "desc_by_value"
        first = data.get(cats[0]) or {}
        options = sorted(options, key=lambda o: (first.get(o) or {}).get(value_field, 0), reverse=reverse)
    cd.categories = options

    all_values: list[float] = []
    for cat in cats:
        row = data.get(cat) or {}
        series_values = [float((row.get(o) or {}).get(value_field, 0) or 0) for o in options]
        cd.add_series(cat, series_values)
        all_values.extend(series_values)
    return cd, all_values


def _apply_series_colors(chart, colors: list[str]) -> None:
    """Apply hex color list to chart series/points."""
    try:
        for series_idx, series in enumerate(chart.series):
            if series_idx < len(colors):
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(colors[series_idx].lstrip("#"))
            # For pie/donut, color individual points
            try:
                for point_idx, point in enumerate(series.points):
                    color = colors[point_idx % len(colors)] if colors else "#7F7F7F"
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#"))
            except Exception:
                pass  # Not all chart types support per-point coloring
    except Exception as exc:
        log.debug("Could not apply series colors: %s", exc)


def _apply_labels(chart, labels_cfg: dict, ctx: RenderContext) -> None:
    """Apply label settings to all plot series.

    Always sets show_value/show_percentage/show_category_name EXPLICITLY so
    python-pptx defaults (which include the raw value for pies) don't leak.
    """
    if not labels_cfg:
        return
    try:
        plot = chart.plots[0]
        show_cat = bool(labels_cfg.get("show_category_name"))
        show_pct = bool(labels_cfg.get("show_percentage"))
        show_val = bool(labels_cfg.get("show_value"))
        plot.has_data_labels = show_cat or show_pct or show_val
        if not plot.has_data_labels:
            return

        dls = plot.data_labels
        dls.show_category_name = show_cat
        dls.show_percentage = show_pct
        dls.show_value = show_val
        # Defensive: explicitly suppress legend key + series name
        try:
            dls.show_legend_key = False
            dls.show_series_name = False
        except Exception:
            pass

        # Number format — controls how percentage is printed (e.g. "0.0%")
        fmt = labels_cfg.get("format")
        if fmt:
            try:
                dls.number_format = fmt
                dls.number_format_is_linked = False
            except Exception:
                pass

        pos_str = labels_cfg.get("position", "outside_end")
        pos = _LABEL_POSITION_MAP.get(pos_str, XL_LABEL_POSITION.OUTSIDE_END)
        try:
            dls.position = pos
        except Exception:
            pass

        font_size = labels_cfg.get("font_size") or ctx.typography.get("label_size", 9)
        try:
            dls.font.size = Pt(font_size)
        except Exception:
            pass
    except Exception as exc:
        log.debug("Could not apply labels: %s", exc)


def _resolve_position(position: dict, ctx: RenderContext) -> tuple[int, int, int, int]:
    """Delegate to pattern_renderer.resolve_position (single source of truth)."""
    from aurum_encuestas.pattern_renderer import resolve_position
    resolved_anchors = getattr(ctx, "resolved_anchors", {}) or {}
    return resolve_position(position, ctx.free_area, resolved_anchors)
