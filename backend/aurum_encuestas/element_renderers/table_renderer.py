"""Table element renderer — dispatches to structure-specific builders."""
from __future__ import annotations

import logging
from typing import TYPE_CHECKING

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

if TYPE_CHECKING:
    from .render_context import RenderContext

log = logging.getLogger(__name__)

_ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

_SEGMENTED_CELLS_FASE_B: dict = {
    "group_header": {
        "style": {"fill": "secondary", "text_color": "background",
                  "font_size": 11, "bold": True, "align_h": "center"},
    },
    "category_header": {
        "style": {"fill": "primary", "text_color": "background",
                  "font_size": 10, "bold": True, "align_h": "center"},
    },
    "counts_row": {
        "style": {"fill": "primary", "text_color": "background",
                  "font_size": 11, "bold": True, "align_h": "center"},
        "label_first_col": "",
    },
    "option_row": {
        "style": {"fill": "primary", "text_color": "#FFFFFF",
                  "font_size": 10, "align_h": "center"},
        "label_col_width_rel": 0.0,
        "value_format": "percentage",
        "value_decimals": 1,
        "minibar": {
            "enabled": True,
            "color_role": "secondary",
            "height_rel_to_cell": 0.25,
            "show_percent_text": True,
            "percent_text_position": "left_of_bar",
        },
    },
}


def render(slide, element: dict, ctx: RenderContext) -> None:
    """Dispatch to the correct table structure builder."""
    structure = element.get("structure", "simple_data")
    if structure == "simple_data":
        _render_simple_data(slide, element, ctx)
    elif structure == "segmented_breakdowns":
        _render_segmented_breakdowns(slide, element, ctx)
    elif structure == "comparison_grid":
        _render_comparison_grid(slide, element, ctx)
    else:
        log.warning("table_renderer: unknown structure %r — falling back to simple_data", structure)
        _render_simple_data(slide, element, ctx)


# ---------------------------------------------------------------------------
# simple_data: header row + data rows
# ---------------------------------------------------------------------------

def _render_simple_data(slide, element: dict, ctx: RenderContext) -> None:
    """Basic table with one header row and N data rows."""
    from .chart_renderer import _resolve_position
    x, y, cx, cy = _resolve_position(element.get("position", {}), ctx)

    data_source = element.get("data_source", {})
    chart_ref_index = data_source.get("chart_ref_index", 0)
    charts_list = getattr(ctx.slide_config, "charts", []) or []
    if chart_ref_index >= len(charts_list):
        log.warning("table_renderer: chart_ref_index %d out of range — skipping", chart_ref_index)
        return

    source_chart = charts_list[chart_ref_index]
    question = getattr(source_chart, "question", None)
    options = question.options if question else []
    data = getattr(source_chart, "data", {}) or {}
    breakdowns = list(data.keys())

    if not breakdowns:
        return

    # Build rows: header + one row per option
    n_cols = len(breakdowns) + 1  # label col + one per breakdown
    n_rows = len(options) + 1     # header + data rows

    try:
        table_shape = slide.shapes.add_table(n_rows, n_cols, Emu(x), Emu(y), Emu(cx), Emu(cy))
        tbl = table_shape.table
    except Exception as exc:
        log.error("table_renderer: failed to add table: %s", exc)
        return

    cells_cfg = element.get("cells", {})
    header_style = (cells_cfg.get("group_header") or {}).get("style", {})

    # Header row
    _set_cell(tbl.cell(0, 0), "Opción", ctx, header_style)
    for col_idx, breakdown in enumerate(breakdowns, start=1):
        _set_cell(tbl.cell(0, col_idx), breakdown, ctx, header_style)

    # Data rows
    option_style = (cells_cfg.get("option_row") or {}).get("style", {})
    for row_idx, option in enumerate(options, start=1):
        _set_cell(tbl.cell(row_idx, 0), option, ctx, option_style)
        for col_idx, breakdown in enumerate(breakdowns, start=1):
            cell_data = (data.get(breakdown) or {}).get(option) or {}
            pct = cell_data.get("pct", 0) or 0
            _set_cell(tbl.cell(row_idx, col_idx), f"{pct * 100:.1f}%", ctx, option_style)


def _set_cell(cell, text: str, ctx: RenderContext, style: dict) -> None:
    """Set cell text and basic style."""
    tf = cell.text_frame
    tf.clear()
    # Disable text wrap so long labels (e.g. "Observaciones") don't push the
    # row height beyond declared cell heights, which breaks compact panels.
    try:
        tf.word_wrap = False
    except Exception:
        pass
    # Anchor text to top of cell so the bottom-anchored minibar overlay never
    # overlaps the percentage label.
    try:
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.TOP
    except Exception:
        pass
    p = tf.paragraphs[0]

    align_h = style.get("align_h", "left")
    p.alignment = _ALIGN_MAP.get(align_h, PP_ALIGN.LEFT)

    run = p.add_run()
    run.text = text

    font_size = style.get("font_size", ctx.typography.get("body_size", 9))
    run.font.size = Pt(font_size)

    if style.get("bold"):
        run.font.bold = True

    font_family = ctx.typography.get("font_family", "Arial")
    run.font.name = font_family

    text_color_role = style.get("text_color", "primary")
    if isinstance(text_color_role, str) and text_color_role.startswith("#"):
        hex_text = text_color_role
    else:
        hex_text = ctx.resolved_colors.get(text_color_role, "#000000")
    try:
        run.font.color.rgb = RGBColor.from_string(hex_text.lstrip("#"))
    except Exception:
        pass

    fill_role = style.get("fill")
    if fill_role:
        if isinstance(fill_role, str) and fill_role.startswith("#"):
            fill_hex = fill_role
        else:
            fill_hex = ctx.resolved_colors.get(fill_role)
        if fill_hex:
            try:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(fill_hex.lstrip("#"))
            except Exception:
                pass


# ---------------------------------------------------------------------------
# segmented_breakdowns — group_header + category_header + counts_row + option rows
# ---------------------------------------------------------------------------

def _render_segmented_breakdowns(slide, element: dict, ctx: RenderContext) -> None:
    """Render one mini-table per breakdown group, auto-wrapping rows.

    Reference layout (Aurora deck): each breakdown group (Edad / Sexo / NSE /
    Punto) becomes its OWN bordered mini-table. Tables are packed into rows by
    weight (= 1 label col + N category cols), wrapping when row weight exceeds
    a tunable limit. Within a row, table widths are proportional to weight.

    Each mini-table rows:
      0  group_header       : group label, spans all category cols
      1  category_header    : one cell per category
      2  counts_row         : "Observaciones" + N per category
      3+ option_row × N     : pct (+ optional minibar) per (option, category)

    Data source: source_chart.all_breakdowns_data — single chart_ref_index
    delivers every demographic, so users don't need N charts.
    """
    from .chart_renderer import _resolve_position
    box_x, box_y, box_cx, box_cy = _resolve_position(element.get("position", {}), ctx)

    # Extend the box vertically to fill the available free area below box_y.
    # Pattern h_rel is often too short for compact multi-panel render — when
    # libreoffice's PDF export inflates small row heights to its internal
    # minimum, declared rows past the bottom get dropped silently.
    fa_y = ctx.free_area.get("y", 0)
    fa_cy = ctx.free_area.get("cy", 1)
    bottom_limit = fa_y + int(fa_cy * 0.97)
    box_cy = max(box_cy, bottom_limit - box_y)

    data_source = element.get("data_source", {})
    chart_ref_index = data_source.get("chart_ref_index", 0)
    breakdown_groups_filter = data_source.get("breakdown_groups", "all_except_general")

    charts_list = getattr(ctx.slide_config, "charts", []) or []
    if chart_ref_index >= len(charts_list):
        log.warning("table_renderer segmented: chart_ref_index out of range — skipping")
        return

    source_chart = charts_list[chart_ref_index]
    question = getattr(source_chart, "question", None)
    options = question.options if question else []
    all_bds: dict = getattr(source_chart, "all_breakdowns_data", {}) or {}

    # Pick which groups to include
    if breakdown_groups_filter == "all_except_general":
        group_ids = [bid for bid in all_bds.keys() if bid.lower() != "general"]
    elif breakdown_groups_filter == "all":
        group_ids = list(all_bds.keys())
    elif isinstance(breakdown_groups_filter, list):
        group_ids = [bid for bid in breakdown_groups_filter if bid in all_bds]
    else:
        group_ids = list(all_bds.keys())

    # Materialize panel descriptors
    panels: list[dict] = []
    for gid in group_ids:
        bd = all_bds.get(gid) or {}
        cats = bd.get("categories") or {}
        if not cats:
            continue
        panels.append({
            "group_id": gid,
            "label": bd.get("label") or gid,
            "cats": list(cats.items()),  # preserve declared order
        })

    if not panels:
        log.debug("table_renderer segmented: no breakdown groups — skipping")
        return

    # Unified Fase B style: single source of truth for all branches.
    group_hdr_cfg = _SEGMENTED_CELLS_FASE_B["group_header"]
    cat_hdr_cfg   = _SEGMENTED_CELLS_FASE_B["category_header"]
    counts_cfg    = _SEGMENTED_CELLS_FASE_B["counts_row"]
    option_cfg    = _SEGMENTED_CELLS_FASE_B["option_row"]

    # Single-panel image-style mode: when only ONE breakdown was requested,
    # render a single, centered, image-faithful mini-table (matches the
    # 'Rango de edad' reference screenshot). Overrides the multi-panel
    # weight-packing logic below.
    if len(panels) == 1:
        only = panels[0]
        _render_panel(
            slide=slide,
            panel=only,
            options=options,
            x=box_x, y=box_y, cx=box_cx, cy=box_cy,
            ctx=ctx,
            group_hdr_cfg=group_hdr_cfg,
            cat_hdr_cfg=cat_hdr_cfg,
            counts_cfg=counts_cfg,
            option_cfg=option_cfg,
            font_cap={"group_header": None, "category_header": None, "counts_row": None, "option_row": None},
        )
        return

    # Pack panels into rows. Weight depends on label_col_width_rel.
    # MAX_ROW_WEIGHT 14 lets Edad(3)+Sexo(3)+NSE(6)=12 fit on row 1, Punto(6) wraps.
    MAX_ROW_WEIGHT = element.get("max_row_weight", 14)
    panel_rows: list[list[dict]] = _pack_panels_into_rows(
        panels, MAX_ROW_WEIGHT,
        label_col_width_rel=_SEGMENTED_CELLS_FASE_B["option_row"]["label_col_width_rel"],
    )

    H_GAP_EMU = int(0.012 * box_cx)  # gap between panels in a row
    V_GAP_EMU = int(0.030 * box_cy)  # gap between rows
    n_rows = len(panel_rows)
    # Each panel renders 5 sub-rows (group_header + cat_header + counts + N option
    # rows). At 8pt body font + small padding, ~0.55cm per sub-row is the floor
    # before libreoffice starts overflowing cells. Allow the table region to
    # extend past the pattern's declared box_cy when needed — better than
    # truncating option rows.
    MIN_SUBROW_EMU = 240000  # ~0.66 cm — enough for minibar + text without clipping
    n_subrows = 3 + max(len(options), 1)
    panel_needed = MIN_SUBROW_EMU * n_subrows
    row_h_fit = (box_cy - V_GAP_EMU * (n_rows - 1)) // max(n_rows, 1)
    row_h = max(panel_needed, row_h_fit)

    use_label = _SEGMENTED_CELLS_FASE_B["option_row"]["label_col_width_rel"] > 0.001
    cur_y = box_y
    for row in panel_rows:
        row_weight = sum((1 + len(p["cats"])) if use_label else len(p["cats"]) for p in row)
        avail_w = box_cx - H_GAP_EMU * (len(row) - 1)
        cur_x = box_x
        for p in row:
            w = (1 + len(p["cats"])) if use_label else len(p["cats"])
            panel_w = int(avail_w * (w / row_weight)) if row_weight else avail_w
            _render_panel(
                slide=slide,
                panel=p,
                options=options,
                x=cur_x, y=cur_y, cx=panel_w, cy=row_h,
                ctx=ctx,
                group_hdr_cfg=group_hdr_cfg,
                cat_hdr_cfg=cat_hdr_cfg,
                counts_cfg=counts_cfg,
                option_cfg=option_cfg,
                font_cap={"group_header": None, "category_header": None,
                          "counts_row": None, "option_row": None},
            )
            cur_x += panel_w + H_GAP_EMU
        cur_y += row_h + V_GAP_EMU


def _pack_panels_into_rows(
    panels: list[dict],
    max_row_weight: int,
    charts_by_bd: dict | None = None,
    label_col_width_rel: float = 0.18,
) -> list[list[dict]]:
    """Greedy row packing by weight.

    Weight per panel:
      - label_col_width_rel > 0: 1 + len(cats)  (label col counts)
      - else (Fase B):           len(cats)
    """
    _ = charts_by_bd  # backward-compat API
    use_label = label_col_width_rel > 0.001

    def weight(p: dict) -> int:
        return (1 + len(p["cats"])) if use_label else len(p["cats"])

    rows: list[list[dict]] = []
    current: list[dict] = []
    current_w = 0
    for p in panels:
        w = weight(p)
        if current and current_w + w > max_row_weight:
            rows.append(current)
            current = [p]
            current_w = w
        else:
            current.append(p)
            current_w += w
    if current:
        rows.append(current)
    return rows


def _render_panel(
    slide, panel: dict, options: list,
    x: int, y: int, cx: int, cy: int,
    ctx: RenderContext,
    group_hdr_cfg: dict, cat_hdr_cfg: dict,
    counts_cfg: dict, option_cfg: dict,
    font_cap: dict | None = None,
) -> None:
    """Render one mini-table (1 breakdown group) at (x,y,cx,cy) EMU."""
    cats: list[tuple[str, dict]] = panel["cats"]  # [(cat_label, opt_cells_dict)]
    n_cat = len(cats)
    N_HEADER_ROWS = 3
    n_rows = N_HEADER_ROWS + len(options)

    table_h_local = cy
    chart_x = chart_y = chart_h = chart_w = 0
    table_w = cx

    # Pre-compute label col width.  When label_col_width_rel <= 0, skip it
    # entirely so data columns start at col 0 (Fase B no-label-col mode).
    label_col_w_rel = option_cfg.get("label_col_width_rel", 0.18)
    if label_col_w_rel <= 0.001:
        n_cols = n_cat
        label_w = 0
    else:
        n_cols = 1 + n_cat
        MIN_LABEL_EMU = 500000  # ~0.55 cm
        label_w = max(MIN_LABEL_EMU, int(label_col_w_rel * table_w))
        if label_w > table_w * 0.4:
            label_w = int(table_w * 0.4)
            chart_x = chart_w = 0  # unused

    first_data_col = 0 if label_w == 0 else 1

    try:
        table_shape = slide.shapes.add_table(n_rows, n_cols, Emu(x), Emu(y), Emu(table_w), Emu(table_h_local))
        tbl = table_shape.table
    except Exception as exc:
        log.error("_render_panel: add_table failed for %s: %s", panel.get("label"), exc)
        return

    # Cap sub-row font sizes so 5 sub-rows fit inside compact panel heights
    # (libreoffice expands cells when content is taller than declared height,
    # pushing later option rows below the panel and breaking the layout).
    # font_cap=None → multi-panel compact caps (9/8/7/7pt).
    # font_cap with None values per role → single-panel mode, bypass caps.
    def _capped(role: str, default_size: int, default_cap: int) -> int:
        if role == "group_header":
            raw = (group_hdr_cfg.get("style") or {}).get("font_size")
        elif role == "category_header":
            raw = (cat_hdr_cfg.get("style") or {}).get("font_size")
        elif role == "counts_row":
            raw = (counts_cfg.get("style") or {}).get("font_size")
        else:
            raw = (option_cfg.get("style") or {}).get("font_size")
        size = raw if raw is not None else default_size
        cap = (font_cap or {}).get(role, default_cap)
        if cap is None:
            return size
        return min(size, cap)

    g_style   = {**(group_hdr_cfg.get("style", {}) or {}), "font_size": _capped("group_header", 9, 9)}
    c_style   = {**(cat_hdr_cfg.get("style", {}) or {}),   "font_size": _capped("category_header", 8, 8)}
    cnt_style = {**(counts_cfg.get("style", {}) or {}),    "font_size": _capped("counts_row", 7, 7)}
    opt_style = {**(option_cfg.get("style", {}) or {}),    "font_size": _capped("option_row", 7, 7)}
    opt_label_style = {**opt_style}

    # Row 0 — group_header (spans all cols)
    _set_cell(tbl.cell(0, 0), panel["label"], ctx, g_style)
    for col_idx in range(1, n_cols):
        _set_cell(tbl.cell(0, col_idx), "", ctx, g_style)
    try:
        tbl.cell(0, 0).merge(tbl.cell(0, n_cols - 1))
    except Exception as exc:
        log.debug("group_header merge failed for %s: %s", panel.get("label"), exc)

    # Row 1 — category_header
    if label_w > 0:
        _set_cell(tbl.cell(1, 0), "", ctx, c_style)
    for c, (cat_label, _) in enumerate(cats):
        _set_cell(tbl.cell(1, first_data_col + c), cat_label, ctx, c_style)

    # Row 2 — counts_row.  Use full "Observaciones" only when label col is wide
    # enough to render it on one line. Narrow panels (Edad/Sexo 2-cat) leave
    # the label cell empty to match the reference deck style.
    if label_w > 0:
        if label_w >= 1600000:  # ~1.7 cm — enough for "Observaciones" single-line
            counts_label = counts_cfg.get("label_first_col", "Observaciones")
        else:
            counts_label = ""
        _set_cell(tbl.cell(2, 0), counts_label, ctx, cnt_style)
    for c, (_, opt_cells) in enumerate(cats):
        total = sum(int((opt_cells.get(opt) or {}).get("count") or 0) for opt in options)
        _set_cell(tbl.cell(2, first_data_col + c), str(total) if total else "", ctx, cnt_style)

    # Rows 3+ — option rows
    value_format = option_cfg.get("value_format", "percentage")
    value_decimals = option_cfg.get("value_decimals", 1)
    minibar_cfg = option_cfg.get("minibar", {})

    for opt_idx, option in enumerate(options):
        row_idx = N_HEADER_ROWS + opt_idx
        if label_w > 0:
            _set_cell(tbl.cell(row_idx, 0), option, ctx, opt_label_style)
        for c, (_, opt_cells) in enumerate(cats):
            cd = opt_cells.get(option) or {}
            pct = cd.get("pct", 0) or 0
            count = cd.get("count", 0) or 0
            if value_format == "percentage":
                val_str = f"{pct * 100:.{value_decimals}f}%"
            elif value_format == "count":
                val_str = str(int(count))
            else:
                val_str = f"{pct * 100:.{value_decimals}f}% ({int(count)})"
            _set_cell(tbl.cell(row_idx, first_data_col + c), val_str, ctx, opt_style)

    # Column widths
    if label_w > 0:
        data_w = (table_w - label_w) // max(n_cat, 1)
        try:
            tbl.columns[0].width = label_w
            for i in range(1, n_cols):
                tbl.columns[i].width = data_w
        except Exception as exc:
            log.debug("Could not set column widths for %s: %s", panel.get("label"), exc)
    else:
        data_w = table_w // max(n_cat, 1)
        try:
            for i in range(n_cols):
                tbl.columns[i].width = data_w
        except Exception as exc:
            log.debug("Could not set column widths for %s: %s", panel.get("label"), exc)

    # Force explicit row heights so libreoffice/PowerPoint don't auto-grow rows
    # beyond the declared table cy (which causes the last option rows to spill
    # below the panel and overlap the next row's panels).
    try:
        per_row = table_h_local // max(n_rows, 1)
        for r in range(n_rows):
            tbl.rows[r].height = Emu(per_row)
    except Exception as exc:
        log.debug("Could not set row heights for %s: %s", panel.get("label"), exc)

    # Minibar overlays — local to this mini-table
    # When label col is absent, data cols start at index 0, so minibar overlay
    # must enumerate starting at first_data_col (0 instead of 1).
    if minibar_cfg.get("enabled", False):
        legacy_data = {}
        legacy_keys = []
        for i, (cat_label, opt_cells) in enumerate(cats):
            k = f"{panel['group_id']}::{cat_label}::{i}"
            legacy_data[k] = opt_cells
            legacy_keys.append(k)
        _render_minibar_overlays(
            slide, tbl, options, legacy_keys, legacy_data,
            minibar_cfg, opt_style, ctx,
            table_x=x, table_y=y, table_cx=table_w, table_cy=cy,
            n_header_rows=N_HEADER_ROWS,
            first_data_col=first_data_col,
        )

    # Mini-chart suppressed — every breakdown renders as table-only for #14
    # fidelity. (chart_x/chart_w/chart_h kept above for future re-enabling.)


def _render_panel_mini_chart(
    slide, panel: dict, options: list, chart,
    x: int, y: int, cx: int, cy: int,
    ctx: RenderContext,
) -> None:
    """Render a small chart visualizing the breakdown's data.

    Categories = breakdown categories. Series = question options. Values = pct.
    chart_type orientation:
      BAR_HORIZONTAL/BAR_*  → horizontal bars
      COLUMN_*              → vertical columns
      PIE/DONUT             → fallback to BAR_HORIZONTAL (pie per-breakdown
                              doesn't make sense for clustered category view)
      else                  → BAR_HORIZONTAL
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    from .chart_renderer import (
        _CHART_TYPE_MAP,
        _apply_labels,
        _apply_series_colors,
    )

    chart_type_raw = getattr(chart, "chart_type", "BAR_HORIZONTAL")
    if chart_type_raw in ("PIE", "DONUT"):
        chart_type = "BAR_HORIZONTAL"
    else:
        chart_type = chart_type_raw
    xl_type = _CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.BAR_CLUSTERED)

    cats: list[tuple[str, dict]] = panel["cats"]
    cat_labels = [c[0] for c in cats]

    cd = CategoryChartData()
    cd.categories = cat_labels
    for opt in options:
        values = []
        for _, opt_cells in cats:
            cell = opt_cells.get(opt) or {}
            pct = cell.get("pct") or 0
            values.append(float(pct) * 100.0)
        cd.add_series(opt, values)

    try:
        cs = slide.shapes.add_chart(xl_type, Emu(x), Emu(y), Emu(cx), Emu(cy), cd)
    except Exception as exc:
        log.error("_render_panel_mini_chart: add_chart failed: %s", exc)
        return
    c = cs.chart

    # Per-series colors from the matching user chart's color picks
    per_chart_colors = list(getattr(chart, "colors", []) or [])
    fallback = ctx.chart_colors or []
    effective: list[str] = []
    for i in range(max(len(options), len(per_chart_colors), len(fallback))):
        if i < len(per_chart_colors) and per_chart_colors[i]:
            effective.append(per_chart_colors[i])
        elif fallback:
            effective.append(fallback[i % len(fallback)])
        else:
            effective.append("#7F7F7F")
    _apply_series_colors(c, effective)

    # Compact: no legend, no title, no axes — pure bars with value labels only.
    c.has_legend = False
    c.has_title = False
    try:
        c.category_axis.visible = False
    except Exception:
        try:
            c.category_axis.tick_labels.font.size = Pt(6)
        except Exception:
            pass
    try:
        c.value_axis.visible = False
    except Exception:
        try:
            c.value_axis.tick_labels.font.size = Pt(6)
        except Exception:
            pass

    _apply_labels(c, {
        "show_value": True,
        "show_percentage": False,
        "show_category_name": False,
        "format": "0.0\\%",
        "position": "outside_end",
        "font_size": 7,
    }, ctx)


def _apply_table_layout(tbl, layout_cfg: dict, total_cx: int, total_cy: int, n_options: int) -> None:
    """Set column widths and row heights from layout config."""
    try:
        col_widths = layout_cfg.get("col_widths", "auto")
        if col_widths == "equal" or col_widths == "auto":
            per_col = total_cx // max(len(tbl.columns), 1)
            for col in tbl.columns:
                col.width = per_col
        elif isinstance(col_widths, list):
            for i, w_rel in enumerate(col_widths):
                if i < len(tbl.columns):
                    tbl.columns[i].width = int(w_rel * total_cx)
    except Exception as exc:
        log.debug("Could not set column widths: %s", exc)


def _render_minibar_overlays(
    slide, tbl, options: list, breakdown_keys: list, data: dict,
    minibar_cfg: dict, opt_style: dict, ctx: RenderContext,
    table_x: int, table_y: int, table_cx: int, table_cy: int,
    n_header_rows: int,
    first_data_col: int = 1,
) -> None:
    """Add MSO rectangle overlays on option row cells to represent minibar values."""
    color_role = minibar_cfg.get("color_role", "primary")
    bar_hex = ctx.resolved_colors.get(color_role, "#7F7F7F")
    height_rel = minibar_cfg.get("height_rel_to_cell", 0.4)
    show_pct_text = minibar_cfg.get("show_percent_text", False)
    pct_text_pos = minibar_cfg.get("percent_text_position", "left_of_bar")

    try:
        # Compute cumulative column X offsets
        col_x_offsets = [0]
        for col in tbl.columns:
            col_x_offsets.append(col_x_offsets[-1] + col.width)

        # Compute cumulative row Y offsets
        row_y_offsets = [0]
        for row in tbl.rows:
            row_y_offsets.append(row_y_offsets[-1] + row.height)

        for opt_idx in range(len(options)):
            row_idx = n_header_rows + opt_idx
            if row_idx >= len(row_y_offsets) - 1:
                continue
            cell_y_abs = table_y + row_y_offsets[row_idx]
            cell_h = tbl.rows[row_idx].height
            bar_h = int(cell_h * height_rel)
            # Anchor bar to BOTTOM so cell value text (rendered in upper half)
            # stays legible above the minibar overlay (reference deck layout).
            bar_y = cell_y_abs + cell_h - bar_h - int(cell_h * 0.05)

            option = options[opt_idx]
            for col_idx, bd_key in enumerate(breakdown_keys, start=first_data_col):
                if col_idx >= len(col_x_offsets):
                    continue
                cell_x_abs = table_x + col_x_offsets[col_idx]
                cell_w = tbl.columns[col_idx].width

                bd_data = data.get(bd_key) or {}
                pct = (bd_data.get(option) or {}).get("pct", 0) or 0
                bar_w = int(cell_w * float(pct))
                if bar_w <= 0:
                    continue

                # Add minibar rectangle
                try:
                    bar_shape = slide.shapes.add_shape(
                        1,  # MSO_SHAPE_TYPE.RECTANGLE
                        Emu(cell_x_abs), Emu(bar_y), Emu(bar_w), Emu(bar_h),
                    )
                    bar_shape.fill.solid()
                    bar_shape.fill.fore_color.rgb = RGBColor.from_string(bar_hex.lstrip("#"))
                    bar_shape.line.fill.background()
                except Exception as exc:
                    log.debug("minibar overlay: could not add rect: %s", exc)
                    continue

                # Percent text overlay
                if show_pct_text:
                    pct_str = f"{pct * 100:.1f}%"
                    if pct_text_pos == "left_of_bar":
                        txt_x = cell_x_abs
                        txt_w = bar_w
                    elif pct_text_pos == "right_of_bar":
                        txt_x = cell_x_abs + bar_w
                        txt_w = cell_w - bar_w
                    else:  # inside_bar
                        txt_x = cell_x_abs
                        txt_w = bar_w
                    try:
                        txt_shape = slide.shapes.add_textbox(
                            Emu(txt_x), Emu(bar_y), Emu(max(txt_w, 1)), Emu(bar_h),
                        )
                        tf = txt_shape.text_frame
                        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                        run = tf.paragraphs[0].add_run()
                        run.text = pct_str
                        run.font.size = Pt(ctx.typography.get("label_size", 8))
                        run.font.name = ctx.typography.get("font_family", "Arial")
                    except Exception as exc:
                        log.debug("minibar text overlay failed: %s", exc)
    except Exception as exc:
        log.warning("minibar overlays: error: %s", exc)


# ---------------------------------------------------------------------------
# comparison_grid — side-by-side breakdown comparison
# ---------------------------------------------------------------------------

def _render_comparison_grid(slide, element: dict, ctx: RenderContext) -> None:
    """Side-by-side comparison of breakdown data for the same question.

    Layout:
      Row 0: header — blank + one col per breakdown
      Row 1+: one row per question option, values side-by-side
    """
    from .chart_renderer import _resolve_position
    x, y, cx, cy = _resolve_position(element.get("position", {}), ctx)

    data_source = element.get("data_source", {})
    chart_ref_index = data_source.get("chart_ref_index", 0)
    charts_list = getattr(ctx.slide_config, "charts", []) or []
    if chart_ref_index >= len(charts_list):
        log.warning("comparison_grid: chart_ref_index out of range — skipping")
        return

    source_chart = charts_list[chart_ref_index]
    question = getattr(source_chart, "question", None)
    options = question.options if question else []
    data = getattr(source_chart, "data", {}) or {}
    breakdown_keys = list(data.keys())

    if not breakdown_keys or not options:
        return

    n_cols = 1 + len(breakdown_keys)
    n_rows = 1 + len(options)

    try:
        table_shape = slide.shapes.add_table(n_rows, n_cols, Emu(x), Emu(y), Emu(cx), Emu(cy))
        tbl = table_shape.table
    except Exception as exc:
        log.error("comparison_grid: failed to add table: %s", exc)
        return

    cells_cfg = element.get("cells", {})
    header_style = (cells_cfg.get("group_header") or {}).get("style", {})
    option_style = (cells_cfg.get("option_row") or {}).get("style", {})
    option_cfg = cells_cfg.get("option_row", {})
    value_format = option_cfg.get("value_format", "percentage")
    value_decimals = option_cfg.get("value_decimals", 1)

    # Header row
    _set_cell(tbl.cell(0, 0), "Opción", ctx, header_style)
    for col_idx, bd_key in enumerate(breakdown_keys, start=1):
        _set_cell(tbl.cell(0, col_idx), bd_key, ctx, header_style)

    # Data rows
    for row_idx, option in enumerate(options, start=1):
        _set_cell(tbl.cell(row_idx, 0), option, ctx, option_style)
        for col_idx, bd_key in enumerate(breakdown_keys, start=1):
            bd_data = data.get(bd_key) or {}
            cell_data = bd_data.get(option) or {}
            pct = cell_data.get("pct", 0) or 0
            count = cell_data.get("count", 0) or 0
            if value_format == "percentage":
                val_str = f"{pct * 100:.{value_decimals}f}%"
            elif value_format == "count":
                val_str = str(int(count))
            else:
                val_str = f"{pct * 100:.{value_decimals}f}% ({int(count)})"
            _set_cell(tbl.cell(row_idx, col_idx), val_str, ctx, option_style)
