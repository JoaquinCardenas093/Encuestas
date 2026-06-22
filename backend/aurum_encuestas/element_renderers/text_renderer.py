"""Text element renderer — handles static, analysis, and computed content sources."""
from __future__ import annotations

import logging
from typing import TYPE_CHECKING

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

if TYPE_CHECKING:
    from .render_context import RenderContext

log = logging.getLogger(__name__)

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def render(slide, element: dict, ctx: RenderContext) -> None:
    """Render a text element onto slide."""
    content_source = element.get("content_source", {})
    text = _resolve_content(content_source, ctx)
    if text is None:
        log.debug("text_renderer: no content resolved for element %r — skipping", element.get("id"))
        return

    from .chart_renderer import _resolve_position
    x, y, cx, cy = _resolve_position(element.get("position", {}), ctx)

    tb_shape = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(cx), Emu(cy))
    tf = tb_shape.text_frame
    tf.word_wrap = True

    style = element.get("style") or {}
    padding = style.get("padding", 0)
    if padding:
        try:
            tf.margin_left = Pt(padding)
            tf.margin_right = Pt(padding)
            tf.margin_top = Pt(padding / 2)
            tf.margin_bottom = Pt(padding / 2)
        except Exception:
            pass

    # Set text + formatting
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text

    font_size = style.get("font_size", ctx.typography.get("body_size", 10))
    run.font.size = Pt(font_size)

    bold = style.get("bold", False)
    if bold:
        run.font.bold = True

    font_family = ctx.typography.get("font_family", "Arial")
    run.font.name = font_family

    text_color_role = style.get("text_color", "primary")
    hex_color = ctx.resolved_colors.get(text_color_role, "#000000")
    try:
        run.font.color.rgb = RGBColor.from_string(hex_color.lstrip("#"))
    except Exception:
        pass

    align_h = style.get("align_h", "left")
    p.alignment = _ALIGN_MAP.get(align_h, PP_ALIGN.LEFT)

    # Fill background
    fill_role = style.get("fill")
    if fill_role:
        fill_hex = ctx.resolved_colors.get(fill_role)
        if fill_hex:
            try:
                tb_shape.fill.solid()
                tb_shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex.lstrip("#"))
            except Exception:
                pass

    # Border left — implemented as a separate thin rectangle shape
    border_left = style.get("border_left")
    if border_left:
        bl_color_role = border_left.get("color", "primary")
        bl_width_pt = border_left.get("width_pt", 3.0)
        bl_hex = ctx.resolved_colors.get(bl_color_role, "#7F7F7F")
        _add_border_left_rect(slide, x, y, cy, bl_hex, bl_width_pt)


def _resolve_content(content_source: dict, ctx: RenderContext) -> str | None:
    source_type = content_source.get("type", "static")
    if source_type == "static":
        return content_source.get("text", "")
    if source_type == "analysis":
        scope = content_source.get("scope", "slide")
        ref_index = content_source.get("ref_index", 0)
        analyses = getattr(ctx.slide_config, "analyses", []) or []
        matching = [a for a in analyses if getattr(a, "scope", None) == scope]
        if not matching:
            return None
        if scope == "chart":
            idx = min(ref_index, len(matching) - 1)
            return matching[idx].text
        return matching[0].text
    if source_type == "computed":
        # `field` reads an attribute off the slide_config (e.g. "title").
        # Returns None if empty so renderer skips drawing.
        field = content_source.get("field")
        if field:
            val = getattr(ctx.slide_config, field, None)
            if val is None or (isinstance(val, str) and not val.strip()):
                return None
            return str(val)
        return content_source.get("text", "")
    return content_source.get("text")


def _add_border_left_rect(slide, x: int, y: int, cy: int, hex_color: str, width_pt: float) -> None:
    from pptx.util import Pt as PtUtil
    try:
        width_emu = int(PtUtil(width_pt))
        rect = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Emu(x), Emu(y), Emu(width_emu), Emu(cy),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))
        rect.line.fill.background()
    except Exception as exc:
        log.debug("Could not add border_left rect: %s", exc)
