import base64
import logging
import tempfile
import uuid
from contextlib import asynccontextmanager
from datetime import UTC
from datetime import datetime as dt
from pathlib import Path
from typing import Literal

from fastapi import BackgroundTasks, FastAPI, File, HTTPException, Request, UploadFile
from fastapi.exceptions import RequestValidationError
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel

from .config import add_recent, get_corpus_dir, get_render_cache_dir, load_recents
from .errors import AurumError
from .llm_client import correct_slide_layout, generate_analysis, suggest_layout
from .models import ProjectState
from .pptx_generator import build_pptx
from .pptx_template import load_template
from .project_store import load_project, save_project
from .render_service import render_slide_to_png
from .style_guide import (
    BUILTIN_STYLE_GUIDE,
    Pattern,
    load_active_style_guide,
    migrate_legacy_files,
    save_style_guide,
)
from .xlsx_parser import parse_xlsx

log = logging.getLogger(__name__)


@asynccontextmanager
async def lifespan(app):
    # Startup
    try:
        migrate_legacy_files()
        log.info("M6 migration check complete")
    except Exception as exc:
        log.warning("migrate_legacy_files failed (non-fatal): %s", exc)
    yield
    # Shutdown (nothing to do)


app = FastAPI(title="AurumEncuestas API", version="0.1.0", lifespan=lifespan)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173"],
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.exception_handler(AurumError)
async def handle_aurum_error(request, exc: AurumError):
    return JSONResponse(status_code=exc.status, content={"code": exc.code, "message": str(exc)})


@app.exception_handler(RequestValidationError)
async def handle_validation_error(request: Request, exc: RequestValidationError):
    print(f"[VALIDATION 422] {request.method} {request.url.path}")
    print(f"  errors: {exc.errors()}")
    try:
        body = await request.body()
        snippet = body[:2000].decode("utf-8", errors="replace")
        print(f"  body[:2000]: {snippet}")
    except Exception as e:
        print(f"  body read failed: {e}")
    return JSONResponse(status_code=422, content={"code": "validation_error", "detail": exc.errors()})


@app.get("/api/health")
async def health():
    return {"status": "ok"}


async def _save_upload_tmp(file: UploadFile, suffix: str) -> str:
    contents = await file.read()
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    tmp.write(contents)
    tmp.close()
    return tmp.name


def _persist_upload(file_bytes: bytes, original_name: str) -> str:
    """Save uploaded file to ~/.aurum/uploads/ keyed by filename. Returns absolute path."""
    from .config import get_aurum_dir
    uploads_dir = get_aurum_dir() / "uploads"
    uploads_dir.mkdir(parents=True, exist_ok=True)
    safe_name = Path(original_name).name  # strip any path component
    dest = uploads_dir / safe_name
    dest.write_bytes(file_bytes)
    return str(dest.resolve())


@app.post("/api/parse-xlsx")
async def parse_xlsx_endpoint(file: UploadFile = File(...)):
    contents = await file.read()
    persisted_path = _persist_upload(contents, file.filename or "uploaded.xlsx")
    db = parse_xlsx(persisted_path)
    result = db.model_dump()
    result["persisted_path"] = persisted_path
    return result


@app.post("/api/parse-template")
async def parse_template_endpoint(file: UploadFile = File(...)):
    contents = await file.read()
    persisted_path = _persist_upload(contents, file.filename or "template.pptx")
    info = load_template(persisted_path)
    result = info.model_dump()
    result["persisted_path"] = persisted_path
    return result


class SaveProjectRequest(BaseModel):
    path: str
    state: dict


@app.post("/api/save-project")
async def save_project_endpoint(req: SaveProjectRequest):
    state = ProjectState.model_validate(req.state)
    save_project(state, req.path)
    add_recent(req.path, state.project_name)
    return {"saved": True, "path": req.path}


class LoadProjectRequest(BaseModel):
    path: str


@app.post("/api/load-project")
async def load_project_endpoint(req: LoadProjectRequest):
    state = load_project(req.path)
    return state.model_dump()


class PreviewSlideRequest(BaseModel):
    state: dict
    slide_index: int = 0


@app.post("/api/preview-slide")
async def preview_slide_endpoint(req: PreviewSlideRequest):
    """Build pptx in-memory from project state, render specified slide as base64 PNG."""
    state = ProjectState.model_validate(req.state)
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp_path = tmp.name
    try:
        build_pptx(state, tmp_path)
        png_bytes = render_slide_to_png(tmp_path, slide_index=req.slide_index)
        return {"png_base64": base64.b64encode(png_bytes).decode("utf-8")}
    finally:
        Path(tmp_path).unlink(missing_ok=True)


class ExportPptxRequest(BaseModel):
    state: dict
    path: str


@app.post("/api/export-pptx")
async def export_pptx_endpoint(req: ExportPptxRequest):
    """Build and export a PPTX file from ProjectState to the given path."""
    state = ProjectState.model_validate(req.state)
    expanded = str(Path(req.path).expanduser())
    Path(expanded).parent.mkdir(parents=True, exist_ok=True)
    build_pptx(state, expanded)
    size = Path(expanded).stat().st_size if Path(expanded).exists() else 0
    return {"exported": True, "path": expanded, "size": size}


class GenerateAnalysisRequest(BaseModel):
    scope: str
    context: dict
    state: dict | None = None  # ProjectState for backend data extraction
    slide_id: str | None = None
    target_id: str | None = None


@app.post("/api/generate-analysis")
async def generate_analysis_endpoint(req: GenerateAnalysisRequest):
    ctx = dict(req.context)
    # Backend extracts real data when state + slide_id provided.
    if req.state and req.slide_id:
        try:
            ctx = _build_analysis_context(req.scope, req.target_id, req.slide_id, req.state, ctx)
        except Exception:
            pass  # keep frontend-provided ctx as fallback
    try:
        text = generate_analysis(req.scope, ctx)
        return {"text": text, "fallback": False}
    except Exception:
        return {"text": "[Análisis no disponible — editar manualmente]", "fallback": True}


def _build_analysis_context(scope: str, target_id: str | None, slide_id: str, state_dict: dict, base_ctx: dict) -> dict:
    """Server-side data extraction for analysis. Aggregates data across charts per scope."""
    from .data_extractor import extract_chart_data, extract_all_breakdowns_data
    state = ProjectState.model_validate(state_dict)
    slide = next((s for s in state.slides if s.id == slide_id), None)
    if slide is None or state.parsed_db is None or not state.inputs:
        return base_ctx
    db_path = state.inputs.db_path
    data_blocks = state.parsed_db.data_blocks or {}

    # Filter charts by scope.
    if scope == "chart":
        charts = [c for c in slide.charts if c.id == target_id]
    elif scope == "question":
        charts = [c for c in slide.charts if c.question_id == target_id]
    else:  # slide → all
        charts = list(slide.charts)
    if not charts:
        return base_ctx

    # Aggregate per chart × per breakdown (every bd, not just primary).
    aggregated: dict = {}
    questions_by_id = {q.id: q for q in state.parsed_db.questions}
    bds_by_id = {b.id: b for b in state.parsed_db.breakdowns}
    for c in charts:
        q = questions_by_id.get(c.question_id)
        if q is None:
            continue
        bd_ids = list(c.breakdown_ids or []) or ["general"]
        for bd_id in bd_ids:
            try:
                data = extract_chart_data(db_path, q, bd_id, data_blocks)
            except Exception:
                data = {}
            bd_obj = bds_by_id.get(bd_id)
            bd_label = bd_obj.label if bd_obj else bd_id
            key = f"{q.code} — {bd_label}"
            aggregated[key] = {
                "question_text": q.text,
                "breakdown_label": bd_label,
                "options": q.options,
                "data": data,
            }

    # Single-bd single-chart: flat shape (current Claude prompt expects this).
    if len(aggregated) == 1:
        first = next(iter(aggregated.values()))
        base_ctx["question_text"] = first["question_text"]
        base_ctx["breakdown_label"] = first["breakdown_label"]
        base_ctx["options"] = first["options"]
        base_ctx["data"] = first["data"]
    else:
        base_ctx["charts"] = aggregated
        # Pass aggregated as the data field too so Claude sees it.
        base_ctx["data"] = aggregated
        # Build human-readable breakdown_label summarizing all bd labels.
        bd_labels = [v["breakdown_label"] for v in aggregated.values()]
        base_ctx["breakdown_label"] = " | ".join(sorted(set(bd_labels)))
        # Use first question_text + options for prompt fields.
        first = next(iter(aggregated.values()))
        base_ctx["question_text"] = first["question_text"]
        base_ctx["options"] = first["options"]
    return base_ctx


# M4 training endpoints (add/list/delete/reprocess/bank) removed in M6.2.
# New corpus/style-guide/cache endpoints added below in M6.8.


# ────────────────────────────────────────────────────────────────────────────
# M6.8 T2: Corpus CRUD endpoints
# ────────────────────────────────────────────────────────────────────────────

def _count_slides_with_charts(pptx_path: Path) -> int:
    """Count slides in PPTX that have at least one chart shape."""
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        return sum(
            1 for slide in prs.slides
            if any(getattr(sh, "has_chart", False) for sh in slide.shapes)
        )
    except Exception:
        return 0


@app.post("/api/training/corpus/add")
async def corpus_add(file: UploadFile = File(...)):
    """Save an uploaded PPT to the corpus directory."""
    filename = file.filename or "upload.pptx"
    if not filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only .pptx files are accepted for the training corpus.")

    corpus_dir = get_corpus_dir()
    dest = corpus_dir / Path(filename).name  # strip any path component
    contents = await file.read()
    dest.write_bytes(contents)

    slides_with_charts = _count_slides_with_charts(dest)
    return {
        "filename": dest.name,
        "slides_with_charts": slides_with_charts,
        "added_at": dt.now(UTC).isoformat(),
    }


@app.get("/api/training/corpus/list")
async def corpus_list():
    """List all PPTs in the corpus directory with metadata."""
    corpus_dir = get_corpus_dir()
    pptxs = []
    for p in sorted(corpus_dir.glob("*.pptx")):
        pptxs.append({
            "filename": p.name,
            "slides_with_charts": _count_slides_with_charts(p),
            "added_at": dt.fromtimestamp(p.stat().st_mtime, UTC).isoformat(),
            "size_bytes": p.stat().st_size,
        })
    return {"pptxs": pptxs}


class CorpusDeleteRequest(BaseModel):
    filename: str


@app.post("/api/training/corpus/delete")
async def corpus_delete(req: CorpusDeleteRequest):
    """Delete a PPT from the corpus by filename."""
    corpus_dir = get_corpus_dir()
    # Safety: strip path components; only allow simple filename
    safe_name = Path(req.filename).name
    target = corpus_dir / safe_name
    # Extra guard: resolved path must be inside corpus_dir
    try:
        target.resolve().relative_to(corpus_dir.resolve())
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid filename — path traversal not allowed.")
    if not target.exists():
        return {"deleted": False, "message": f"{safe_name} not found in corpus"}
    target.unlink()
    return {"deleted": True}


class SuggestLayoutRequest(BaseModel):
    n_charts: int
    chart_types: list[str]
    n_chart_an: int = 0
    n_q_an: int = 0
    has_slide_an: bool = False
    free_area: dict


@app.post("/api/suggest-layout")
async def suggest_layout_endpoint(req: SuggestLayoutRequest):
    return suggest_layout(
        n_charts=req.n_charts,
        chart_types=req.chart_types,
        n_chart_an=req.n_chart_an,
        n_q_an=req.n_q_an,
        has_slide_an=req.has_slide_an,
        free_area=req.free_area,
    )


def _enforce_multi_row(positions: dict, payload_shapes: list) -> None:
    """If Sonnet returned all charts at same y_cm but sum cx > safe canvas width,
    redistribute into 2 rows in-place. Mutates positions dict."""
    EMU = 360000
    SAFE_W_EMU = int(30.6 * EMU)
    SAFE_X_MIN_EMU = int(1.3 * EMU)

    chart_ids = [s["id"][len("chart_"):] for s in payload_shapes if s.get("kind") == "chart" and s["id"][len("chart_"):] in positions]
    if len(chart_ids) < 2:
        return
    # Check if all charts share same y (single row).
    ys = [positions[cid]["y_emu"] for cid in chart_ids]
    if max(ys) - min(ys) > int(0.5 * EMU):
        return  # Already multi-row.
    # Sum cx + gaps.
    total_cx = sum(positions[cid]["cx_emu"] for cid in chart_ids)
    gap = int(0.4 * EMU)
    if total_cx + gap * (len(chart_ids) - 1) <= SAFE_W_EMU:
        return  # Fits.

    # Split into 2 rows. Greedy: fill row 1 until next chart would overflow.
    base_y = ys[0]
    row_h_emu = int(6.5 * EMU)  # ~6.5cm per row
    row1_y = base_y
    row2_y = base_y + row_h_emu + int(0.3 * EMU)

    row1_ids, row2_ids = [], []
    cur_w = 0
    for cid in chart_ids:
        cx = positions[cid]["cx_emu"]
        if cur_w + cx + (gap if cur_w > 0 else 0) <= SAFE_W_EMU:
            row1_ids.append(cid)
            cur_w += cx + (gap if cur_w > 1 else 0)
        else:
            row2_ids.append(cid)
    # If row1 is full and row2 empty, force split.
    if not row2_ids and row1_ids:
        row2_ids.append(row1_ids.pop())

    # Position row 1.
    cur_x = SAFE_X_MIN_EMU
    for cid in row1_ids:
        positions[cid]["x_emu"] = cur_x
        positions[cid]["y_emu"] = row1_y
        positions[cid]["cy_emu"] = row_h_emu
        cur_x += positions[cid]["cx_emu"] + gap
    cur_x = SAFE_X_MIN_EMU
    for cid in row2_ids:
        positions[cid]["x_emu"] = cur_x
        positions[cid]["y_emu"] = row2_y
        positions[cid]["cy_emu"] = row_h_emu
        cur_x += positions[cid]["cx_emu"] + gap


def _rects_overlap(r1: dict, r2: dict) -> bool:
    return not (
        r1["x_emu"] + r1["cx_emu"] <= r2["x_emu"]
        or r2["x_emu"] + r2["cx_emu"] <= r1["x_emu"]
        or r1["y_emu"] + r1["cy_emu"] <= r2["y_emu"]
        or r2["y_emu"] + r2["cy_emu"] <= r1["y_emu"]
    )


def _reposition_analyses(positions: dict, payload_shapes: list) -> None:
    """Slide-scope analyses: if overlap any chart, move to safe band above charts
    (y between question subtitle and first chart row)."""
    EMU = 360000
    SAFE_X_MIN_EMU = int(1.3 * EMU)
    SAFE_W_EMU = int(30.6 * EMU)
    SAFE_Y_MIN_EMU = int(2.8 * EMU)  # below question

    chart_rects = []
    for s in payload_shapes:
        if s.get("kind") != "chart":
            continue
        cid = s["id"][len("chart_"):]
        if cid in positions:
            chart_rects.append(positions[cid])

    if not chart_rects:
        return
    top_chart_y = min(r["y_emu"] for r in chart_rects)
    available_h = top_chart_y - SAFE_Y_MIN_EMU - int(0.2 * EMU)

    for s in payload_shapes:
        if s.get("kind") != "analysis" or s.get("scope") != "slide":
            continue
        aid = s["id"][len("analysis_"):]
        if aid not in positions:
            continue
        a = positions[aid]
        # Check overlap with any chart.
        if not any(_rects_overlap(a, c) for c in chart_rects):
            continue
        # Move to top band, full width.
        a["x_emu"] = SAFE_X_MIN_EMU
        a["y_emu"] = SAFE_Y_MIN_EMU
        a["cx_emu"] = SAFE_W_EMU
        a["cy_emu"] = max(int(1.5 * EMU), available_h)


class SuggestSlideLayoutRequest(BaseModel):
    state: dict
    slide_id: str


@app.post("/api/suggest-slide-layout")
async def suggest_slide_layout_endpoint(req: SuggestSlideLayoutRequest):
    """AI layout corrector. Returns {elements: [{id, x_cm, y_cm, w_cm, h_cm, font_pt?}], changes: [...]}."""
    from .element_renderers.xlsx_builder import compute_xlsx_natural_dim_emu
    state = ProjectState.model_validate(req.state)
    slide = next((s for s in state.slides if s.id == req.slide_id), None)
    if slide is None:
        return {"error": "slide not found", "elements": [], "changes": []}

    EMU_PER_CM = 360000
    payload_shapes = []
    for c in slide.charts:
        # Compute natural dim for TABLE_WITH_MINIBARS; PIE/BAR use 10cm × 7cm default.
        if c.chart_type == "TABLE_WITH_MINIBARS":
            # Need source_chart context; use breakdown_ids real (non-general)
            try:
                # Build minimal source_chart-like object via classifier path
                from .pattern_classifier import enrich_slide_config
                from .data_extractor import extract_all_breakdowns_data
                bds_real = [b for b in c.breakdown_ids if b and b.lower() != "general"]
                q = next((qq for qq in state.parsed_db.questions if qq.id == c.question_id), None) if state.parsed_db else None
                if q and state.parsed_db and state.inputs:
                    all_bds_data = extract_all_breakdowns_data(state.inputs.db_path, q, state.parsed_db.breakdowns, state.parsed_db.data_blocks or {})
                    from types import SimpleNamespace
                    src = SimpleNamespace(question=q, all_breakdowns_data=all_bds_data, breakdown_ids=c.breakdown_ids, show_legend=c.show_legend)
                    w_emu, h_emu = compute_xlsx_natural_dim_emu(src, bds_real)
                else:
                    w_emu, h_emu = 10 * EMU_PER_CM, 7 * EMU_PER_CM
            except Exception:
                w_emu, h_emu = 10 * EMU_PER_CM, 7 * EMU_PER_CM
        else:
            w_emu, h_emu = 10 * EMU_PER_CM, 7 * EMU_PER_CM
        payload_shapes.append({
            "id": f"chart_{c.id}",
            "kind": "chart",
            "chart_type": c.chart_type,
            "title": c.title,
            "w_cm": round(w_emu / EMU_PER_CM, 2),
            "h_cm": round(h_emu / EMU_PER_CM, 2),
        })
    for a in slide.analyses:
        payload_shapes.append({
            "id": f"analysis_{a.id}",
            "kind": "analysis",
            "scope": a.scope,
            "text_chars": len(a.text or ""),
            "text_preview": (a.text or "")[:100],
        })

    # Canvas hint + sum natural widths so Sonnet can detect overflow risk.
    total_natural_w = sum(s.get("w_cm", 0) for s in payload_shapes if s.get("kind") == "chart")
    slide_payload = {
        "title": slide.title,
        "type": slide.type,
        "canvas": {"safe_x_cm": [1.3, 31.9], "safe_y_cm": [3.2, 16.5], "total_w_cm": 30.6, "total_h_cm": 13.3},
        "sum_natural_chart_w_cm": round(total_natural_w, 1),
        "needs_multi_row": total_natural_w > 30.0,
        "shapes": payload_shapes,
    }
    try:
        raw = correct_slide_layout(slide_payload)
    except Exception as e:
        return {"error": str(e), "positions": {}, "changes": []}

    # Convert AI cm coords → EMU positions keyed by element id (stripped of prefix).
    EMU = 360000
    positions = {}
    for el in raw.get("elements", []) or []:
        eid = str(el.get("id", "")).strip()
        if not eid:
            continue
        is_analysis = eid.startswith("analysis_")
        if eid.startswith("chart_"):
            key = eid[len("chart_"):]
        elif is_analysis:
            key = eid[len("analysis_"):]
        else:
            key = eid
        # Default font_pt: analyses=11 (per prompt rule 10.5–11), charts None.
        font_pt = el.get("font_pt")
        if font_pt is None and is_analysis:
            font_pt = 11
        if font_pt is not None:
            try:
                font_pt = float(font_pt)
            except (TypeError, ValueError):
                font_pt = 11 if is_analysis else None
        positions[key] = {
            "x_emu": int(float(el.get("x_cm", 0)) * EMU),
            "y_emu": int(float(el.get("y_cm", 0)) * EMU),
            "cx_emu": int(float(el.get("w_cm", 0)) * EMU),
            "cy_emu": int(float(el.get("h_cm", 0)) * EMU),
            "font_pt": font_pt,
        }
    # Post-process: enforce multi-row + reposition analyses below chart band.
    _enforce_multi_row(positions, payload_shapes)
    _reposition_analyses(positions, payload_shapes)

    # Convert extras (new shapes) cm → EMU.
    extras_emu = []
    for ex in raw.get("extras", []) or []:
        kind = ex.get("kind")
        if kind not in ("line", "callout", "text"):
            continue
        extras_emu.append({
            "kind": kind,
            "x_emu": int(float(ex.get("x_cm", 0)) * EMU),
            "y_emu": int(float(ex.get("y_cm", 0)) * EMU),
            "cx_emu": int(float(ex.get("w_cm", 0)) * EMU),
            "cy_emu": int(float(ex.get("h_cm", 0)) * EMU),
            "text": ex.get("text"),
            "font_pt": ex.get("font_pt"),
            "bold": bool(ex.get("bold", False)),
            "style": ex.get("style"),
            "color": ex.get("color"),
            "fill": ex.get("fill"),
        })
    return {"positions": positions, "extras": extras_emu, "changes": raw.get("changes", [])}


@app.get("/api/recents")
async def recents_endpoint():
    return {"recents": load_recents()}


@app.get("/api/config/recent-colors")
async def recent_colors_endpoint():
    from .color_resolver import get_recent_colors
    return {"recent_colors": get_recent_colors()}


class UpdateRecentColorRequest(BaseModel):
    hex: str


@app.post("/api/config/recent-colors")
async def update_recent_color_endpoint(req: UpdateRecentColorRequest):
    from .color_resolver import update_recent
    update_recent(req.hex)
    from .color_resolver import get_recent_colors
    return {"recent_colors": get_recent_colors()}


# ────────────────────────────────────────────────────────────────────────────
# M6.7: Async AI analysis job endpoints
# ────────────────────────────────────────────────────────────────────────────

_analysis_jobs: dict[str, dict] = {}


@app.post("/api/training/analyze-with-ai")
async def analyze_with_ai(background_tasks: BackgroundTasks):
    """Start async AI analysis job. Returns job_id immediately."""
    job_id = str(uuid.uuid4())
    _analysis_jobs[job_id] = {"progress": 0, "status": "running", "message": "Iniciando..."}

    def _run():
        from .style_guide import load_active_style_guide
        from .style_guide_analyzer import run_full_analysis_pipeline
        try:
            existing_manual = load_active_style_guide().manual_edits or {}
        except Exception:
            existing_manual = {}
        result = run_full_analysis_pipeline(
            progress_dict=_analysis_jobs[job_id],
            existing_manual_edits=existing_manual,
        )
        _analysis_jobs[job_id].update(result)

    background_tasks.add_task(_run)
    return {"job_id": job_id}


@app.get("/api/training/analysis-status/{job_id}")
async def analysis_status(job_id: str):
    """Get progress of an async analysis job."""
    job = _analysis_jobs.get(job_id)
    if job is None:
        raise HTTPException(status_code=404, detail=f"Job {job_id!r} not found")
    return job


# ────────────────────────────────────────────────────────────────────────────
# M6.8 T4: Style guide read + manual pattern edit endpoints
# ────────────────────────────────────────────────────────────────────────────

@app.get("/api/training/style-guide")
async def get_style_guide():
    """Return the active style guide (AI-generated or built-in fallback)."""
    try:
        sg = load_active_style_guide()
    except Exception:
        sg = BUILTIN_STYLE_GUIDE
    return sg.model_dump(by_alias=True)


class PatternUpdateRequest(BaseModel):
    id: str
    priority: int = 0
    trigger: dict = {}
    extends: str | None = None
    best_example: str | None = None
    why_picked: str | None = None
    implementation: dict = {}


@app.put("/api/training/style-guide/pattern/{pattern_id}")
async def update_style_guide_pattern(pattern_id: str, req: PatternUpdateRequest):
    """Manually edit a single pattern in the active style guide.

    Marks manual_edits[pattern_id] = current timestamp so future AI re-analysis
    can detect which patterns have been user-modified.
    """
    try:
        sg = load_active_style_guide()
    except Exception:
        sg = BUILTIN_STYLE_GUIDE

    # Find the pattern to update
    pattern_idx = next(
        (i for i, p in enumerate(sg.patterns) if p.id == pattern_id),
        None,
    )
    if pattern_idx is None:
        raise HTTPException(status_code=404, detail=f"Pattern {pattern_id!r} not found in active style guide.")

    # Update the pattern fields from request
    existing = sg.patterns[pattern_idx]
    updated_data = existing.model_dump(by_alias=True)
    req_dict = req.model_dump()
    # Only override fields that were explicitly provided (non-None)
    updated_data.update({k: v for k, v in req_dict.items() if v is not None})

    try:
        sg.patterns[pattern_idx] = Pattern.model_validate(updated_data)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Invalid pattern data: {exc}")

    # Mark as manually edited
    if sg.manual_edits is None:
        sg.manual_edits = {}
    sg.manual_edits[pattern_id] = dt.now(UTC).isoformat()

    save_style_guide(sg)
    return {"ok": True, "pattern_id": pattern_id, "edited_at": sg.manual_edits[pattern_id]}


# ────────────────────────────────────────────────────────────────────────────
# M6.8 T5: Cache clear endpoint
# ────────────────────────────────────────────────────────────────────────────

class ClearCacheRequest(BaseModel):
    cache_type: Literal["render", "classifier", "all"]


@app.post("/api/training/clear-cache")
async def clear_cache(req: ClearCacheRequest):
    """Clear one or all backend caches.

    cache_type options:
      - "render": deletes all PNG files in ~/.aurum/training/render_cache/
      - "classifier": clears in-memory pattern classifier LRU dict
      - "all": both of the above
    """
    cleared: dict[str, int] = {}

    if req.cache_type in ("render", "all"):
        cache_dir = get_render_cache_dir()
        png_files = list(cache_dir.glob("*.png"))
        for f in png_files:
            try:
                f.unlink()
            except Exception:
                pass
        cleared["render"] = len(png_files)

    if req.cache_type in ("classifier", "all"):
        try:
            from .pattern_classifier import _classifier_cache
            count = len(_classifier_cache)
            _classifier_cache.clear()
            cleared["classifier"] = count
        except (ImportError, AttributeError):
            cleared["classifier"] = 0

    return {"cleared": cleared, "cache_type": req.cache_type}
