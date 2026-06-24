import json as _json
import logging
import re

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu, Pt

from .config import get_layout_bank_path
from .data_extractor import extract_chart_data
from .models import Chart, LayoutBank, ProjectState, Slide

CHART_TYPE_MAP = {
    "PIE": XL_CHART_TYPE.PIE,
    "DONUT": XL_CHART_TYPE.DOUGHNUT,
    "BAR": XL_CHART_TYPE.BAR_CLUSTERED,
    "COLUMN": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "BAR_STACKED": XL_CHART_TYPE.BAR_STACKED,
    "COLUMN_STACKED": XL_CHART_TYPE.COLUMN_STACKED,
    "LINE": XL_CHART_TYPE.LINE,
    "AREA": XL_CHART_TYPE.AREA,
    "RADAR": XL_CHART_TYPE.RADAR,
}


def _load_bank() -> LayoutBank:
    p = get_layout_bank_path()
    if not p.exists():
        return LayoutBank()
    try:
        return LayoutBank.model_validate(_json.loads(p.read_text(encoding="utf-8")))
    except Exception:
        return LayoutBank()

PLACEHOLDER_RE = re.compile(r"@(\w+)")


def _slide_has_placeholder(slide, marker: str) -> bool:
    for sh in slide.shapes:
        if sh.has_text_frame and marker in (sh.text_frame.text or ""):
            return True
    return False


def _titulo_y(slide) -> int:
    """Return the Y (EMU) of the textbox containing @Titulo. -1 if not found."""
    for sh in slide.shapes:
        if sh.has_text_frame and "@Titulo" in (sh.text_frame.text or ""):
            return sh.top or 0
    return -1


def _detect_shell_separator_indices(prs) -> tuple[int, int]:
    """Heuristics ordered by reliability:
    1. If only one slide has @Notas → that's shell.
    2. Else: shell = slide whose @Titulo textbox is highest (smallest Y); separator = the other.
    3. Fallback: (0, 1).
    """
    slides_list = list(prs.slides)
    if len(slides_list) < 2:
        return 0, 1
    has_notas_0 = _slide_has_placeholder(slides_list[0], "@Notas")
    has_notas_1 = _slide_has_placeholder(slides_list[1], "@Notas")
    if has_notas_0 and not has_notas_1:
        return 0, 1
    if has_notas_1 and not has_notas_0:
        return 1, 0
    y0 = _titulo_y(slides_list[0])
    y1 = _titulo_y(slides_list[1])
    if y0 >= 0 and y1 >= 0 and y0 != y1:
        return (0, 1) if y0 < y1 else (1, 0)
    return 0, 1


def _duplicate_slide(prs, src_slide):
    """Add a new slide at end of presentation, cloning all shapes AND rels from src_slide.
    Image rels (and other r:embed/r:link references) get NEW rIds in the destination part,
    so we must rewrite the rId attributes inside the copied shape XML to match."""
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    # remove placeholders created by layout
    for sp in list(new_slide.shapes):
        sp._element.getparent().remove(sp._element)

    # Build map old_rId → new_rId by re-creating rels on new slide (skip notesSlide + slideLayout)
    rid_map: dict[str, str] = {}
    for rel in src_slide.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        if rel.is_external:
            new_rid = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rel.rId] = new_rid

    # Copy each shape XML, remapping rIds where they appear as attribute values
    for shape in src_slide.shapes:
        xml_str = etree.tostring(shape._element, encoding="unicode")
        # Sort by length descending to avoid partial substring collisions (rId10 vs rId1)
        for old_rid in sorted(rid_map.keys(), key=len, reverse=True):
            xml_str = xml_str.replace(f'"{old_rid}"', f'"{rid_map[old_rid]}"')
        new_el = etree.fromstring(xml_str)
        new_slide.shapes._spTree.append(new_el)

    return new_slide


def build_pptx(state: ProjectState, out_path: str) -> None:
    """Build final pptx. Open template, detect which slide is shell vs separator by @Notas
    placeholder, duplicate appropriate source per user slide (preserving image rels),
    substitute placeholders, insert charts, then remove the template originals."""
    template_path = state.inputs.template_path
    prs = Presentation(template_path)

    shell_idx, sep_idx = _detect_shell_separator_indices(prs)
    shell_src = prs.slides[shell_idx]
    separator_src = prs.slides[sep_idx]

    # Compute free_area from shell src before any mutation
    from .pptx_template import _compute_free_area
    free_area = _compute_free_area(shell_src, prs.slide_width, prs.slide_height)

    sep_counter = 0
    for slide_def in state.slides:
        if slide_def.type == "separator":
            sep_counter += 1
            new_slide = _duplicate_slide(prs, separator_src)
            _substitute_placeholders(new_slide, {"@Titulo": f"{sep_counter}. {slide_def.title or ''}", "@Notas": ""})
        else:
            new_slide = _duplicate_slide(prs, shell_src)
            notes_text = slide_def.auto_notes or f"Respuesta única. Número de observaciones: {state.parsed_db.sample_size if state.parsed_db else 500}."
            # @Titulo = section title (slide.title from UI). @Subtitulo = question text.
            unique_q_ids = {c.question_id for c in slide_def.charts} if slide_def.charts else set()
            if len(unique_q_ids) == 1 and state.parsed_db:
                qid = next(iter(unique_q_ids))
                q = next((q for q in state.parsed_db.questions if q.id == qid), None)
                subtitle_text = f"{q.code}. {q.text}" if q else ""
            else:
                subtitle_text = ""
            title_text = slide_def.title or ""
            _substitute_placeholders(new_slide, {
                "@Titulo": title_text,
                "@Subtitulo": subtitle_text,
                "@Notas": notes_text,
            })
            _add_slide_content(new_slide, slide_def, state, free_area)

    # Remove the 2 template source slides (they're at positions 0 and 1)
    xml_slides = prs.slides._sldIdLst
    to_remove = [xml_slides[0], xml_slides[1]] if len(xml_slides) >= 2 else list(xml_slides)
    for sld in to_remove:
        rId = sld.rId
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        xml_slides.remove(sld)

    prs.save(out_path)


_log = logging.getLogger(__name__)


def _add_slide_content(slide, slide_def: Slide, state: ProjectState, free_area: dict) -> None:
    """Add charts + analyses to a shell slide using M6.6 classify→render pipeline.

    Pipeline:
      1. Load active style guide (fallback to BUILTIN_STYLE_GUIDE)
      2. Build slide_config from slide_def
      3. Classify → find matching pattern
      4. Build RenderContext
      5. Dispatch to pattern_renderer.render_pattern

    If no pattern matched and BUILTIN_STYLE_GUIDE has patterns, use the first one.
    Falls back to legacy heuristic only when no pattern is available at all.
    """
    from .color_resolver import build_render_context
    from .pattern_classifier import build_slide_config, classify
    from .pattern_renderer import render_pattern
    from .style_guide import BUILTIN_STYLE_GUIDE, load_active

    # 1. Load style guide
    try:
        style_guide = load_active()
    except Exception:
        style_guide = BUILTIN_STYLE_GUIDE

    # 2. Build slide_config (enriches charts with .question and .data)
    db_path = state.inputs.db_path if state.inputs else ""
    slide_config = build_slide_config(slide_def, state.parsed_db, db_path=db_path)

    # 3. Classify — pass real parsed_db dict so extract_context can resolve
    # question_type / n_options. Empty dict here caused everything to be typed
    # as "open", breaking every AI binary/multi pattern trigger.
    matched_pattern = None
    try:
        _sc_dict = {
            "charts": [c.model_dump() for c in slide_def.charts],
            "analyses": [a.model_dump() for a in slide_def.analyses],
        }
        _parsed_db_dict = state.parsed_db.model_dump() if state.parsed_db else {}
        matched_pattern = classify(_sc_dict, _parsed_db_dict, style_guide, require_chart=bool(slide_def.charts))
    except Exception as exc:
        _log.warning("_add_slide_content: classify failed: %s", exc)

    # Fallback strategy:
    #   1. Try classify against BUILTIN_STYLE_GUIDE (hand-curated, covers common cases
    #      like binary + general-only that AI corpus may miss).
    #   2. Else first style_guide pattern that has a chart element.
    #   3. Else first BUILTIN pattern with a chart element.
    if matched_pattern is None:
        try:
            matched_pattern = classify(
                _sc_dict, _parsed_db_dict, BUILTIN_STYLE_GUIDE,
                require_chart=bool(slide_def.charts),
            )
        except Exception as exc:
            _log.warning("_add_slide_content: BUILTIN classify failed: %s", exc)

    if matched_pattern is None:
        candidates = list(style_guide.patterns or []) + list(BUILTIN_STYLE_GUIDE.patterns or [])
        if slide_def.charts:
            matched_pattern = next(
                (p for p in candidates
                 if any(getattr(e, "kind", None) == "chart" for e in (p.implementation.elements or []))),
                None,
            )
        if matched_pattern is None and candidates:
            matched_pattern = candidates[0]

    if matched_pattern is None:
        _log.warning("_add_slide_content: no pattern matched and no fallback — using legacy insertion")
        _add_slide_content_legacy(slide, slide_def, state, free_area)
        return

    # 4. Build RenderContext
    chart_colors_override = state.palette or {}
    ctx = build_render_context(
        style_guide=style_guide,
        slide_config=slide_config,
        chart_colors_override=chart_colors_override,
        free_area=free_area,
    )

    # 5. Render
    try:
        render_pattern(
            pattern=matched_pattern,
            slide=slide,
            ctx=ctx,
            style_guide=style_guide,
            all_patterns=style_guide.patterns,
        )
    except Exception as exc:
        _log.error("_add_slide_content: render_pattern failed: %s", exc, exc_info=True)

    # 6. Render analyses as textboxes (pattern_renderer doesn't handle analyses).
    _add_analyses_textboxes(slide, slide_def, free_area, state.inputs.font_override if state.inputs else None)

    # 7. Render AI-created extras (lines, callouts, subtitle text).
    if slide_def.layout and slide_def.layout.extras:
        _add_layout_extras(slide, slide_def.layout.extras)


def _add_layout_extras(slide, extras) -> None:
    """Render AI-generated extra shapes: separator lines, callout boxes, text."""
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_AUTO_SIZE
    for ex in extras:
        try:
            if ex.kind == "line":
                # Horizontal: cx>0, cy=0 → end at (x+cx, y). Vertical: cx=0, cy>0.
                x1, y1 = ex.x_emu, ex.y_emu
                x2, y2 = ex.x_emu + ex.cx_emu, ex.y_emu + ex.cy_emu
                from pptx.enum.shapes import MSO_CONNECTOR
                line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(x1), Emu(y1), Emu(x2), Emu(y2))
                if ex.color:
                    line.line.color.rgb = RGBColor.from_string(ex.color)
                if ex.style == "dotted":
                    from pptx.oxml.ns import qn
                    ln = line.line._get_or_add_ln()
                    prstDash = ln.makeelement(qn("a:prstDash"), {"val": "sysDot"})
                    for old in ln.findall(qn("a:prstDash")):
                        ln.remove(old)
                    ln.append(prstDash)
                elif ex.style == "dashed":
                    from pptx.oxml.ns import qn
                    ln = line.line._get_or_add_ln()
                    prstDash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
                    for old in ln.findall(qn("a:prstDash")):
                        ln.remove(old)
                    ln.append(prstDash)
        except Exception as exc:
            _log.warning("_add_layout_extras: failed extra %s: %s", ex.kind, exc)


def _add_analyses_textboxes(slide, slide_def: Slide, free_area: dict, font_override: str | None) -> None:
    """Append analysis textboxes. Uses AI layout positions when present,
    else stacks in bottom band of free_area."""
    if not slide_def.analyses:
        return
    ai_positions = {}
    if slide_def.layout and slide_def.layout.positions:
        ai_positions = slide_def.layout.positions

    fa_x = free_area.get("x", 0)
    fa_y = free_area.get("y", 0)
    fa_cx = free_area.get("cx", 1)
    fa_cy = free_area.get("cy", 1)

    # Fallback band for analyses without AI position.
    band_h_frac = 0.20
    band_y = fa_y + int(fa_cy * (1.0 - band_h_frac))
    band_cy = int(fa_cy * band_h_frac)
    band_x = fa_x + int(fa_cx * 0.04)
    band_cx = int(fa_cx * 0.92)

    no_ai_analyses = [a for a in slide_def.analyses if a.id not in ai_positions]
    n_no_ai = len(no_ai_analyses)
    per_cy = band_cy // n_no_ai if n_no_ai else band_cy

    no_ai_index = 0
    for a in slide_def.analyses:
        font_pt = None
        callout = False
        if a.id in ai_positions:
            box = ai_positions[a.id]
            el = {"x": box.x_emu, "y": box.y_emu, "cx": box.cx_emu, "cy": box.cy_emu}
            font_pt = box.font_pt
            callout = box.callout
        else:
            el = {"x": band_x, "y": band_y + no_ai_index * per_cy, "cx": band_cx, "cy": per_cy}
            no_ai_index += 1
        try:
            if callout:
                _add_callout(slide, a.text, el, font_override, font_pt=font_pt)
            else:
                _add_textbox(slide, a.text, el, font_override, font_pt=font_pt)
        except Exception as exc:
            _log.warning("_add_analyses_textboxes: failed analysis %s: %s", a.id, exc)


def _add_callout(slide, text: str, el: dict, font_name: str | None = None, font_pt: float | None = None) -> None:
    """Render analysis as styled callout box: rounded rect, fill D9D9D9, black text."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.dml.color import RGBColor
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Emu(el["x"]), Emu(el["y"]), Emu(el["cx"]), Emu(el["cy"]))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    p = tf.paragraphs[0]
    for run in list(p.runs):
        run.text = ""
    run = p.add_run()
    run.text = text or ""
    if font_name:
        run.font.name = font_name
    if font_pt:
        run.font.size = Pt(font_pt)
    run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)


def _add_slide_content_legacy(slide, slide_def: Slide, state: ProjectState, free_area: dict) -> None:
    """Legacy heuristic chart insertion — fallback when no pattern available."""
    n_chart_an = sum(1 for a in slide_def.analyses if a.scope == "chart")
    n_q_an = sum(1 for a in slide_def.analyses if a.scope == "question")
    has_slide_an = any(a.scope == "slide" for a in slide_def.analyses)

    from .llm_client import _compute_layout_heuristic
    layout_result = _compute_layout_heuristic(
        n_charts=len(slide_def.charts),
        chart_types=[c.chart_type for c in slide_def.charts],
        n_chart_analyses=n_chart_an,
        n_question_analyses=n_q_an,
        has_slide_analysis=has_slide_an,
        free_area=free_area,
    )
    layout = {"elements": layout_result["elements"]}

    for el in layout["elements"]:
        role = el["role"]
        if role.startswith("chart_") and not role.startswith("chart_analysis"):
            i = int(role.split("_")[1])
            if i < len(slide_def.charts):
                _add_chart(slide, slide_def.charts[i], state, el, None)
        elif role.startswith("chart_analysis_"):
            i = int(role.split("_")[2])
            chart_analyses = [a for a in slide_def.analyses if a.scope == "chart"]
            if i < len(chart_analyses):
                _add_textbox(slide, chart_analyses[i].text, el, state.inputs.font_override)
        elif role.startswith("question_analysis_"):
            i = int(role.split("_")[2])
            q_analyses = [a for a in slide_def.analyses if a.scope == "question"]
            if i < len(q_analyses):
                _add_textbox(slide, q_analyses[i].text, el, state.inputs.font_override)
        elif role == "slide_analysis":
            slide_an = next((a for a in slide_def.analyses if a.scope == "slide"), None)
            if slide_an:
                _add_textbox(slide, slide_an.text, el, state.inputs.font_override)


def _add_chart(slide, chart_def: Chart, state: ProjectState, el: dict, specific_style: dict | None = None) -> None:
    primary_bd = chart_def.breakdown_ids[0] if chart_def.breakdown_ids else "general"
    data = extract_chart_data(state.inputs.db_path, _find_question(state, chart_def.question_id),
                              primary_bd, state.parsed_db.data_blocks if state.parsed_db else {})
    cd = CategoryChartData()
    # Categories = options, Series = breakdown categories (or single "Total" if general)
    options = _find_question(state, chart_def.question_id).options
    cd.categories = options

    def _value(cell: dict) -> float:
        """Prefer percentage (×100), fallback to count if pct missing."""
        pct = cell.get("pct")
        if pct is not None:
            try:
                return round(float(pct) * 100, 2)
            except (TypeError, ValueError):
                pass
        try:
            return float(cell.get("count") or 0)
        except (TypeError, ValueError):
            return 0.0

    is_general = (not chart_def.breakdown_ids) or chart_def.breakdown_ids[0].lower() == "general"
    if is_general or len(data) <= 1:
        primary_cat = next((c for c in data if c.lower() in ("total", "general")), None) or (next(iter(data), None) if data else None)
        series = [_value(data.get(primary_cat, {}).get(opt, {})) for opt in options] if primary_cat else [0.0] * len(options)
        cd.add_series("Total", series)
    else:
        cats_to_plot = [c for c in data if c.lower() not in ("total", "general")] or list(data.keys())
        for cat in cats_to_plot:
            values = [_value(data[cat].get(opt, {})) for opt in options]
            cd.add_series(cat, values)

    chart_type_xl = CHART_TYPE_MAP.get(chart_def.chart_type, XL_CHART_TYPE.BAR_CLUSTERED)
    chart_shape = slide.shapes.add_chart(chart_type_xl, Emu(el["x"]), Emu(el["y"]), Emu(el["cx"]), Emu(el["cy"]), cd)
    _apply_training_style(chart_shape.chart, chart_def.chart_type, specific_style)


def _apply_training_style(chart, chart_type: str, specific_style: dict | None = None) -> None:
    """Apply colors / font / data labels / legend from training to a fresh chart.
    If specific_style provided (from a matched layout), use it. Else aggregate from bank."""
    # M6.2: training_extractor removed; chart_style comes from specific_style only.
    # Full color resolution from style_guide planned in M6.4 (color_resolver).
    style = dict(specific_style or {})
    # Auto-defaults so legacy fallback preserves PIE/DONUT labels.
    if chart_type in ("PIE", "DONUT") and not style:
        style = {
            "has_data_labels": True,
            "show_value": False,
            "show_percent": True,
            "show_cat_name": True,
        }
    if not style:
        return

    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_LEGEND_POSITION
    from pptx.util import Pt

    colors = style.get("colors") or []
    font_name = style.get("font")
    font_size = style.get("font_size")
    legend_pos = style.get("legend")
    has_dl = style.get("has_data_labels")

    # Series colors
    if colors:
        try:
            for plot in chart.plots:
                series_list = list(plot.series)
                if chart_type in ("PIE", "DONUT") and series_list:
                    # pie: each data point gets a color
                    ser = series_list[0]
                    pts = list(ser.points)
                    for i, pt in enumerate(pts):
                        c = colors[i % len(colors)]
                        try:
                            fill = pt.format.fill
                            fill.solid()
                            fill.fore_color.rgb = RGBColor.from_string(c)
                        except Exception:
                            pass
                else:
                    for i, ser in enumerate(series_list):
                        c = colors[i % len(colors)]
                        try:
                            fill = ser.format.fill
                            fill.solid()
                            fill.fore_color.rgb = RGBColor.from_string(c)
                        except Exception:
                            pass
        except Exception:
            pass

    # Font + size on chart-level text
    try:
        if font_name:
            chart.font.name = font_name
        if font_size:
            chart.font.size = Pt(font_size)
    except Exception:
        pass

    # Legend
    try:
        if legend_pos:
            chart.has_legend = True
            pos_map = {
                "r": XL_LEGEND_POSITION.RIGHT,
                "right": XL_LEGEND_POSITION.RIGHT,
                "l": XL_LEGEND_POSITION.LEFT,
                "left": XL_LEGEND_POSITION.LEFT,
                "t": XL_LEGEND_POSITION.TOP,
                "top": XL_LEGEND_POSITION.TOP,
                "b": XL_LEGEND_POSITION.BOTTOM,
                "bottom": XL_LEGEND_POSITION.BOTTOM,
                "tr": XL_LEGEND_POSITION.CORNER,
            }
            chart.legend.position = pos_map.get(str(legend_pos).lower(), XL_LEGEND_POSITION.RIGHT)
            chart.legend.include_in_layout = False
    except Exception:
        pass

    # Data labels — use extracted flags so format matches training (e.g., Sí 91.6% vs 458; 92%)
    try:
        if has_dl:
            for plot in chart.plots:
                plot.has_data_labels = True
                dl = plot.data_labels
                dl.show_value = bool(style.get("show_value"))
                dl.show_percentage = bool(style.get("show_percent"))
                dl.show_category_name = bool(style.get("show_cat_name"))
                # if extractor didn't capture any flag, default to percent for pie/donut
                if not any([style.get("show_value"), style.get("show_percent"), style.get("show_cat_name")]):
                    if chart_type in ("PIE", "DONUT"):
                        dl.show_percentage = True
                    else:
                        dl.show_value = True
    except Exception:
        pass


def _add_textbox(slide, text: str, el: dict, font_name: str | None = None, font_pt: int | None = None) -> None:
    tb = slide.shapes.add_textbox(Emu(el["x"]), Emu(el["y"]), Emu(el["cx"]), Emu(el["cy"]))
    tf = tb.text_frame
    tf.word_wrap = True
    # Disable autofit-shrink (AI sets explicit font_pt → don't let pptx scale).
    from pptx.enum.text import MSO_AUTO_SIZE
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    # Build paragraph with explicit run so font props stick (tf.text creates
    # a paragraph but the runs lose set font.size on serialization in some cases).
    p = tf.paragraphs[0]
    # Clear any default empty run.
    for run in list(p.runs):
        run.text = ""
    run = p.add_run()
    run.text = text
    if font_name:
        run.font.name = font_name
    if font_pt:
        run.font.size = Pt(font_pt)


def _substitute_placeholders(slide, mapping: dict[str, str]) -> None:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            full_text = "".join(run.text or "" for run in para.runs)
            for key, val in mapping.items():
                if key in full_text:
                    new_text = full_text.replace(key, val)
                    # rewrite the paragraph as a single run
                    for run in list(para.runs):
                        run.text = ""
                    if para.runs:
                        para.runs[0].text = new_text
                    else:
                        para.add_run().text = new_text


def _find_question(state: ProjectState, qid: str):
    return next(q for q in state.parsed_db.questions if q.id == qid)
