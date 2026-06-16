import re
from pathlib import Path

from pptx import Presentation

from .errors import TemplateInvalidError
from .models import TemplateInfo


PLACEHOLDER_RE = re.compile(r"@\w+")


def load_template(path: str) -> TemplateInfo:
    p = Path(path)
    if not p.exists():
        raise TemplateInvalidError(f"Archivo no encontrado: {path}")
    try:
        prs = Presentation(path)
    except Exception as e:
        raise TemplateInvalidError(f"No se pudo abrir el pptx: {e}") from e

    if len(prs.slides) != 2:
        raise TemplateInvalidError(
            f"Template requiere exactamente 2 slides (shell + separador), tiene {len(prs.slides)}"
        )

    placeholders_by_slide = []
    for idx, slide in enumerate(prs.slides):
        found = set()
        for sh in slide.shapes:
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        for m in PLACEHOLDER_RE.findall(run.text or ""):
                            found.add(m)
                    for m in PLACEHOLDER_RE.findall(para.text or ""):
                        found.add(m)
        if "@Titulo" not in found:
            raise TemplateInvalidError(f"Slide {idx + 1} no tiene placeholder @Titulo")
        placeholders_by_slide.append(found)

    all_placeholders = sorted(placeholders_by_slide[0] | placeholders_by_slide[1])
    free_area = _compute_free_area(prs.slides[0], prs.slide_width, prs.slide_height)
    default_font = _detect_default_font(prs.slides[0])

    return TemplateInfo(
        shell_slide_index=0,
        separator_slide_index=1,
        free_area=free_area,
        placeholders=all_placeholders,
        default_font=default_font,
    )


def _compute_free_area(slide, slide_w_emu: int, slide_h_emu: int) -> dict:
    """Largest contiguous rect not covered by any shape. Approximation: bbox below all shapes' bottom and above their top."""
    if not slide.shapes:
        return {"x": 0, "y": 0, "cx": slide_w_emu, "cy": slide_h_emu}

    shapes_bottom = max((sh.top or 0) + (sh.height or 0) for sh in slide.shapes if sh.top is not None)
    shapes_top = min(sh.top for sh in slide.shapes if sh.top is not None)

    # heuristic: free area is between top-most shape's bottom and bottom-most shape's top
    # if multiple shapes: use the middle gap
    tops_bottoms = []
    for sh in slide.shapes:
        if sh.top is None or sh.height is None:
            continue
        tops_bottoms.append((sh.top, sh.top + sh.height))
    tops_bottoms.sort()

    # find largest vertical gap
    cursor = 0
    best_y = 0
    best_h = 0
    for top, bot in tops_bottoms:
        if top - cursor > best_h:
            best_y = cursor
            best_h = top - cursor
        cursor = max(cursor, bot)
    if slide_h_emu - cursor > best_h:
        best_y = cursor
        best_h = slide_h_emu - cursor

    margin = int(slide_w_emu * 0.03)  # 3% margin
    return {
        "x": margin,
        "y": best_y + margin,
        "cx": slide_w_emu - 2 * margin,
        "cy": max(0, best_h - 2 * margin),
    }


def _detect_default_font(slide) -> str | None:
    for sh in slide.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.font and run.font.name:
                        return run.font.name
    return None
