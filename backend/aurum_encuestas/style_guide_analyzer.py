"""AI Style Guide Analyzer — renders training corpus slides to PNG + calls Claude Sonnet 4.6.

Pipeline:
  1. render_corpus_slides(corpus_dir) → list of slide metadata + PNG bytes (with cache)
  2. build_vision_message(slides) → Anthropic vision content array
  3. analyze_training_corpus(slides) via llm_client → raw JSON string
  4. _validate_and_repair(raw_json, existing_style_guide) → StyleGuide
  5. Save + log
"""
from __future__ import annotations

import base64
import copy
import hashlib
import json
import logging
import shutil
import subprocess
import tempfile
import time
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import Path

from pptx import Presentation

log = logging.getLogger(__name__)

MAX_SLIDES_PER_PPT = 15
MAX_SLIDES_TOTAL = 30


@dataclass
class SlideRenderResult:
    pptx_name: str
    slide_idx: int
    png_bytes: bytes
    metadata: dict  # shape counts, chart types, key text


def _pptx_hash(pptx_path: Path) -> str:
    """Return first 16 hex chars of SHA256 of file content."""
    h = hashlib.sha256(pptx_path.read_bytes()).hexdigest()
    return h[:16]


def _render_slide_to_png(pptx_path: Path, slide_idx: int, output_dir: Path) -> bytes | None:
    """Render a single slide via render_service (libreoffice + pdftoppm pipeline).

    Reuses the production render_service which handles soffice path detection
    + multi-slide via PDF intermediate. Returns PNG bytes or None on failure.
    """
    from .render_service import _PLACEHOLDER_PNG, render_slide_to_png

    try:
        png = render_slide_to_png(str(pptx_path), slide_index=slide_idx)
        # If render_service returns the placeholder, treat as failure (libreoffice missing)
        if png == _PLACEHOLDER_PNG or len(png) < 200:
            log.warning("render returned placeholder for %s slide %d", pptx_path.name, slide_idx)
            return None
        return png
    except Exception as exc:
        log.warning("render_slide_to_png failed for %s slide %d: %s", pptx_path.name, slide_idx, exc)
        return None


def _evict_render_cache_if_needed(cache_dir: Path, max_bytes: int) -> None:
    """LRU eviction: remove oldest files until total size is below max_bytes."""
    files = sorted(cache_dir.glob("*.png"), key=lambda f: f.stat().st_mtime)
    total = sum(f.stat().st_size for f in files)
    while total > max_bytes and files:
        oldest = files.pop(0)
        total -= oldest.stat().st_size
        oldest.unlink()
        log.debug("render_cache evict: removed %s", oldest.name)


def render_corpus_slides(corpus_dir: Path) -> list[SlideRenderResult]:
    """Render slides with charts from all PPTs in corpus_dir.

    Uses disk cache at ~/.aurum/training/render_cache/{pptx_hash}_{slide_idx}.png.
    Samples max 15 slides per PPT and 30 total.
    """
    from .config import RENDER_CACHE_MAX_BYTES, get_render_cache_dir

    cache_dir = get_render_cache_dir()
    results: list[SlideRenderResult] = []
    total_rendered = 0

    pptx_files = sorted(corpus_dir.glob("*.pptx"))
    for pptx_path in pptx_files:
        if total_rendered >= MAX_SLIDES_TOTAL:
            break

        try:
            prs = Presentation(str(pptx_path))
        except Exception as exc:
            log.warning("Could not open %s: %s", pptx_path.name, exc)
            continue

        # Identify slides with at least one chart shape
        slides_with_charts = [
            idx for idx, slide in enumerate(prs.slides)
            if any(getattr(sh, "has_chart", False) for sh in slide.shapes)
        ]

        if not slides_with_charts:
            continue

        # Sample max MAX_SLIDES_PER_PPT uniformly
        if len(slides_with_charts) > MAX_SLIDES_PER_PPT:
            step = len(slides_with_charts) / MAX_SLIDES_PER_PPT
            slides_with_charts = [slides_with_charts[int(i * step)] for i in range(MAX_SLIDES_PER_PPT)]

        pptx_hash = _pptx_hash(pptx_path)

        for slide_idx in slides_with_charts:
            if total_rendered >= MAX_SLIDES_TOTAL:
                break

            cache_key = f"{pptx_hash}_{slide_idx}.png"
            cache_file = cache_dir / cache_key

            if cache_file.exists():
                png_bytes = cache_file.read_bytes()
                log.debug("render_cache HIT: %s", cache_key)
            else:
                log.debug("render_cache MISS: %s — calling libreoffice", cache_key)
                png_bytes = _render_slide_to_png(pptx_path, slide_idx, cache_dir)
                if png_bytes is None:
                    continue
                cache_file.write_bytes(png_bytes)
                _evict_render_cache_if_needed(cache_dir, RENDER_CACHE_MAX_BYTES)

            metadata = _extract_slide_metadata(prs.slides[slide_idx])
            results.append(SlideRenderResult(
                pptx_name=pptx_path.name,
                slide_idx=slide_idx,
                png_bytes=png_bytes,
                metadata=metadata,
            ))
            total_rendered += 1

    return results


def _extract_slide_metadata(slide) -> dict:
    """Extract lightweight XML metadata from a slide for the vision prompt context."""
    shapes_info = []
    for sh in slide.shapes:
        info = {"type": str(sh.shape_type), "has_chart": getattr(sh, "has_chart", False)}
        if sh.has_text_frame:
            text = sh.text_frame.text[:80]
            if text.strip():
                info["text_preview"] = text
        if getattr(sh, "has_chart", False):
            try:
                info["chart_type"] = str(sh.chart.chart_type)
            except Exception:
                pass
        shapes_info.append(info)
    return {
        "shape_count": len(list(slide.shapes)),
        "chart_count": sum(1 for sh in slide.shapes if getattr(sh, "has_chart", False)),
        "shapes": shapes_info[:10],  # cap to keep prompt manageable
    }


# ────────────────────────────────────────────────────────────────────────────
# T2: Build vision user message
# ────────────────────────────────────────────────────────────────────────────

def build_vision_message(slides: list[SlideRenderResult]) -> list[dict]:
    """Build Anthropic vision content array from slide render results.

    Format: [header_text, image_1, slide_1_metadata_text, image_2, ...]
    Capped at MAX_SLIDES_TOTAL images.
    """
    content: list[dict] = []

    # Header text block
    header = (
        f"Analizá estas {min(len(slides), MAX_SLIDES_TOTAL)} slides de entrenamiento del corpus Aurum. "
        "Identificá patterns de presentación, tipos de elementos usados, y sintetizá un style guide JSON "
        "siguiendo EXACTAMENTE el schema especificado en el system prompt."
    )
    content.append({"type": "text", "text": header})

    for slide in slides[:MAX_SLIDES_TOTAL]:
        # PNG image block
        png_b64 = base64.standard_b64encode(slide.png_bytes).decode("ascii")
        content.append({
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": "image/png",
                "data": png_b64,
            },
        })

        # Metadata context text block
        meta_text = (
            f"[{slide.pptx_name} — slide {slide.slide_idx + 1}] "
            f"Formas: {slide.metadata.get('shape_count', '?')}, "
            f"Charts: {slide.metadata.get('chart_count', '?')}"
        )
        shapes_with_text = [s for s in slide.metadata.get("shapes", []) if "text_preview" in s]
        if shapes_with_text:
            texts = [s["text_preview"] for s in shapes_with_text[:3]]
            meta_text += f". Textos: {' | '.join(texts)}"
        content.append({"type": "text", "text": meta_text})

    return content


# ────────────────────────────────────────────────────────────────────────────
# T4: Validation pipeline — JSON parse + pydantic + semantic repair
# ────────────────────────────────────────────────────────────────────────────

from .style_guide import StyleGuide  # noqa: E402


def _validate_and_repair(
    raw_json: str,
    existing_manual_edits: dict[str, str],
) -> tuple[StyleGuide | None, list[str], list[str]]:
    """Parse, pydantic-validate, and semantically repair AI-generated style guide JSON.

    Returns:
        (StyleGuide | None, repairs: list[str], errors: list[str])
    """
    repairs: list[str] = []
    errors: list[str] = []

    # Stage 1: JSON parse
    try:
        data = json.loads(raw_json)
    except json.JSONDecodeError as exc:
        errors.append(f"JSON parse failed: {exc}")
        return None, repairs, errors

    # Stage 2: Semantic pre-repair before pydantic validation
    data = _semantic_repair(data, repairs, errors)

    # Stage 3: Pydantic schema validation
    try:
        sg = StyleGuide.model_validate(data)
    except Exception as exc:
        errors.append(f"Pydantic schema validation failed: {exc}")
        return None, repairs, errors

    # Stage 4: Merge existing manual edits (preserve user overrides)
    if existing_manual_edits:
        # Existing manual edits take precedence (not overwritten by AI re-analysis)
        sg.manual_edits = dict(existing_manual_edits)
        repairs.append(f"Preserved {len(existing_manual_edits)} manual edit(s) from previous analysis")

    return sg, repairs, errors


def _semantic_repair(data: dict, repairs: list[str], errors: list[str]) -> dict:
    """Apply semantic repairs to raw parsed dict before pydantic validation."""
    data = copy.deepcopy(data)

    patterns = data.get("patterns", []) or []

    # Drop duplicate pattern ids (keep first occurrence)
    seen_ids: set[str] = set()
    unique_patterns = []
    for p in patterns:
        pid = p.get("id")
        if pid in seen_ids:
            repairs.append(f"Dropped duplicate pattern id: {pid!r}")
            continue
        if pid:
            seen_ids.add(pid)
        unique_patterns.append(p)
    data["patterns"] = unique_patterns

    # Validate and repair each pattern
    valid_patterns = []
    available_chart_types = set(data.get("available_chart_types", []))
    for p in unique_patterns:
        pid = p.get("id", "<no-id>")
        impl = p.get("implementation") or {}
        elements = impl.get("elements") or []

        repaired_elements = []
        for el in elements:
            # Clamp relative positions to [0, 1]
            position = el.get("position") or {}
            for key in ("x_rel", "y_rel", "w_rel", "h_rel"):
                if key in position:
                    orig = position[key]
                    clamped = max(0.0, min(1.0, float(orig)))
                    if clamped != orig:
                        repairs.append(f"Pattern {pid!r} element {el.get('id')!r}: clamped {key} {orig} → {clamped}")
                        position[key] = clamped
            el["position"] = position

            # Map unsupported chart types to BAR_HORIZONTAL
            if el.get("kind") == "chart":
                ct = el.get("chart_type", "")
                if available_chart_types and ct not in available_chart_types:
                    repairs.append(f"Pattern {pid!r}: chart_type {ct!r} not in available_chart_types → BAR_HORIZONTAL")
                    el["chart_type"] = "BAR_HORIZONTAL"

            repaired_elements.append(el)

        if "implementation" not in p:
            p["implementation"] = {}
        p["implementation"]["elements"] = repaired_elements

        # Validate extends ref (must be present in seen_ids if not null)
        extends = p.get("extends")
        if extends and extends not in seen_ids:
            repairs.append(f"Pattern {pid!r}: extends ref {extends!r} not found — clearing extends")
            p["extends"] = None

        valid_patterns.append(p)

    data["patterns"] = valid_patterns
    return data


# ────────────────────────────────────────────────────────────────────────────
# T5: Full pipeline + async job support
# ────────────────────────────────────────────────────────────────────────────

def run_full_analysis_pipeline(
    progress_dict: dict,
    existing_manual_edits: dict[str, str],
) -> dict:
    """Execute the full AI analysis pipeline synchronously.

    Designed to be called from a background task. Updates progress_dict in-place.

    Returns summary dict: {status, patterns_valid, patterns_dropped, patterns_repaired, ...}
    """
    from .config import get_analysis_logs_dir, get_corpus_dir
    from .llm_client import analyze_training_corpus
    from .style_guide import save_style_guide

    start_time = time.monotonic()

    def _update_progress(pct: int, message: str = "") -> None:
        progress_dict["progress"] = pct
        progress_dict["status"] = "running"
        progress_dict["message"] = message

    _update_progress(5, "Listando corpus...")
    corpus_dir = get_corpus_dir()

    _update_progress(10, "Renderizando slides a PNG...")
    try:
        slides = render_corpus_slides(corpus_dir)
    except Exception as exc:
        log.error("render_corpus_slides failed: %s", exc)
        slides = []

    _update_progress(40, f"Construyendo mensaje vision ({len(slides)} slides)...")
    slides_content = build_vision_message(slides)

    _update_progress(50, "Llamando Claude Sonnet 4.6...")
    raw_json = ""
    llm_result: dict = {}
    all_errors: list[str] = []

    # Retry logic: up to 2 attempts
    for attempt in range(2):
        try:
            llm_result = analyze_training_corpus(slides_content)
            raw_json = llm_result.get("raw_json", "")
            break
        except Exception as exc:
            all_errors.append(f"LLM attempt {attempt + 1} failed: {exc}")
            log.warning("analyze_training_corpus attempt %d failed: %s", attempt + 1, exc)
            if attempt == 1:
                progress_dict["status"] = "error"
                progress_dict["message"] = f"LLM failed after 2 attempts: {all_errors[-1]}"
                return {"status": "error", "errors": all_errors}

    _update_progress(70, "Validando y reparando style guide...")
    sg, repairs, errors = _validate_and_repair(raw_json, existing_manual_edits)
    all_errors.extend(errors)

    if sg is None:
        # Retry once with error feedback
        error_feedback_content = list(slides_content) + [
            {"type": "text", "text": f"Tu respuesta anterior falló validación: {errors}. Corregí y devolvé JSON válido."}
        ]
        try:
            llm_result2 = analyze_training_corpus(error_feedback_content)
            raw_json2 = llm_result2.get("raw_json", "")
            sg, repairs2, errors2 = _validate_and_repair(raw_json2, existing_manual_edits)
            repairs.extend(repairs2)
            all_errors.extend(errors2)
        except Exception as exc:
            all_errors.append(f"Retry failed: {exc}")

    _update_progress(85, "Guardando style guide...")

    patterns_valid = 0
    patterns_dropped = 0

    # ALWAYS save raw response for debug — even if validation passed
    try:
        raw_path = get_corpus_dir().parent / ".last_ai_raw.json"
        raw_path.write_text(raw_json or "(empty)", encoding="utf-8")
    except Exception:
        pass

    if sg is not None:
        # Save validated style guide
        save_style_guide(sg)
        patterns_valid = len(sg.patterns)
        log.info("✓ AI style guide saved: %d patterns valid, %d repairs", patterns_valid, len(repairs))
    else:
        # Fallback to built-in
        from .style_guide import BUILTIN_STYLE_GUIDE
        save_style_guide(BUILTIN_STYLE_GUIDE)
        patterns_valid = len(BUILTIN_STYLE_GUIDE.patterns)
        log.warning("AI analysis failed — falling back to built-in style guide")
        log.warning("Slides rendered: %d", len(slides))
        log.warning("Raw JSON length: %d chars", len(raw_json or ""))
        log.warning("Raw JSON head: %s", (raw_json or "")[:500])
        log.warning("All errors: %s", all_errors)
        log.warning("Raw response saved to %s for inspection", raw_path)

    duration = time.monotonic() - start_time
    corpus_pptxs = [p.name for p in get_corpus_dir().glob("*.pptx")]

    # Save analysis log
    log_entry = {
        "timestamp": datetime.now(UTC).isoformat(),
        "duration_seconds": round(duration, 1),
        "corpus_pptxs": corpus_pptxs,
        "slides_analyzed": len(slides),
        "prompt_version": "v1.0",
        "input_tokens": llm_result.get("input_tokens", 0),
        "output_tokens": llm_result.get("output_tokens", 0),
        "cached_input_tokens": llm_result.get("cached_input_tokens", 0),
        "estimated_cost_usd": llm_result.get("estimated_cost_usd", 0),
        "validation_errors": all_errors,
        "repairs": repairs,
        "patterns_valid": patterns_valid,
        "patterns_dropped": patterns_dropped,
    }
    try:
        logs_dir = get_analysis_logs_dir()
        ts = datetime.now(UTC).strftime("%Y%m%dT%H%M%S")
        (logs_dir / f"{ts}.json").write_text(json.dumps(log_entry, indent=2, ensure_ascii=False), encoding="utf-8")
    except Exception as exc:
        log.warning("Could not save analysis log: %s", exc)

    _update_progress(100, "Análisis completo")
    progress_dict["status"] = "done"

    return {
        "status": "done",
        "patterns_valid": patterns_valid,
        "patterns_dropped": patterns_dropped,
        "repairs": repairs,
        "errors": all_errors,
        "estimated_cost_usd": llm_result.get("estimated_cost_usd", 0),
        "duration_seconds": round(duration, 1),
    }


# ────────────────────────────────────────────────────────────────────────────
# Legacy stub interface (used by other M6.2 stubs)
# ────────────────────────────────────────────────────────────────────────────

def analyze_corpus(corpus_pptx_paths: list[str], existing_style_guide=None) -> StyleGuide:
    """Legacy stub: analyze training corpus with AI vision.

    Returns active StyleGuide.
    """
    from .style_guide import load_active
    return load_active()


def get_render_cache_path(pptx_hash: str, slide_idx: int) -> str:
    """Return the expected PNG cache path for a slide render."""
    from .config import get_render_cache_dir
    return str(get_render_cache_dir() / f"{pptx_hash}_{slide_idx}.png")
