"""M6 Training Sets — corpus CRUD helpers.

Flat corpus (not sets): all PPTs live in ~/.aurum/training/corpus/.
Replaces the training_extractor.py concept from M4/M5.
Full CRUD implemented in M6.8 alongside API endpoints.
"""
from __future__ import annotations

from pathlib import Path


def get_corpus_pptxs() -> list[Path]:
    """Return all PPT files in the corpus directory."""
    from .config import get_corpus_dir
    return sorted(get_corpus_dir().glob("*.pptx"))


def add_pptx_to_corpus(source_path: Path, filename: str) -> Path:
    """Copy a PPT into the corpus directory. Returns destination path."""
    from .config import get_corpus_dir
    dest = get_corpus_dir() / filename
    import shutil
    shutil.copy2(source_path, dest)
    return dest


def delete_pptx_from_corpus(filename: str) -> bool:
    """Delete a PPT from corpus by filename. Returns True if deleted."""
    from .config import get_corpus_dir
    p = get_corpus_dir() / filename
    if p.exists():
        p.unlink()
        return True
    return False


def count_chart_slides(pptx_path: Path) -> int:
    """Count slides that contain at least one chart shape."""
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        return sum(
            1 for slide in prs.slides
            if any(getattr(sh, "has_chart", False) for sh in slide.shapes)
        )
    except Exception:
        return 0
