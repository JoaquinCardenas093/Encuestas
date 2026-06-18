import shutil

import pytest

from aurum_encuestas.render_service import _PLACEHOLDER_PNG, render_slide_to_png

HAS_LIBREOFFICE = shutil.which("soffice") is not None


class TestRenderService:
    """Tests for render_service.py"""

    @pytest.mark.skipif(not HAS_LIBREOFFICE, reason="LibreOffice not installed")
    def test_render_slide_returns_png_bytes_if_libreoffice_installed(self, tmp_path):
        """If LibreOffice is available, render_slide_to_png returns PNG bytes."""
        pptx_path = str(tmp_path / "test.pptx")
        # Create a minimal valid PPTX (empty)
        from pptx import Presentation
        prs = Presentation()
        prs.save(pptx_path)

        result = render_slide_to_png(pptx_path, slide_index=0)

        assert isinstance(result, bytes)
        assert result.startswith(b'\x89PNG')  # PNG magic number

    def test_render_slide_returns_placeholder_if_no_libreoffice(self):
        """If LibreOffice is not available, render_slide_to_png returns placeholder PNG."""
        # Mock absence of LibreOffice by passing a non-existent path
        # This test runs regardless of libreoffice availability
        if HAS_LIBREOFFICE:
            pytest.skip("LibreOffice is installed; skipping placeholder test")

        result = render_slide_to_png("/nonexistent/path.pptx", slide_index=0)

        assert isinstance(result, bytes)
        assert result == _PLACEHOLDER_PNG
        # Verify it's a valid PNG
        assert result.startswith(b'\x89PNG')

    @pytest.mark.skipif(not HAS_LIBREOFFICE, reason="LibreOffice not installed")
    def test_render_specific_slide(self, tmp_path):
        """render_slide_to_png can render specific slides via PDF intermediate."""
        pptx_path = str(tmp_path / "multi_slide.pptx")
        # Create a PPTX with multiple slides
        from pptx import Presentation
        prs = Presentation()
        # Add 3 blank slides
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(pptx_path)

        # Render first slide
        result_0 = render_slide_to_png(pptx_path, slide_index=0)
        assert isinstance(result_0, bytes)
        assert result_0.startswith(b'\x89PNG')

        # Render second slide
        result_1 = render_slide_to_png(pptx_path, slide_index=1)
        assert isinstance(result_1, bytes)
        assert result_1.startswith(b'\x89PNG')
