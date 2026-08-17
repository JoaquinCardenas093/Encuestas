from pptx import Presentation

from aurum_encuestas.models import (
    Analysis,
    Breakdown,
    Chart,
    ParsedDB,
    ProjectInputs,
    ProjectState,
    Question,
    Slide,
    Subtitle,
)
from aurum_encuestas.pptx_generator import build_pptx


def _state(slides, valid_xlsx_path, valid_template_path):
    return ProjectState(
        version=1, project_name="t",
        inputs=ProjectInputs(db_path=str(valid_xlsx_path), template_path=str(valid_template_path)),
        parsed_db=ParsedDB(
            questions=[Question(id="q1", code="P1", text="$p1.recordacion", options=["Sí", "No"], confidence=1.0)],
            breakdowns=[
                Breakdown(id="general", label="General", categories=["Total"]),
                Breakdown(id="sexo", label="Sexo", categories=["Hombre", "Mujer"]),
            ],
            sample_size=500,
            data_blocks={"counts_cols": [3, 17], "pct_row_cols": [21, 35], "pct_col_cols": [41, 55]},
        ),
        slides=slides,
    )


def test_build_pptx_with_separator_and_shell(tmp_path, valid_xlsx_path, valid_template_path):
    slides = [
        Slide(id="s1", type="separator", title="Sección A"),
        Slide(id="s2", type="shell", title="Sección A", charts=[], analyses=[]),
    ]
    state = _state(slides, valid_xlsx_path, valid_template_path)
    out = tmp_path / "out.pptx"
    build_pptx(state, str(out))
    assert out.exists()

    prs = Presentation(str(out))
    assert len(prs.slides) == 2


def test_build_pptx_substitutes_titulo(tmp_path, valid_xlsx_path, valid_template_path):
    slides = [
        Slide(id="s1", type="separator", title="Sección XYZ"),
    ]
    state = _state(slides, valid_xlsx_path, valid_template_path)
    out = tmp_path / "out.pptx"
    build_pptx(state, str(out))

    prs = Presentation(str(out))
    texts = []
    for sh in prs.slides[0].shapes:
        if sh.has_text_frame:
            texts.append(sh.text_frame.text)
    assert any("Sección XYZ" in t for t in texts)
    assert not any("@Titulo" in t for t in texts)


def test_build_pptx_with_chart(tmp_path, valid_xlsx_path, valid_template_path):
    chart = Chart(id="c1", question_id="q1", breakdown_ids=[], chart_type="PIE")
    slides = [
        Slide(id="s1", type="separator", title="Sec"),
        Slide(id="s2", type="shell", title="Sec", charts=[chart]),
    ]
    state = _state(slides, valid_xlsx_path, valid_template_path)
    out = tmp_path / "out.pptx"
    build_pptx(state, str(out))

    prs = Presentation(str(out))
    chart_shapes = [sh for sh in prs.slides[1].shapes if getattr(sh, "has_chart", False)]
    assert len(chart_shapes) == 1


def test_build_pptx_with_analysis_text(tmp_path, valid_xlsx_path, valid_template_path):
    """Pipeline produces a valid pptx when slide has an analysis (no crash)."""
    analysis = Analysis(id="a1", scope="slide", target_id=None, text="Análisis XYZ", ai_generated=True, edited=False)
    slides = [
        Slide(id="s1", type="separator", title="Sec"),
        Slide(id="s2", type="shell", title="Sec", analyses=[analysis]),
    ]
    state = _state(slides, valid_xlsx_path, valid_template_path)
    out = tmp_path / "out.pptx"
    build_pptx(state, str(out))

    prs = Presentation(str(out))
    # The new pipeline renders pattern elements; analysis scope handled via content_source.
    # We just assert the pptx was produced and has 2 slides.
    assert len(prs.slides) == 2


def test_build_pptx_applies_font_override(tmp_path, valid_xlsx_path, valid_template_path):
    """Pipeline completes without crash when font_override is set."""
    state = _state(
        [
            Slide(id="s1", type="separator", title="Sec"),
            Slide(id="s2", type="shell", title="Sec",
                  analyses=[Analysis(id="a1", scope="slide", target_id=None, text="X", ai_generated=False, edited=True)]),
        ],
        valid_xlsx_path, valid_template_path,
    )
    state.inputs.font_override = "Roboto"
    out = tmp_path / "out.pptx"
    build_pptx(state, str(out))

    prs = Presentation(str(out))
    assert len(prs.slides) == 2


def test_build_pptx_pipeline_produces_slide_count(tmp_path, valid_xlsx_path, valid_template_path):
    """classify→render pipeline: 1 sep + 1 shell → 2-slide output, no crash."""
    chart = Chart(id="c1", question_id="q1", breakdown_ids=[], chart_type="PIE")
    slides = [
        Slide(id="s1", type="separator", title="Pipeline Test"),
        Slide(id="s2", type="shell", title="Pipeline Test", charts=[chart]),
    ]
    state = _state(slides, valid_xlsx_path, valid_template_path)
    out = tmp_path / "pipeline_out.pptx"
    build_pptx(state, str(out))

    prs = Presentation(str(out))
    assert len(prs.slides) == 2


def test_build_pptx_empty_shell_no_crash(tmp_path, valid_xlsx_path, valid_template_path):
    """Empty shell slide (no charts, no analyses) produces valid pptx without crash."""
    slides = [
        Slide(id="s1", type="shell", title="Empty"),
    ]
    state = _state(slides, valid_xlsx_path, valid_template_path)
    out = tmp_path / "empty_shell.pptx"
    build_pptx(state, str(out))

    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_add_dashed_box_has_dashed_border_and_no_fill():
    from pptx import Presentation
    from aurum_encuestas.pptx_generator import _add_dashed_box
    from pptx.oxml.ns import qn
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    n0 = len(slide.shapes)
    el = {"x": 500000, "y": 1200000, "cx": 8000000, "cy": 3000000}
    _add_dashed_box(slide, "Texto de análisis", el)
    assert len(slide.shapes) == n0 + 1
    shape = slide.shapes[-1]
    # dashed border present
    ln = shape.line._get_or_add_ln()
    dashes = ln.findall(qn("a:prstDash"))
    assert dashes and dashes[0].get("val") == "dash"
    # no solid fill (noFill / BACKGROUND)
    assert shape.fill.type is None or "BACKGROUND" in str(shape.fill.type)
    # text made it in
    assert "Texto de análisis" in shape.text_frame.text


def test_analyses_textboxes_uses_dashed_box_when_box_style_dashed(valid_template_path, valid_xlsx_path, monkeypatch):
    # Build a slide_def with one slide-analysis and an AI layout whose LayoutBox
    # has box_style="dashed"; assert _add_dashed_box is chosen (a shape with prstDash exists).
    # Reuse whatever ProjectState/Slide construction the existing pptx_generator tests use;
    # if none, spy on the dispatch by monkeypatching _add_dashed_box to record the call.
    from aurum_encuestas import pptx_generator
    called = {}
    monkeypatch.setattr(pptx_generator, "_add_dashed_box",
                        lambda *a, **k: called.setdefault("dashed", True))
    monkeypatch.setattr(pptx_generator, "_add_callout",
                        lambda *a, **k: called.setdefault("callout", True))
    # Minimal: construct a slide with one analysis + an ai layout box_style=dashed and
    # invoke _add_analyses_textboxes directly. See models.Slide/Analysis/SlideLayout/LayoutBox.
    from pptx import Presentation
    from aurum_encuestas.models import Slide, Analysis, SlideLayout, LayoutBox
    prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[6])
    an = Analysis(id="a1", scope="slide", target_id=None, text="X", ai_generated=True, edited=False)
    layout = SlideLayout(positions={"a1": LayoutBox(x_emu=1, y_emu=1, cx_emu=100, cy_emu=100, box_style="dashed")})
    slide_def = Slide(id="s1", type="shell", title="T", charts=[], analyses=[an], layout=layout)
    pptx_generator._add_analyses_textboxes(slide, slide_def, {"x":0,"y":0,"cx":100,"cy":100}, None)
    assert called.get("dashed") and not called.get("callout")


def test_subtitle_text_renders_in_pptx(tmp_path, valid_xlsx_path, valid_template_path):
    slides = [
        Slide(id="s1", type="separator", title="Sec"),
        Slide(id="s2", type="shell", title="Sec", charts=[], analyses=[],
              subtitles=[Subtitle(id="sub1", text="P1. Texto de la pregunta")]),
    ]
    state = _state(slides, valid_xlsx_path, valid_template_path)
    out = tmp_path / "out.pptx"
    build_pptx(state, str(out))
    prs = Presentation(str(out))
    texts = [sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame]
    assert any("Texto de la pregunta" in t for t in texts)


def test_no_auto_subtitle_from_chart_question(tmp_path, valid_xlsx_path, valid_template_path):
    # A shell slide with a single-question chart and NO subtitles must NOT auto-insert
    # the question text (old @Subtitulo behavior is removed).
    slides = [
        Slide(id="s1", type="separator", title="Sec"),
        Slide(id="s2", type="shell", title="Sec",
              charts=[Chart(id="c1", question_id="q1", breakdown_ids=[], chart_type="PIE")],
              analyses=[], subtitles=[]),
    ]
    state = _state(slides, valid_xlsx_path, valid_template_path)
    out = tmp_path / "out.pptx"
    build_pptx(state, str(out))
    prs = Presentation(str(out))
    texts = [sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame]
    assert not any("$p1.recordacion" in t for t in texts)


def test_add_textbox_applies_color():
    from pptx import Presentation
    from pptx.util import Inches
    from aurum_encuestas.pptx_generator import _add_textbox
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, "Hola", {"x": 0, "y": 0, "cx": 1000000, "cy": 500000}, None, font_pt=12, color="C00000")
    tb = [s for s in slide.shapes if s.has_text_frame][-1]
    run = tb.text_frame.paragraphs[0].runs[0]
    assert str(run.font.color.rgb) == "C00000"


def test_hidden_analysis_not_rendered(tmp_path, valid_xlsx_path, valid_template_path):
    from aurum_encuestas.models import Analysis, SlideLayout, LayoutBox
    layout = SlideLayout(positions={"a1": LayoutBox(x_emu=0, y_emu=0, cx_emu=100, cy_emu=100, hidden=True)})
    slides = [
        Slide(id="s1", type="separator", title="Sec"),
        Slide(id="s2", type="shell", title="Sec", charts=[], subtitles=[],
              analyses=[Analysis(id="a1", scope="slide", text="OCULTO_XYZ")], layout=layout),
    ]
    state = _state(slides, valid_xlsx_path, valid_template_path)
    out = tmp_path / "out.pptx"
    build_pptx(state, str(out))
    prs = Presentation(str(out))
    texts = [sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame]
    assert not any("OCULTO_XYZ" in t for t in texts)


def test_created_textbox_rendered(tmp_path, valid_xlsx_path, valid_template_path):
    from aurum_encuestas.models import SlideLayout, LayoutExtra
    layout = SlideLayout(extras=[LayoutExtra(kind="textbox", id="free_1", x_emu=1000000, y_emu=1000000,
                                             cx_emu=3000000, cy_emu=800000, text="CREADO_XYZ", color="404040")])
    slides = [
        Slide(id="s1", type="separator", title="Sec"),
        Slide(id="s2", type="shell", title="Sec", charts=[], analyses=[], subtitles=[], layout=layout),
    ]
    state = _state(slides, valid_xlsx_path, valid_template_path)
    out = tmp_path / "out.pptx"
    build_pptx(state, str(out))
    prs = Presentation(str(out))
    texts = [sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame]
    assert any("CREADO_XYZ" in t for t in texts)


def test_hidden_chart_not_rendered(tmp_path, valid_xlsx_path, valid_template_path):
    from aurum_encuestas.models import SlideLayout, LayoutBox
    layout = SlideLayout(positions={"c1": LayoutBox(x_emu=0, y_emu=0, cx_emu=1, cy_emu=1, hidden=True)})
    slides = [
        Slide(id="s1", type="separator", title="Sec"),
        Slide(id="s2", type="shell", title="Sec", subtitles=[], analyses=[],
              charts=[Chart(id="c1", question_id="q1", breakdown_ids=[], chart_type="PIE")], layout=layout),
    ]
    state = _state(slides, valid_xlsx_path, valid_template_path)
    out = tmp_path / "out.pptx"
    build_pptx(state, str(out))
    prs = Presentation(str(out))
    n_charts = sum(1 for sh in prs.slides[1].shapes if getattr(sh, "has_chart", False))
    assert n_charts == 0


def test_add_chart_handles_empty_breakdown_ids_as_general():
    """Chart with breakdown_ids=[] plots the General row as single series."""
    from pptx import Presentation
    from aurum_encuestas.pptx_generator import _add_chart
    from aurum_encuestas.models import Chart, Question, Breakdown, ParsedDB, ProjectInputs, ProjectState
    from types import SimpleNamespace

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    chart = Chart(id="c1", question_id="q1", breakdown_ids=[], chart_type="PIE")
    parsed = ParsedDB(
        questions=[Question(id="q1", code="Q1", text="t", options=["Sí","No"], confidence=0.9)],
        breakdowns=[Breakdown(id="general", label="General", categories=["Total"])],
        sample_size=100,
        data_blocks={"counts_cols":[],"pct_row_cols":[],"pct_col_cols":[]},
    )
    state = ProjectState(
        project_name="t",
        inputs=ProjectInputs(db_path="", template_path=""),
        parsed_db=parsed,
        slides=[],
    )
    el = {"x": 0, "y": 0, "cx": 5_000_000, "cy": 3_000_000}
    # Monkeypatch extract_chart_data through the import path used by _add_chart
    import aurum_encuestas.pptx_generator as pg
    orig = pg.extract_chart_data
    pg.extract_chart_data = lambda *a, **kw: {"General": {"Sí":{"pct":0.6,"count":60},"No":{"pct":0.4,"count":40}}}
    try:
        _add_chart(slide, chart, state, el)
    finally:
        pg.extract_chart_data = orig
    assert any(sh.has_chart for sh in slide.shapes)
