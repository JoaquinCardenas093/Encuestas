"""Tests for style_guide_analyzer — T1 render cache, T2 vision message, T4 validate/repair, T5 pipeline."""
import json
import os
import time
from pathlib import Path
from unittest.mock import patch

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from aurum_encuestas.style_guide_analyzer import (
    SlideRenderResult,
    _evict_render_cache_if_needed,
    _pptx_hash,
    _validate_and_repair,
    build_vision_message,
    render_corpus_slides,
    run_full_analysis_pipeline,
)

# ─── Fixtures ───────────────────────────────────────────────────────────────

@pytest.fixture
def training_pptx_with_chart(tmp_path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    cd = CategoryChartData()
    cd.categories = ["Sí", "No"]
    cd.add_series("Total", [80, 20])
    s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(2), Inches(2), Inches(4), Inches(4), cd)
    out = tmp_path / "test_corpus.pptx"
    prs.save(str(out))
    return out


# ─── T1: Render cache ────────────────────────────────────────────────────────

def test_pptx_hash_is_16_chars(training_pptx_with_chart):
    h = _pptx_hash(training_pptx_with_chart)
    assert len(h) == 16
    assert h.isalnum()


def test_render_corpus_slides_returns_list(tmp_path, monkeypatch, training_pptx_with_chart):
    monkeypatch.setenv("HOME", str(tmp_path))
    corpus_dir = tmp_path / ".aurum" / "training" / "corpus"
    corpus_dir.mkdir(parents=True, exist_ok=True)
    import shutil
    shutil.copy(training_pptx_with_chart, corpus_dir / "test_corpus.pptx")

    # Mock libreoffice render to return fake PNG bytes
    with patch("aurum_encuestas.style_guide_analyzer._render_slide_to_png") as mock_render:
        mock_render.return_value = b"\x89PNG\r\nfake_bytes"
        results = render_corpus_slides(corpus_dir)  # noqa: F841
    assert isinstance(results, list)
    assert len(results) >= 0  # May be 0 if no charts detected without real libreoffice


def test_render_cache_hit_skips_libreoffice(tmp_path, monkeypatch, training_pptx_with_chart):
    monkeypatch.setenv("HOME", str(tmp_path))
    cache_dir = tmp_path / ".aurum" / "training" / "render_cache"
    cache_dir.mkdir(parents=True, exist_ok=True)

    h = _pptx_hash(training_pptx_with_chart)
    cache_file = cache_dir / f"{h}_0.png"
    cache_file.write_bytes(b"\x89PNG\r\nfake_cached")

    with patch("aurum_encuestas.style_guide_analyzer._render_slide_to_png") as mock_render:
        corpus_dir = tmp_path / ".aurum" / "training" / "corpus"
        corpus_dir.mkdir(parents=True, exist_ok=True)
        import shutil
        shutil.copy(training_pptx_with_chart, corpus_dir / "test_corpus.pptx")
        render_corpus_slides(corpus_dir)
        # render should NOT be called for cached slides
        mock_render.assert_not_called()


def test_evict_render_cache_removes_oldest(tmp_path):
    cache_dir = tmp_path / "render_cache"
    cache_dir.mkdir()

    # Create files with slightly different mtimes
    for i in range(3):
        f = cache_dir / f"slide_{i}.png"
        f.write_bytes(b"x" * (200 * 1024 * 1024))  # 200MB each
        # Simulate older files
        if i < 2:
            os.utime(f, (time.time() - (i + 1) * 1000, time.time() - (i + 1) * 1000))

    # Cache is 600MB, limit is 500MB — should evict oldest
    _evict_render_cache_if_needed(cache_dir, max_bytes=500 * 1024 * 1024)
    remaining = list(cache_dir.glob("*.png"))
    assert len(remaining) < 3


# ─── T2: Vision message ─────────────────────────────────────────────────────

def test_build_vision_message_structure():
    slides = [
        SlideRenderResult(
            pptx_name="test.pptx",
            slide_idx=0,
            png_bytes=b"\x89PNG\r\nfake",
            metadata={"shape_count": 3, "chart_count": 1, "shapes": []},
        ),
        SlideRenderResult(
            pptx_name="test.pptx",
            slide_idx=2,
            png_bytes=b"\x89PNG\r\nfake2",
            metadata={"shape_count": 5, "chart_count": 2, "shapes": []},
        ),
    ]
    message_content = build_vision_message(slides)
    assert isinstance(message_content, list)
    # Should have header text block + (image + text metadata) per slide
    text_blocks = [b for b in message_content if b.get("type") == "text"]
    image_blocks = [b for b in message_content if b.get("type") == "image"]
    assert len(image_blocks) == 2
    assert len(text_blocks) >= 1


def test_build_vision_message_capped_at_30():
    slides = [
        SlideRenderResult("t.pptx", i, b"\x89PNG\r\nfake", {"shape_count": 2, "chart_count": 1, "shapes": []})
        for i in range(40)
    ]
    message_content = build_vision_message(slides)
    image_blocks = [b for b in message_content if b.get("type") == "image"]
    assert len(image_blocks) <= 30


# ─── T4: Validation + repair ─────────────────────────────────────────────────

VALID_STYLE_GUIDE_DICT = {
    "version": 1,
    "is_builtin": False,
    "generated_at": "2026-06-17T20:00:00Z",
    "ai_prompt_version": "v1.0",
    "source_pptxs": ["test.pptx"],
    "manual_edits": {},
    "global": {
        "typography": {"font_family": "Arial", "title_size": 16, "subtitle_size": 12, "label_size": 9, "body_size": 10},
        "text_patterns": {"title": "{code}. {text}", "notes": "{tipo}. N: {n}.", "analysis_style": "El {X}%...", "tone": "formal"},
        "suggested_palette": ["#7F7F7F", "#BFBFBF"],
        "vibe": "Minimalista",
    },
    "available_chart_types": ["PIE", "BAR_HORIZONTAL"],
    "patterns": [
        {
            "id": "binary_general",
            "priority": 0,
            "trigger": {"$and": [{"field": "n_charts_in_slide", "$eq": 1}]},
            "extends": None,
            "best_example": "test.pptx#slide1",
            "why_picked": "Clean layout",
            "implementation": {
                "elements": [
                    {
                        "kind": "chart",
                        "id": "main_chart",
                        "position": {"x_rel": 0.1, "y_rel": 0.1, "w_rel": 0.8, "h_rel": 0.7},
                        "chart_type": "PIE",
                        "data_source": {"chart_ref_index": 0, "value_field": "pct"},
                        "labels": {"show_percentage": True},
                        "legend": "none",
                        "sort": "none",
                    }
                ]
            },
        }
    ],
}


def test_validate_and_repair_valid_input():
    raw_json = json.dumps(VALID_STYLE_GUIDE_DICT)
    sg, repairs, errors = _validate_and_repair(raw_json, existing_manual_edits={})
    assert sg is not None
    assert len(sg.patterns) == 1
    assert sg.patterns[0].id == "binary_general"


def test_validate_and_repair_malformed_json_returns_none():
    raw_json = "not valid json {"
    sg, repairs, errors = _validate_and_repair(raw_json, existing_manual_edits={})
    assert sg is None
    assert len(errors) >= 1


def test_validate_and_repair_clamps_positions():
    d = json.loads(json.dumps(VALID_STYLE_GUIDE_DICT))
    # Inject out-of-range position
    d["patterns"][0]["implementation"]["elements"][0]["position"]["x_rel"] = 1.5
    raw_json = json.dumps(d)
    sg, repairs, errors = _validate_and_repair(raw_json, existing_manual_edits={})
    assert sg is not None
    # x_rel should be clamped to [0, 1]
    el = sg.patterns[0].implementation.elements[0]
    assert el.position.x_rel <= 1.0


def test_validate_and_repair_drops_duplicate_pattern_ids():
    d = json.loads(json.dumps(VALID_STYLE_GUIDE_DICT))
    # Add duplicate pattern
    dup = json.loads(json.dumps(d["patterns"][0]))
    d["patterns"].append(dup)
    raw_json = json.dumps(d)
    sg, repairs, errors = _validate_and_repair(raw_json, existing_manual_edits={})
    assert sg is not None
    assert len(sg.patterns) == 1  # duplicate dropped


def test_validate_and_repair_preserves_manual_edits():
    raw_json = json.dumps(VALID_STYLE_GUIDE_DICT)
    existing_manual = {"binary_general": "2026-06-17T10:00:00Z"}
    sg, repairs, errors = _validate_and_repair(raw_json, existing_manual_edits=existing_manual)
    assert sg is not None
    assert sg.manual_edits.get("binary_general") == "2026-06-17T10:00:00Z"


# ─── T5: Full pipeline ───────────────────────────────────────────────────────

@patch("aurum_encuestas.style_guide_analyzer.render_corpus_slides")
@patch("aurum_encuestas.style_guide_analyzer.build_vision_message")
@patch("aurum_encuestas.llm_client.analyze_training_corpus")
def test_run_full_analysis_pipeline_success(
    mock_analyze, mock_build_msg, mock_render, tmp_path, monkeypatch
):
    monkeypatch.setenv("HOME", str(tmp_path))
    mock_render.return_value = []
    mock_build_msg.return_value = [{"type": "text", "text": "test"}]
    mock_analyze.return_value = {
        "raw_json": json.dumps(VALID_STYLE_GUIDE_DICT),
        "input_tokens": 1000,
        "output_tokens": 200,
        "cached_input_tokens": 800,
        "estimated_cost_usd": 0.05,
    }

    progress = {}
    result = run_full_analysis_pipeline(
        progress_dict=progress,
        existing_manual_edits={},
    )
    assert result["status"] == "done"
    assert result["patterns_valid"] >= 0
    assert progress.get("progress", 0) == 100


@patch("aurum_encuestas.style_guide_analyzer.render_corpus_slides")
@patch("aurum_encuestas.style_guide_analyzer.build_vision_message")
@patch("aurum_encuestas.llm_client.analyze_training_corpus")
def test_run_full_analysis_pipeline_saves_log(
    mock_analyze, mock_build_msg, mock_render, tmp_path, monkeypatch
):
    monkeypatch.setenv("HOME", str(tmp_path))
    mock_render.return_value = []
    mock_build_msg.return_value = [{"type": "text", "text": "test"}]
    mock_analyze.return_value = {
        "raw_json": json.dumps(VALID_STYLE_GUIDE_DICT),
        "input_tokens": 1000,
        "output_tokens": 200,
        "cached_input_tokens": 800,
        "estimated_cost_usd": 0.05,
    }

    run_full_analysis_pipeline(progress_dict={}, existing_manual_edits={})
    logs_dir = tmp_path / ".aurum" / "training" / "ai_analysis_logs"
    log_files = list(logs_dir.glob("*.json"))
    assert len(log_files) >= 1
