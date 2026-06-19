"""Tests for pattern_renderer — M6.6."""
from unittest.mock import MagicMock

from pptx import Presentation

from aurum_encuestas.element_renderers.render_context import RenderContext
from aurum_encuestas.pattern_renderer import merge_implementations, render_pattern

FREE_AREA = {"x": 500000, "y": 1200000, "cx": 11000000, "cy": 5500000}


def _make_ctx():
    ctx = MagicMock(spec=RenderContext)
    ctx.free_area = FREE_AREA
    ctx.chart_colors = ["#7F7F7F", "#BFBFBF"]
    ctx.resolved_colors = {"primary": "#7F7F7F", "secondary": "#BFBFBF", "background": "#FFFFFF"}
    ctx.typography = {"font_family": "Arial", "label_size": 9, "body_size": 10}
    ctx.slide_config = MagicMock()
    ctx.slide_config.charts = []
    ctx.slide_config.analyses = []
    ctx.slide_config.template_shapes = {}
    ctx.resolved_anchors = {}
    return ctx


def _make_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_merge_implementations_child_overrides_parent():
    parent_impl = {
        "elements": [
            {"kind": "chart", "id": "pie", "position": {"x_rel": 0.05, "y_rel": 0.1, "w_rel": 0.4, "h_rel": 0.6}, "chart_type": "PIE", "legend": "none"},
        ]
    }
    child_impl = {
        "elements": [
            {"kind": "chart", "id": "pie", "chart_type": "DONUT"},  # override chart_type
        ]
    }
    merged = merge_implementations(parent_impl, child_impl)
    pie_el = next(e for e in merged["elements"] if e["id"] == "pie")
    assert pie_el["chart_type"] == "DONUT"
    # Parent position preserved
    assert "position" in pie_el
    assert pie_el["position"]["x_rel"] == 0.05


def test_merge_implementations_child_adds_new_element():
    parent_impl = {"elements": [{"kind": "text", "id": "title", "position": {}, "content_source": {"type": "static", "text": "T"}, "style": {}}]}
    child_impl = {"elements": [{"kind": "shape", "id": "divider", "position": {}, "shape_type": "line", "style": {}}]}
    merged = merge_implementations(parent_impl, child_impl)
    ids = [e["id"] for e in merged["elements"]]
    assert "title" in ids
    assert "divider" in ids


def test_render_pattern_dispatches_all_elements():
    slide = _make_slide()
    ctx = _make_ctx()

    pattern = MagicMock()
    pattern.extends = None
    pattern.implementation = MagicMock()
    pattern.implementation.elements = [
        {"kind": "shape", "id": "box", "position": {"x_rel": 0.0, "y_rel": 0.0, "w_rel": 0.5, "h_rel": 0.5}, "shape_type": "rectangle", "style": {"fill": "primary", "color": "primary"}},
        {"kind": "text", "id": "lbl", "position": {"x_rel": 0.0, "y_rel": 0.6, "w_rel": 0.5, "h_rel": 0.2}, "content_source": {"type": "static", "text": "Hello"}, "style": {}},
    ]

    initial_shapes = len(slide.shapes)
    render_pattern(pattern, slide, ctx, style_guide=None, all_patterns=[])
    assert len(slide.shapes) > initial_shapes


def test_render_pattern_topological_order_resolves_anchor():
    """An anchored element must be placed after its anchor is resolved."""
    slide = _make_slide()
    ctx = _make_ctx()

    pattern = MagicMock()
    pattern.extends = None
    pattern.implementation = MagicMock()
    pattern.implementation.elements = [
        # Element B anchors to element A — B must render after A
        {"kind": "text", "id": "B", "position": {"anchor": "A", "relative": "right_of", "offset_rel": 0.01, "w_rel": 0.3, "h_rel": 0.5}, "content_source": {"type": "static", "text": "Right"}, "style": {}},
        {"kind": "shape", "id": "A", "position": {"x_rel": 0.0, "y_rel": 0.0, "w_rel": 0.3, "h_rel": 0.5}, "shape_type": "rectangle", "style": {"fill": "primary", "color": "primary"}},
    ]
    # Should not raise despite B listed before A
    render_pattern(pattern, slide, ctx, style_guide=None, all_patterns=[])


def test_render_pattern_extends_merges_parent_elements():
    slide = _make_slide()
    ctx = _make_ctx()

    parent_pattern = MagicMock()
    parent_pattern.id = "base_chart"
    parent_pattern.extends = None
    parent_pattern.implementation = MagicMock()
    parent_pattern.implementation.elements = [
        {"kind": "shape", "id": "background", "position": {"x_rel": 0.0, "y_rel": 0.0, "w_rel": 1.0, "h_rel": 1.0}, "shape_type": "rectangle", "style": {"fill": "secondary", "color": "secondary"}},
    ]

    child_pattern = MagicMock()
    child_pattern.extends = "base_chart"
    child_pattern.implementation = MagicMock()
    child_pattern.implementation.elements = [
        {"kind": "text", "id": "label", "position": {"x_rel": 0.1, "y_rel": 0.1, "w_rel": 0.8, "h_rel": 0.2}, "content_source": {"type": "static", "text": "Child"}, "style": {}},
    ]

    initial = len(slide.shapes)
    render_pattern(child_pattern, slide, ctx, style_guide=None, all_patterns=[parent_pattern])
    # Both parent shape + child text should be rendered
    assert len(slide.shapes) >= initial + 2


# ── T2: resolve_position tests ──────────────────────────────────────────────

from aurum_encuestas.pattern_renderer import resolve_position  # noqa: E402

FREE_AREA_SIMPLE = {"x": 0, "y": 0, "cx": 10_000_000, "cy": 5_000_000}


def test_resolve_position_relative():
    pos = {"x_rel": 0.1, "y_rel": 0.2, "w_rel": 0.5, "h_rel": 0.4}
    x, y, cx, cy = resolve_position(pos, FREE_AREA_SIMPLE, resolved_anchors={})
    assert x == 1_000_000
    assert y == 1_000_000
    assert cx == 5_000_000
    assert cy == 2_000_000


def test_resolve_position_anchored_right_of():
    resolved_anchors = {"box_a": {"x": 1_000_000, "y": 500_000, "cx": 3_000_000, "cy": 2_000_000}}
    pos = {"anchor": "box_a", "relative": "right_of", "offset_rel": 0.01, "w_rel": 0.3, "h_rel": 0.4}
    x, y, cx, cy = resolve_position(pos, FREE_AREA_SIMPLE, resolved_anchors)
    # x should be box_a.x + box_a.cx + offset
    expected_x = 1_000_000 + 3_000_000 + int(0.01 * 10_000_000)
    assert x == expected_x
    assert y == 500_000  # same y as anchor
    assert cx == int(0.3 * 10_000_000)
    assert cy == int(0.4 * 5_000_000)


def test_resolve_position_anchored_below():
    resolved_anchors = {"box_a": {"x": 0, "y": 0, "cx": 5_000_000, "cy": 2_000_000}}
    pos = {"anchor": "box_a", "relative": "below", "offset_rel": 0.02, "w_rel": 0.5, "h_rel": 0.3}
    x, y, cx, cy = resolve_position(pos, FREE_AREA_SIMPLE, resolved_anchors)
    expected_y = 0 + 2_000_000 + int(0.02 * 5_000_000)  # offset_rel uses cy of free_area
    assert y == expected_y


def test_resolve_position_defaults_for_missing_keys():
    pos = {}  # all defaults
    x, y, cx, cy = resolve_position(pos, FREE_AREA_SIMPLE, resolved_anchors={})
    assert x == 0
    assert y == 0
    assert cx == int(0.5 * 10_000_000)
    assert cy == int(0.5 * 5_000_000)


# ── T4: resolve_data_source tests ────────────────────────────────────────────

from aurum_encuestas.pattern_renderer import resolve_data_source  # noqa: E402


def _make_slide_config_with_charts():
    config = MagicMock()
    config.charts = [
        MagicMock(
            question=MagicMock(options=["Sí", "No"]),
            data={
                "General": {"Sí": {"count": 80, "pct": 0.8}, "No": {"count": 20, "pct": 0.2}},
                "Masculino": {"Sí": {"count": 60, "pct": 0.6}, "No": {"count": 40, "pct": 0.4}},
                "Femenino": {"Sí": {"count": 85, "pct": 0.85}, "No": {"count": 15, "pct": 0.15}},
            },
        )
    ]
    config.analyses = []
    return config


def test_resolve_data_source_chart_ref_index():
    slide_config = _make_slide_config_with_charts()
    ds = {"chart_ref_index": 0, "value_field": "pct"}
    result = resolve_data_source(ds, slide_config)
    assert result is not None
    assert result["chart"] is slide_config.charts[0]
    assert result["value_field"] == "pct"


def test_resolve_data_source_out_of_range_returns_none():
    slide_config = _make_slide_config_with_charts()
    ds = {"chart_ref_index": 99}
    result = resolve_data_source(ds, slide_config)
    assert result is None


def test_resolve_data_source_breakdown_groups_all_except_general():
    slide_config = _make_slide_config_with_charts()
    ds = {"chart_ref_index": 0, "breakdown_groups": "all_except_general"}
    result = resolve_data_source(ds, slide_config)
    assert result is not None
    assert "General" not in result["breakdown_keys"]
    assert "Masculino" in result["breakdown_keys"]


def test_resolve_data_source_breakdown_groups_all():
    slide_config = _make_slide_config_with_charts()
    ds = {"chart_ref_index": 0, "breakdown_groups": "all"}
    result = resolve_data_source(ds, slide_config)
    assert "General" in result["breakdown_keys"]
    assert "Masculino" in result["breakdown_keys"]


def test_resolve_data_source_breakdown_groups_explicit_list():
    slide_config = _make_slide_config_with_charts()
    ds = {"chart_ref_index": 0, "breakdown_groups": ["Masculino"]}
    result = resolve_data_source(ds, slide_config)
    assert result["breakdown_keys"] == ["Masculino"]


# ── Task 4: n_charts_grid fan-out ───────────────────────────────────────────

def test_n_charts_grid_renders_three_chart_shapes():
    from pptx import Presentation
    from aurum_encuestas.pattern_renderer import render_pattern
    from aurum_encuestas.style_guide import BUILTIN_STYLE_GUIDE
    from aurum_encuestas.element_renderers.render_context import RenderContext
    from types import SimpleNamespace

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    q = SimpleNamespace(options=["A", "B", "C"])
    charts = [
        SimpleNamespace(question=q, breakdown_id="general", chart_type="PIE",
                        data={"General": {"A": {"pct": 0.5}, "B": {"pct": 0.3}, "C": {"pct": 0.2}}}, colors=[]),
        SimpleNamespace(question=q, breakdown_id="general", chart_type="BAR_CLUSTERED",
                        data={"General": {"A": {"pct": 0.4}, "B": {"pct": 0.4}, "C": {"pct": 0.2}}}, colors=[]),
        SimpleNamespace(question=q, breakdown_id="general", chart_type="BAR_CLUSTERED",
                        data={"General": {"A": {"pct": 0.6}, "B": {"pct": 0.3}, "C": {"pct": 0.1}}}, colors=[]),
    ]
    slide_config = SimpleNamespace(charts=charts, analyses=[], n_charts=3)
    ctx = RenderContext(slide_config=slide_config, chart_colors=["#7F7F7F"],
                        resolved_colors={"primary": "#7F7F7F", "secondary": "#BFBFBF", "background": "#FFFFFF"},
                        free_area={"x": 0, "y": 0, "cx": 12_192_000, "cy": 6_858_000},
                        typography={"label_size": 9}, resolved_anchors={})

    pattern = next(p for p in BUILTIN_STYLE_GUIDE.patterns if p.id == "n_charts_grid")
    render_pattern(pattern, slide, ctx, BUILTIN_STYLE_GUIDE, list(BUILTIN_STYLE_GUIDE.patterns))

    n_chart_shapes = sum(1 for sh in slide.shapes if sh.has_chart)
    assert n_chart_shapes == 3, f"expected 3 chart shapes, got {n_chart_shapes}"


def test_chart_with_table_type_routes_to_table_renderer():
    """When source_chart.chart_type == TABLE_WITH_MINIBARS, the dispatch hook
    must produce a table shape (has_table True) instead of a chart shape."""
    from pptx import Presentation
    from types import SimpleNamespace
    from aurum_encuestas.pattern_renderer import render_pattern
    from aurum_encuestas.style_guide import BUILTIN_STYLE_GUIDE
    from aurum_encuestas.element_renderers.render_context import RenderContext

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    q = SimpleNamespace(options=["Sí", "No"])
    source_chart = SimpleNamespace(
        question=q,
        breakdown_id="edad",
        chart_type="TABLE_WITH_MINIBARS",
        colors=[],
        data={"General": {"Sí": {"pct": 0.92, "count": 460}, "No": {"pct": 0.08, "count": 40}}},
        all_breakdowns_data={
            "edad": {
                "label": "Rango de edad",
                "categories": {
                    "De 18 a 39 años": {"Sí": {"pct": 0.92, "count": 230}, "No": {"pct": 0.08, "count": 20}},
                    "De 40 a 59 años": {"Sí": {"pct": 0.912, "count": 228}, "No": {"pct": 0.088, "count": 22}},
                },
            },
        },
    )
    slide_config = SimpleNamespace(charts=[source_chart], analyses=[], n_charts=1)
    ctx = RenderContext(
        slide_config=slide_config,
        chart_colors=["#7F7F7F", "#404040", "#EEC245", "#C00000", "#FFC000"],
        resolved_colors={
            "primary": "#7F7F7F", "secondary": "#404040", "background": "#EEC245",
            "accent": "#C00000", "dark": "#FFC000", "light": "#7F7F7F",
        },
        free_area={"x": 487680, "y": 1097280, "cx": 11216640, "cy": 5212080},
        typography={"label_size": 9, "body_size": 10, "title_size": 16, "font_family": "Calibri"},
        style_guide=BUILTIN_STYLE_GUIDE,
        resolved_anchors={},
    )

    # binary_general matches: 1 chart, binary question, no breakdowns
    pattern = next(p for p in BUILTIN_STYLE_GUIDE.patterns if p.id == "binary_general")
    render_pattern(pattern, slide, ctx, BUILTIN_STYLE_GUIDE, list(BUILTIN_STYLE_GUIDE.patterns))

    has_table = any(sh.has_table for sh in slide.shapes)
    has_chart = any(sh.has_chart for sh in slide.shapes)
    assert has_table, f"expected a table shape, got shapes: {[str(sh.shape_type) for sh in slide.shapes]}"
    assert not has_chart, "expected NO chart shape when chart_type is TABLE_WITH_MINIBARS"


# ── Task 5: breakdown_ids list (peek + synthesize) ───────────────────────────

def test_synthesize_table_element_uses_full_breakdown_ids_list():
    """_synthesize_table_element passes ALL real breakdown_ids as
    breakdown_groups, not just the first."""
    from aurum_encuestas.pattern_renderer import _synthesize_table_element
    from types import SimpleNamespace

    chart_el = {
        "kind": "chart",
        "id": "main",
        "position": {"x_rel": 0.1, "y_rel": 0.1, "w_rel": 0.8, "h_rel": 0.8},
        "data_source": {"chart_ref_index": 0, "value_field": "pct"},
    }
    src = SimpleNamespace(breakdown_ids=["edad", "sexo", "general", "nse"], chart_type="TABLE_WITH_MINIBARS")
    el = _synthesize_table_element(chart_el, src)
    assert el["kind"] == "table"
    assert el["structure"] == "segmented_breakdowns"
    assert el["data_source"]["breakdown_groups"] == ["edad", "sexo", "nse"]


def test_dispatch_does_not_fire_when_breakdown_ids_empty():
    """chart_type=TABLE_WITH_MINIBARS with breakdown_ids=[] (general) must
    NOT route to table_renderer — falls through to chart_renderer (which
    will warn about an unmapped chart_type)."""
    from pptx import Presentation
    from types import SimpleNamespace
    from aurum_encuestas.pattern_renderer import render_pattern
    from aurum_encuestas.style_guide import BUILTIN_STYLE_GUIDE
    from aurum_encuestas.element_renderers.render_context import RenderContext

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    q = SimpleNamespace(options=["Sí","No"])
    src = SimpleNamespace(
        question=q, breakdown_ids=[], chart_type="TABLE_WITH_MINIBARS", colors=[],
        data={"General": {"Sí":{"pct":0.6},"No":{"pct":0.4}}}, all_breakdowns_data={},
    )
    slide_config = SimpleNamespace(charts=[src], analyses=[], n_charts=1)
    ctx = RenderContext(
        slide_config=slide_config,
        chart_colors=["#7F7F7F","#404040","#EEC245","#C00000","#FFC000"],
        resolved_colors={"primary":"#7F7F7F","secondary":"#404040","background":"#EEC245"},
        free_area={"x":0,"y":0,"cx":12_192_000,"cy":6_858_000},
        typography={"label_size":9,"body_size":10,"title_size":16,"font_family":"Calibri"},
        style_guide=BUILTIN_STYLE_GUIDE, resolved_anchors={},
    )
    pattern = next(p for p in BUILTIN_STYLE_GUIDE.patterns if p.id == "binary_general")
    render_pattern(pattern, slide, ctx, BUILTIN_STYLE_GUIDE, list(BUILTIN_STYLE_GUIDE.patterns))
    assert not any(sh.has_table for sh in slide.shapes), "should not synthesize table for empty breakdown_ids"


def test_dispatch_fires_for_multi_breakdown_ids():
    """chart_type=TABLE_WITH_MINIBARS with breakdown_ids=['edad','sexo']
    → table shape rendered with both breakdowns in breakdown_groups."""
    from pptx import Presentation
    from types import SimpleNamespace
    from aurum_encuestas.pattern_renderer import render_pattern
    from aurum_encuestas.style_guide import BUILTIN_STYLE_GUIDE
    from aurum_encuestas.element_renderers.render_context import RenderContext

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    q = SimpleNamespace(options=["Sí","No"])
    src = SimpleNamespace(
        question=q, breakdown_ids=["edad","sexo"], chart_type="TABLE_WITH_MINIBARS", colors=[],
        data={"General": {"Sí":{"pct":0.6},"No":{"pct":0.4}}},
        all_breakdowns_data={
            "edad": {"label":"Edad","categories":{
                "18-39":{"Sí":{"pct":0.9,"count":40},"No":{"pct":0.1,"count":5}},
                "40-59":{"Sí":{"pct":0.3,"count":15},"No":{"pct":0.7,"count":35}},
            }},
            "sexo": {"label":"Sexo","categories":{
                "F":{"Sí":{"pct":0.5,"count":30},"No":{"pct":0.5,"count":30}},
                "M":{"Sí":{"pct":0.7,"count":35},"No":{"pct":0.3,"count":15}},
            }},
        },
    )
    slide_config = SimpleNamespace(charts=[src], analyses=[], n_charts=1)
    ctx = RenderContext(
        slide_config=slide_config,
        chart_colors=["#7F7F7F","#404040","#EEC245","#C00000","#FFC000"],
        resolved_colors={"primary":"#7F7F7F","secondary":"#404040","background":"#EEC245"},
        free_area={"x":0,"y":0,"cx":12_192_000,"cy":6_858_000},
        typography={"label_size":9,"body_size":10,"title_size":16,"font_family":"Calibri"},
        style_guide=BUILTIN_STYLE_GUIDE, resolved_anchors={},
    )
    pattern = next(p for p in BUILTIN_STYLE_GUIDE.patterns if p.id == "binary_general")
    render_pattern(pattern, slide, ctx, BUILTIN_STYLE_GUIDE, list(BUILTIN_STYLE_GUIDE.patterns))
    assert any(sh.has_table for sh in slide.shapes), "expected a table for multi-bd TABLE_WITH_MINIBARS"
