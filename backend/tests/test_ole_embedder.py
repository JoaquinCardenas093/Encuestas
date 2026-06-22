from io import BytesIO

from pptx import Presentation
from pptx.util import Inches


def _make_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _png_bytes():
    from PIL import Image
    img = Image.new("RGB", (100, 50), (255, 255, 255))
    buf = BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _xlsx_bytes():
    from openpyxl import Workbook
    wb = Workbook()
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def test_embedded_xlsx_part_added():
    from aurum_encuestas.element_renderers.ole_embedder import embed_ole_xlsx_with_preview
    _prs, slide = _make_slide()
    embed_ole_xlsx_with_preview(
        slide, x=914_400, y=914_400, w=4_572_000, h=2_286_000,
        xlsx_bytes=_xlsx_bytes(), png_bytes=_png_bytes(),
    )
    package = slide.part.package
    partnames = [str(p.partname) for p in package.iter_parts()]
    assert any(
        p.startswith("/ppt/embeddings/Microsoft_Excel_Worksheet") and p.endswith(".xlsx")
        for p in partnames
    )


def test_image_part_added():
    from aurum_encuestas.element_renderers.ole_embedder import embed_ole_xlsx_with_preview
    _prs, slide = _make_slide()
    embed_ole_xlsx_with_preview(
        slide, x=0, y=0, w=4_572_000, h=2_286_000,
        xlsx_bytes=_xlsx_bytes(), png_bytes=_png_bytes(),
    )
    package = slide.part.package
    partnames = [str(p.partname) for p in package.iter_parts()]
    assert any(p.startswith("/ppt/media/image") and p.endswith(".png") for p in partnames)


def test_slide_rels_contain_package_and_image_types():
    from aurum_encuestas.element_renderers.ole_embedder import embed_ole_xlsx_with_preview
    _prs, slide = _make_slide()
    embed_ole_xlsx_with_preview(
        slide, x=0, y=0, w=4_572_000, h=2_286_000,
        xlsx_bytes=_xlsx_bytes(), png_bytes=_png_bytes(),
    )
    types = {rel.reltype for rel in slide.part.rels.values()}
    assert any(t.endswith("/package") for t in types), types
    assert any(t.endswith("/image") for t in types), types


def test_graphic_frame_appended_with_oleObj_and_blip():
    from aurum_encuestas.element_renderers.ole_embedder import embed_ole_xlsx_with_preview
    _prs, slide = _make_slide()
    embed_ole_xlsx_with_preview(
        slide, x=0, y=0, w=4_572_000, h=2_286_000,
        xlsx_bytes=_xlsx_bytes(), png_bytes=_png_bytes(),
    )
    spTree = slide.shapes._spTree
    xml = spTree.xml if hasattr(spTree, "xml") else __import__("lxml.etree", fromlist=["tostring"]).tostring(spTree, encoding="unicode")
    assert "graphicFrame" in xml
    assert "oleObj" in xml
    assert 'progId="Excel.Sheet.12"' in xml
    assert "<a:blip" in xml


def test_multiple_embeds_get_distinct_partnames():
    from aurum_encuestas.element_renderers.ole_embedder import embed_ole_xlsx_with_preview
    _prs, slide = _make_slide()
    embed_ole_xlsx_with_preview(slide, 0, 0, 4_572_000, 2_286_000, _xlsx_bytes(), _png_bytes())
    embed_ole_xlsx_with_preview(slide, 0, 3_000_000, 4_572_000, 2_286_000, _xlsx_bytes(), _png_bytes())
    package = slide.part.package
    xlsx_parts = [
        str(p.partname) for p in package.iter_parts()
        if str(p.partname).startswith("/ppt/embeddings/Microsoft_Excel_Worksheet")
    ]
    png_parts = [str(p.partname) for p in package.iter_parts() if str(p.partname).endswith(".png")]
    assert len(set(xlsx_parts)) == 2
    assert len(set(png_parts)) == 2


def test_round_trip_save_and_reopen_succeeds():
    """Save+reopen round trip — regression guard for partname/blob shape."""
    from aurum_encuestas.element_renderers.ole_embedder import embed_ole_xlsx_with_preview
    prs, slide = _make_slide()
    embed_ole_xlsx_with_preview(
        slide, x=0, y=0, w=4_572_000, h=2_286_000,
        xlsx_bytes=_xlsx_bytes(), png_bytes=_png_bytes(),
    )
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    prs2 = Presentation(buf)
    xlsx_parts = [
        p for p in prs2.part.package.iter_parts()
        if str(p.partname).startswith("/ppt/embeddings/Microsoft_Excel_Worksheet")
    ]
    assert len(xlsx_parts) == 1
    assert xlsx_parts[0].blob[:4] == b"PK\x03\x04", "embedded blob must be raw xlsx PK zip"


def test_graphic_frame_contains_mc_alternate_content():
    from aurum_encuestas.element_renderers.ole_embedder import embed_ole_xlsx_with_preview
    _prs, slide = _make_slide()
    embed_ole_xlsx_with_preview(
        slide, x=0, y=0, w=4_572_000, h=2_286_000,
        xlsx_bytes=_xlsx_bytes(), png_bytes=_png_bytes(),
    )
    from lxml.etree import tostring
    xml = tostring(slide.shapes._spTree, encoding="unicode")
    assert "AlternateContent" in xml
    assert "Choice" in xml
    assert "Fallback" in xml


def test_choice_has_xmlns_v_and_spid():
    from aurum_encuestas.element_renderers.ole_embedder import embed_ole_xlsx_with_preview
    _prs, slide = _make_slide()
    embed_ole_xlsx_with_preview(
        slide, x=0, y=0, w=4_572_000, h=2_286_000,
        xlsx_bytes=_xlsx_bytes(), png_bytes=_png_bytes(),
    )
    from lxml.etree import tostring
    xml = tostring(slide.shapes._spTree, encoding="unicode")
    assert 'xmlns:v="urn:schemas-microsoft-com:vml"' in xml
    assert 'spid="_x0000_s' in xml


def test_fallback_branch_has_no_spid():
    from aurum_encuestas.element_renderers.ole_embedder import embed_ole_xlsx_with_preview
    from lxml import etree
    _prs, slide = _make_slide()
    embed_ole_xlsx_with_preview(
        slide, x=0, y=0, w=4_572_000, h=2_286_000,
        xlsx_bytes=_xlsx_bytes(), png_bytes=_png_bytes(),
    )
    root = etree.fromstring(etree.tostring(slide.shapes._spTree))
    nsmap = {
        "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    fallback = root.xpath(".//mc:Fallback//p:oleObj", namespaces=nsmap)
    assert len(fallback) == 1
    assert fallback[0].get("spid") is None


def test_both_branches_reference_same_xlsx_rid():
    from aurum_encuestas.element_renderers.ole_embedder import embed_ole_xlsx_with_preview
    from lxml import etree
    _prs, slide = _make_slide()
    embed_ole_xlsx_with_preview(
        slide, x=0, y=0, w=4_572_000, h=2_286_000,
        xlsx_bytes=_xlsx_bytes(), png_bytes=_png_bytes(),
    )
    root = etree.fromstring(etree.tostring(slide.shapes._spTree))
    nsmap = {
        "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }
    choice_oleobj = root.xpath(".//mc:Choice/p:oleObj", namespaces=nsmap)
    fallback_oleobj = root.xpath(".//mc:Fallback/p:oleObj", namespaces=nsmap)
    rid_key = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    assert len(choice_oleobj) == 1 and len(fallback_oleobj) == 1
    assert choice_oleobj[0].get(rid_key) == fallback_oleobj[0].get(rid_key)


def test_embedded_part_content_type_is_xlsx():
    from aurum_encuestas.element_renderers.ole_embedder import embed_ole_xlsx_with_preview
    _prs, slide = _make_slide()
    embed_ole_xlsx_with_preview(
        slide, x=0, y=0, w=4_572_000, h=2_286_000,
        xlsx_bytes=_xlsx_bytes(), png_bytes=_png_bytes(),
    )
    package = slide.part.package
    for p in package.iter_parts():
        name = str(p.partname)
        if name.startswith("/ppt/embeddings/Microsoft_Excel_Worksheet") and name.endswith(".xlsx"):
            assert p.content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            return
    raise AssertionError("no xlsx embedding part found")


def test_embedded_blob_is_pk_zip():
    from aurum_encuestas.element_renderers.ole_embedder import embed_ole_xlsx_with_preview
    _prs, slide = _make_slide()
    embed_ole_xlsx_with_preview(
        slide, x=0, y=0, w=4_572_000, h=2_286_000,
        xlsx_bytes=_xlsx_bytes(), png_bytes=_png_bytes(),
    )
    package = slide.part.package
    for p in package.iter_parts():
        name = str(p.partname)
        if name.startswith("/ppt/embeddings/Microsoft_Excel_Worksheet") and name.endswith(".xlsx"):
            assert p.blob[:4] == b"PK\x03\x04"
            return
    raise AssertionError("no xlsx embedding part found")
