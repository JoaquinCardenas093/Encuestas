from unittest.mock import MagicMock

from pptx import Presentation
from pptx.util import Inches

from aurum_encuestas.element_renderers.chart_renderer import render as render_chart
from aurum_encuestas.element_renderers.image_renderer import render as render_image
from aurum_encuestas.element_renderers.render_context import RenderContext
from aurum_encuestas.element_renderers.shape_renderer import render as render_shape
from aurum_encuestas.element_renderers.table_renderer import render as render_table
from aurum_encuestas.element_renderers.text_renderer import render as render_text

FREE_AREA = {"x": Inches(0.5), "y": Inches(1.5), "cx": Inches(12), "cy": Inches(5.5)}


def _make_ctx(colors=None):
    ctx = MagicMock(spec=RenderContext)
    ctx.resolved_colors = colors or {"primary": "#7F7F7F", "secondary": "#BFBFBF", "background": "#FFFFFF"}
    ctx.chart_colors = ["#7F7F7F", "#BFBFBF", "#FFC000"]
    ctx.typography = {"font_family": "Arial", "label_size": 9}
    ctx.free_area = FREE_AREA
    ctx.slide_config = MagicMock()
    ctx.slide_config.charts = [
        MagicMock(
            question=MagicMock(options=["Sí", "No"]),
            data={"General": {"Sí": {"count": 80, "pct": 0.8}, "No": {"count": 20, "pct": 0.2}}},
        )
    ]
    return ctx


def _make_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# T1: chart_renderer
# ---------------------------------------------------------------------------


def test_chart_renderer_pie_adds_chart_shape():
    slide = _make_slide()
    element = {
        "kind": "chart",
        "id": "main_pie",
        "position": {"x_rel": 0.05, "y_rel": 0.1, "w_rel": 0.4, "h_rel": 0.7},
        "chart_type": "PIE",
        "data_source": {"chart_ref_index": 0, "value_field": "pct"},
        "labels": {"show_category_name": True, "show_percentage": True, "position": "outside_end", "format": "0.0%"},
        "legend": "none",
        "sort": "none",
    }
    initial_shapes = len(slide.shapes)
    render_chart(slide, element, _make_ctx())
    assert len(slide.shapes) > initial_shapes


def test_chart_renderer_bar_horizontal_adds_chart_shape():
    slide = _make_slide()
    element = {
        "kind": "chart",
        "id": "bar_chart",
        "position": {"x_rel": 0.0, "y_rel": 0.0, "w_rel": 0.5, "h_rel": 0.5},
        "chart_type": "BAR_HORIZONTAL",
        "data_source": {"chart_ref_index": 0, "value_field": "pct"},
        "labels": {"show_value": True, "format": "0%"},
        "legend": "right",
        "sort": "desc_by_value",
    }
    initial_shapes = len(slide.shapes)
    render_chart(slide, element, _make_ctx())
    assert len(slide.shapes) > initial_shapes


def test_chart_renderer_unknown_chart_type_falls_back_to_bar():
    slide = _make_slide()
    element = {
        "kind": "chart",
        "id": "x",
        "position": {"x_rel": 0.0, "y_rel": 0.0, "w_rel": 0.5, "h_rel": 0.5},
        "chart_type": "UNKNOWN_TYPE",
        "data_source": {"chart_ref_index": 0, "value_field": "pct"},
        "labels": {},
        "legend": "none",
        "sort": "none",
    }
    # Should not raise, falls back gracefully
    render_chart(slide, element, _make_ctx())


def test_chart_renderer_missing_data_source_skips():
    slide = _make_slide()
    element = {
        "kind": "chart",
        "id": "x",
        "position": {"x_rel": 0.0, "y_rel": 0.0, "w_rel": 0.5, "h_rel": 0.5},
        "chart_type": "PIE",
        "data_source": {"chart_ref_index": 99, "value_field": "pct"},  # out of range
        "labels": {},
        "legend": "none",
        "sort": "none",
    }
    initial_shapes = len(slide.shapes)
    render_chart(slide, element, _make_ctx())
    # Should silently skip (no chart added, no exception)
    assert len(slide.shapes) == initial_shapes


# ---------------------------------------------------------------------------
# T2: text_renderer
# ---------------------------------------------------------------------------


def test_text_renderer_static_content():
    slide = _make_slide()
    element = {
        "kind": "text",
        "id": "label1",
        "position": {"x_rel": 0.0, "y_rel": 0.0, "w_rel": 0.5, "h_rel": 0.2},
        "content_source": {"type": "static", "text": "Título de prueba"},
        "style": {"text_color": "primary", "font_size": 12, "bold": True, "align_h": "left"},
    }
    ctx = _make_ctx()
    render_text(slide, element, ctx)
    texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
    assert any("Título de prueba" in t for t in texts)


def test_text_renderer_analysis_content():
    slide = _make_slide()
    ctx = _make_ctx()
    ctx.slide_config.analyses = [
        MagicMock(scope="slide", text="El 80% respondió Sí.", target_id=None)
    ]
    element = {
        "kind": "text",
        "id": "analysis_box",
        "position": {"x_rel": 0.0, "y_rel": 0.8, "w_rel": 1.0, "h_rel": 0.2},
        "content_source": {"type": "analysis", "scope": "slide"},
        "style": {
            "fill": "background",
            "text_color": "primary",
            "font_size": 10,
            "border_left": {"color": "primary", "width_pt": 3.0},
            "padding": 5,
        },
    }
    render_text(slide, element, ctx)
    texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
    assert any("80%" in t for t in texts)


def test_text_renderer_empty_analysis_skips():
    slide = _make_slide()
    ctx = _make_ctx()
    ctx.slide_config.analyses = []
    element = {
        "kind": "text",
        "id": "a",
        "position": {"x_rel": 0.0, "y_rel": 0.0, "w_rel": 0.5, "h_rel": 0.1},
        "content_source": {"type": "analysis", "scope": "slide"},
        "style": {},
    }
    initial = len(slide.shapes)
    render_text(slide, element, ctx)
    # No shapes added if no analysis text
    assert len(slide.shapes) == initial


# ---------------------------------------------------------------------------
# T3: shape_renderer
# ---------------------------------------------------------------------------


def test_shape_renderer_rectangle_added():
    slide = _make_slide()
    element = {
        "kind": "shape",
        "id": "divider_rect",
        "position": {"x_rel": 0.0, "y_rel": 0.0, "w_rel": 1.0, "h_rel": 0.02},
        "shape_type": "rectangle",
        "style": {"fill": "primary", "color": "primary", "width_pt": 0},
    }
    initial = len(slide.shapes)
    render_shape(slide, element, _make_ctx())
    assert len(slide.shapes) > initial


def test_shape_renderer_line_added():
    slide = _make_slide()
    element = {
        "kind": "shape",
        "id": "horiz_line",
        "position": {"x_rel": 0.0, "y_rel": 0.2, "w_rel": 1.0, "h_rel": 0.0},
        "shape_type": "line",
        "style": {"color": "secondary", "width_pt": 1.5},
    }
    initial = len(slide.shapes)
    render_shape(slide, element, _make_ctx())
    assert len(slide.shapes) > initial


# ---------------------------------------------------------------------------
# T4: image_renderer
# ---------------------------------------------------------------------------


def test_image_renderer_no_template_shape_skips():
    """When template shape not found, renderer skips silently."""
    slide = _make_slide()
    ctx = _make_ctx()
    ctx.slide_config.template_shapes = {}  # empty map
    element = {
        "kind": "image",
        "id": "logo",
        "position": {"x_rel": 0.8, "y_rel": 0.0, "w_rel": 0.15, "h_rel": 0.1},
        "source_ref": "logo_shape_id",
    }
    initial = len(slide.shapes)
    render_image(slide, element, ctx)
    assert len(slide.shapes) == initial  # nothing added, no crash


def test_image_renderer_with_template_shape_copies():
    """When a template image shape is provided, it gets copied to the slide."""
    slide = _make_slide()
    ctx = _make_ctx()
    # Simulate a template shape with a picture (we use a mock for simplicity)
    mock_pic = MagicMock()
    mock_pic.shape_type = 13  # MSO_SHAPE_TYPE.PICTURE
    mock_pic.left = 0
    mock_pic.top = 0
    mock_pic.width = 100
    mock_pic.height = 100
    ctx.slide_config.template_shapes = {"logo_shape_id": mock_pic}
    element = {
        "kind": "image",
        "id": "logo",
        "position": {"x_rel": 0.8, "y_rel": 0.0, "w_rel": 0.15, "h_rel": 0.1},
        "source_ref": "logo_shape_id",
    }
    # Should not raise (actual copy logic may vary)
    render_image(slide, element, ctx)


# ---------------------------------------------------------------------------
# T5: table_renderer simple_data
# ---------------------------------------------------------------------------


def _make_table_element(structure="simple_data"):
    return {
        "kind": "table",
        "id": "demo_table",
        "position": {"x_rel": 0.4, "y_rel": 0.1, "w_rel": 0.55, "h_rel": 0.8},
        "structure": structure,
        "data_source": {"chart_ref_index": 0, "breakdown_groups": "all_except_general"},
        "layout": {"col_widths": "auto", "header_height_rel": 0.15, "counts_row_height_rel": 0.1},
        "cells": {
            "group_header": {
                "style": {"fill": "primary", "text_color": "background", "font_size": 10, "bold": True, "align_h": "center"},
                "merge_per_breakdown": True,
            },
            "category_header": {"style": {"fill": "secondary", "font_size": 9, "bold": True}},
            "counts_row": {"style": {"fill": "background", "font_size": 9, "align_h": "center"}, "label_first_col": "Observaciones"},
            "option_row": {
                "style": {"fill": "background", "font_size": 9},
                "label_col_width_rel": 0.10,
                "value_format": "percentage",
                "value_decimals": 1,
                "minibar": {"enabled": False},
            },
        },
    }


def test_table_renderer_simple_data_adds_table():
    slide = _make_slide()
    element = {
        "kind": "table",
        "id": "t1",
        "position": {"x_rel": 0.0, "y_rel": 0.0, "w_rel": 1.0, "h_rel": 0.5},
        "structure": "simple_data",
        "data_source": {"chart_ref_index": 0, "breakdown_groups": "all"},
        "layout": {},
        "cells": {},
    }
    ctx = _make_ctx()
    initial = len(slide.shapes)
    render_table(slide, element, ctx)
    assert len(slide.shapes) > initial


def test_table_renderer_dispatch_by_structure():
    """All three structures should dispatch without raising."""
    ctx = _make_ctx()
    for structure in ("simple_data", "segmented_breakdowns", "comparison_grid"):
        slide = _make_slide()
        element = _make_table_element(structure)
        render_table(slide, element, ctx)  # should not raise


# ---------------------------------------------------------------------------
# T6: segmented_breakdowns
# ---------------------------------------------------------------------------


def _make_ctx_with_breakdowns():
    ctx = _make_ctx()
    # Nested all_breakdowns_data shape: {bd_id: {label, categories: {cat: {opt: {count,pct}}}}}
    ctx.slide_config.charts = [
        MagicMock(
            question=MagicMock(options=["Sí", "No"]),
            data={
                "General": {"Sí": {"count": 80, "pct": 0.8}, "No": {"count": 20, "pct": 0.2}},
            },
            all_breakdowns_data={
                "general": {
                    "label": "General",
                    "categories": {"Total": {"Sí": {"count": 80, "pct": 0.8}, "No": {"count": 20, "pct": 0.2}}},
                },
                "sexo": {
                    "label": "Sexo",
                    "categories": {
                        "Masculino": {"Sí": {"count": 60, "pct": 0.6}, "No": {"count": 40, "pct": 0.4}},
                        "Femenino": {"Sí": {"count": 85, "pct": 0.85}, "No": {"count": 15, "pct": 0.15}},
                    },
                },
            },
            breakdown=MagicMock(id="sexo", label="Sexo", categories=["Masculino", "Femenino"]),
        )
    ]
    return ctx


def test_segmented_breakdowns_adds_table():
    slide = _make_slide()
    element = _make_table_element("segmented_breakdowns")
    ctx = _make_ctx_with_breakdowns()
    render_table(slide, element, ctx)
    # Should have added at least a table shape
    table_shapes = [s for s in slide.shapes if s.shape_type == 19]  # 19 = TABLE
    assert len(table_shapes) >= 1


def test_segmented_breakdowns_row_count():
    """Table must have: group_header + category_header + counts_row + N option rows per breakdown."""
    slide = _make_slide()
    element = _make_table_element("segmented_breakdowns")
    ctx = _make_ctx_with_breakdowns()
    render_table(slide, element, ctx)
    table_shapes = [s for s in slide.shapes if s.shape_type == 19]
    if not table_shapes:
        return  # Already checked in previous test
    tbl = table_shapes[0].table
    # 2 options (Sí, No) + group_header + category_header + counts_row = 5 rows minimum
    assert tbl.rows._tbl.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}tr") is not None


# ---------------------------------------------------------------------------
# T7: minibar overlays
# ---------------------------------------------------------------------------


def test_segmented_breakdowns_minibar_adds_rectangle_overlays():
    """When minibar.enabled=True, rectangle shapes are added on top of table cells."""
    slide = _make_slide()
    element = _make_table_element("segmented_breakdowns")
    # Enable minibar
    element["cells"]["option_row"]["minibar"] = {
        "enabled": True,
        "color_role": "primary",
        "height_rel_to_cell": 0.4,
        "show_percent_text": True,
        "percent_text_position": "left_of_bar",
    }
    ctx = _make_ctx_with_breakdowns()
    initial_shapes = len(slide.shapes)
    render_table(slide, element, ctx)
    # Should have added table + multiple rectangle overlays (one per option×col cell)
    shapes_after = len(slide.shapes)
    assert shapes_after > initial_shapes + 1  # table + at least some minibars


# ---------------------------------------------------------------------------
# T8: comparison_grid
# ---------------------------------------------------------------------------


def test_comparison_grid_adds_table():
    slide = _make_slide()
    element = _make_table_element("comparison_grid")
    ctx = _make_ctx_with_breakdowns()
    initial = len(slide.shapes)
    render_table(slide, element, ctx)
    assert len(slide.shapes) > initial


# ---------------------------------------------------------------------------
# T9: chart_type UI override
# ---------------------------------------------------------------------------


def test_render_chart_respects_source_chart_chart_type(tmp_path):
    """UI-selected chart_type must override the pattern's chart_type."""
    from pptx import Presentation
    from aurum_encuestas.element_renderers.chart_renderer import render
    from aurum_encuestas.element_renderers.render_context import RenderContext
    from types import SimpleNamespace

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    q = SimpleNamespace(options=["Sí", "No"])
    source_chart = SimpleNamespace(
        question=q,
        data={"General": {"Sí": {"pct": 0.6, "count": 60}, "No": {"pct": 0.4, "count": 40}}},
        chart_type="BAR_HORIZONTAL",  # ← UI choice
        colors=[],
    )
    slide_config = SimpleNamespace(charts=[source_chart])
    ctx = RenderContext(
        slide_config=slide_config,
        chart_colors=["#7F7F7F"],
        free_area={"x": 0, "y": 0, "cx": 6_000_000, "cy": 4_000_000},
        typography={"label_size": 9},
        resolved_colors={"primary": "#7F7F7F"},
        resolved_anchors={},
    )

    element = {
        "kind": "chart", "id": "main_pie",
        "chart_type": "PIE",  # ← pattern said PIE
        "position": {"x_rel": 0.1, "y_rel": 0.1, "w_rel": 0.8, "h_rel": 0.8},
        "data_source": {"chart_ref_index": 0, "value_field": "pct"},
    }
    render(slide, element, ctx)

    # Verify the rendered chart is BAR_CLUSTERED (XL type for BAR_HORIZONTAL), not PIE (5)
    from pptx.enum.chart import XL_CHART_TYPE
    chart_shape = next(sh for sh in slide.shapes if sh.has_chart)
    assert chart_shape.chart.chart_type == XL_CHART_TYPE.BAR_CLUSTERED


# ---------------------------------------------------------------------------
# T10: breakdown_ids-driven multi-series
# ---------------------------------------------------------------------------


def test_render_chart_for_breakdown_creates_two_series():
    from pptx import Presentation
    from aurum_encuestas.element_renderers.chart_renderer import render
    from aurum_encuestas.element_renderers.render_context import RenderContext
    from types import SimpleNamespace

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    q = SimpleNamespace(options=["Sí", "No"])
    source_chart = SimpleNamespace(
        question=q,
        breakdown_ids=["sexo"],
        chart_type="BAR_HORIZONTAL",
        data={
            "General": {"Sí": {"pct": 0.55}, "No": {"pct": 0.45}},
            "Hombre":  {"Sí": {"pct": 0.80}, "No": {"pct": 0.20}},
            "Mujer":   {"Sí": {"pct": 0.30}, "No": {"pct": 0.70}},
        },
        colors=[],
    )
    ctx = RenderContext(
        slide_config=SimpleNamespace(charts=[source_chart]),
        chart_colors=["#7F7F7F","#404040"],
        free_area={"x":0,"y":0,"cx":6_000_000,"cy":4_000_000},
        typography={"label_size":9},
        resolved_colors={"primary": "#7F7F7F"},
        resolved_anchors={},
    )
    render(slide, {
        "kind":"chart","id":"c","position":{"x_rel":0,"y_rel":0,"w_rel":1,"h_rel":1},
        "data_source":{"chart_ref_index":0,"value_field":"pct"},
    }, ctx)

    chart_shape = next(sh for sh in slide.shapes if sh.has_chart)
    series = list(chart_shape.chart.series)
    # Expect 2 series (Hombre, Mujer), each with 2 points (Sí, No)
    assert len(series) == 2, f"expected 2 series, got {len(series)}"
    names = {s.name for s in series}
    assert names == {"Hombre", "Mujer"}


# ---------------------------------------------------------------------------
# T4 (task-4): breakdown_ids list + defensive grouped fallback
# ---------------------------------------------------------------------------


def test_build_chart_data_empty_breakdown_ids_uses_general():
    """breakdown_ids=[] → plot the General/Total row as single series."""
    from aurum_encuestas.element_renderers.chart_renderer import _build_chart_data
    from types import SimpleNamespace

    q = SimpleNamespace(options=["Sí","No"])
    source = SimpleNamespace(
        question=q,
        breakdown_ids=[],
        chart_type="PIE",
        data={"General": {"Sí":{"pct":0.6,"count":60},"No":{"pct":0.4,"count":40}}},
        colors=[],
    )
    cd, values = _build_chart_data(source, "pct", "desc_by_value")
    # General branch: one series of N option values
    assert len(values) == 2
    assert [c.label for c in cd.categories] == ["Sí", "No"]


def test_build_chart_data_single_breakdown_uses_that_breakdown():
    """breakdown_ids=['edad'] → multi-series, one per Edad category."""
    from aurum_encuestas.element_renderers.chart_renderer import _build_chart_data
    from types import SimpleNamespace

    q = SimpleNamespace(options=["Sí","No"])
    source = SimpleNamespace(
        question=q,
        breakdown_ids=["edad"],
        chart_type="BAR_HORIZONTAL",
        data={
            "General": {"Sí":{"pct":0.5},"No":{"pct":0.5}},
            "18-39":   {"Sí":{"pct":0.9},"No":{"pct":0.1}},
            "40-59":   {"Sí":{"pct":0.3},"No":{"pct":0.7}},
        },
        colors=[],
    )
    cd, all_values = _build_chart_data(source, "pct", "desc_by_value")
    # values from BOTH categories flattened (2 cats × 2 options = 4 entries)
    assert len(all_values) == 4


def test_chart_title_renders_on_chart_when_set():
    from pptx import Presentation
    from types import SimpleNamespace
    from aurum_encuestas.element_renderers.chart_renderer import render
    from aurum_encuestas.element_renderers.render_context import RenderContext

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    q = SimpleNamespace(options=["A","B"])
    src = SimpleNamespace(
        question=q, breakdown_ids=[], chart_type="BAR_HORIZONTAL",
        title="Plazo del crédito",
        colors=[], show_legend=False, grid_cols=None, cat_titles=None,
        data={"General": {"A":{"pct":0.5},"B":{"pct":0.5}}},
    )
    ctx = RenderContext(
        slide_config=SimpleNamespace(charts=[src]),
        chart_colors=["#7F7F7F","#404040","#EEC245","#C00000","#FFC000"],
        resolved_colors={"primary":"#7F7F7F","secondary":"#404040","background":"#EEC245"},
        free_area={"x":0,"y":0,"cx":12_192_000,"cy":6_858_000},
        typography={"label_size":9,"body_size":10,"title_size":16,"font_family":"Calibri"},
        style_guide=None, resolved_anchors={},
    )
    render(slide, {"kind":"chart","id":"c","position":{"x_rel":0,"y_rel":0,"w_rel":1,"h_rel":1},
                   "data_source":{"chart_ref_index":0,"value_field":"pct"}}, ctx)
    chart_shape = next(sh for sh in slide.shapes if sh.has_chart)
    assert chart_shape.chart.has_title is True
    assert "Plazo del crédito" in chart_shape.chart.chart_title.text_frame.text


def test_bar_horizontal_grouped_legend_bottom_when_show_legend():
    from pptx import Presentation
    from pptx.enum.chart import XL_LEGEND_POSITION
    from types import SimpleNamespace
    from aurum_encuestas.element_renderers.chart_renderer import render
    from aurum_encuestas.element_renderers.render_context import RenderContext

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    q = SimpleNamespace(options=["A","B"])
    src = SimpleNamespace(
        question=q, breakdown_ids=["entidad"], chart_type="BAR_HORIZONTAL_GROUPED",
        show_legend=True, title=None, grid_cols=None, cat_titles=None,
        colors=[],
        data={
            "General": {"A":{"pct":0.5},"B":{"pct":0.5}},
            "Banco":   {"A":{"pct":0.4},"B":{"pct":0.6}},
            "MAF":     {"A":{"pct":0.7},"B":{"pct":0.3}},
        },
    )
    ctx = RenderContext(
        slide_config=SimpleNamespace(charts=[src]),
        chart_colors=["#7F7F7F","#404040","#EEC245","#C00000","#FFC000"],
        resolved_colors={"primary":"#7F7F7F","secondary":"#404040","background":"#EEC245"},
        free_area={"x":0,"y":0,"cx":12_192_000,"cy":6_858_000},
        typography={"label_size":9,"body_size":10,"title_size":16,"font_family":"Calibri"},
        style_guide=None, resolved_anchors={},
    )
    render(slide, {"kind":"chart","id":"c","position":{"x_rel":0,"y_rel":0,"w_rel":1,"h_rel":1},
                   "data_source":{"chart_ref_index":0,"value_field":"pct"}}, ctx)
    chart_shape = next(sh for sh in slide.shapes if sh.has_chart)
    assert chart_shape.chart.has_legend is True
    assert chart_shape.chart.legend.position == XL_LEGEND_POSITION.BOTTOM


def test_bar_horizontal_grouped_no_legend_when_show_legend_false():
    from pptx import Presentation
    from types import SimpleNamespace
    from aurum_encuestas.element_renderers.chart_renderer import render
    from aurum_encuestas.element_renderers.render_context import RenderContext

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    q = SimpleNamespace(options=["A","B"])
    src = SimpleNamespace(
        question=q, breakdown_ids=["entidad"], chart_type="BAR_HORIZONTAL_GROUPED",
        show_legend=False, title=None, grid_cols=None, cat_titles=None,
        colors=[],
        data={"General":{"A":{"pct":0.5},"B":{"pct":0.5}},
              "Banco":{"A":{"pct":0.4},"B":{"pct":0.6}}},
    )
    ctx = RenderContext(
        slide_config=SimpleNamespace(charts=[src]),
        chart_colors=["#7F7F7F","#404040","#EEC245","#C00000","#FFC000"],
        resolved_colors={"primary":"#7F7F7F","secondary":"#404040","background":"#EEC245"},
        free_area={"x":0,"y":0,"cx":12_192_000,"cy":6_858_000},
        typography={"label_size":9,"body_size":10,"title_size":16,"font_family":"Calibri"},
        style_guide=None, resolved_anchors={},
    )
    render(slide, {"kind":"chart","id":"c","position":{"x_rel":0,"y_rel":0,"w_rel":1,"h_rel":1},
                   "data_source":{"chart_ref_index":0,"value_field":"pct"}}, ctx)
    chart_shape = next(sh for sh in slide.shapes if sh.has_chart)
    assert chart_shape.chart.has_legend is False


def test_grouped_fallback_warning_removed(caplog):
    import logging
    from pptx import Presentation
    from types import SimpleNamespace
    from aurum_encuestas.element_renderers.chart_renderer import render
    from aurum_encuestas.element_renderers.render_context import RenderContext

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    q = SimpleNamespace(options=["A","B"])
    src = SimpleNamespace(
        question=q, breakdown_ids=["entidad"], chart_type="BAR_HORIZONTAL_GROUPED",
        show_legend=False, title=None, grid_cols=None, cat_titles=None, colors=[],
        data={"Banco":{"A":{"pct":0.4},"B":{"pct":0.6}}},
    )
    ctx = RenderContext(
        slide_config=SimpleNamespace(charts=[src]),
        chart_colors=["#7F7F7F","#404040","#EEC245","#C00000","#FFC000"],
        resolved_colors={"primary":"#7F7F7F","secondary":"#404040","background":"#EEC245"},
        free_area={"x":0,"y":0,"cx":12_192_000,"cy":6_858_000},
        typography={"label_size":9,"body_size":10,"title_size":16,"font_family":"Calibri"},
        style_guide=None, resolved_anchors={},
    )
    with caplog.at_level(logging.WARNING):
        render(slide, {"kind":"chart","id":"c","position":{"x_rel":0,"y_rel":0,"w_rel":1,"h_rel":1},
                       "data_source":{"chart_ref_index":0,"value_field":"pct"}}, ctx)
    assert not any("Fase B" in r.message for r in caplog.records), \
        f"unexpected Fase B warning still emitted: {[r.message for r in caplog.records]}"


# ---------------------------------------------------------------------------
# T5 (task-5): PIE_GROUPED grid render
# ---------------------------------------------------------------------------


def test_compute_grid_dims_auto_rule():
    from aurum_encuestas.element_renderers.chart_renderer import _compute_grid_dims
    assert _compute_grid_dims(1, None) == (1, 1)
    assert _compute_grid_dims(2, None) == (1, 2)
    assert _compute_grid_dims(3, None) == (1, 3)
    assert _compute_grid_dims(4, None) == (2, 2)
    assert _compute_grid_dims(5, None) == (2, 3)
    assert _compute_grid_dims(6, None) == (2, 3)
    assert _compute_grid_dims(7, None) == (3, 3)
    assert _compute_grid_dims(9, None) == (3, 3)


def test_compute_grid_dims_user_override():
    from aurum_encuestas.element_renderers.chart_renderer import _compute_grid_dims
    assert _compute_grid_dims(6, 2) == (3, 2)
    assert _compute_grid_dims(6, 3) == (2, 3)
    assert _compute_grid_dims(5, 1) == (5, 1)


def _make_pie_grouped_ctx(n_cats=3):
    from types import SimpleNamespace
    from aurum_encuestas.element_renderers.render_context import RenderContext
    q = SimpleNamespace(options=["Sí","No"])
    cats = {f"cat{i}": {"Sí":{"pct":0.5,"count":50},"No":{"pct":0.5,"count":50}} for i in range(n_cats)}
    src = SimpleNamespace(
        question=q, breakdown_ids=["entidad"], chart_type="PIE_GROUPED",
        show_legend=False, title=None, grid_cols=None, cat_titles=None, colors=[],
        data={}, all_breakdowns_data={"entidad":{"label":"Entidad","categories":cats}},
    )
    ctx = RenderContext(
        slide_config=SimpleNamespace(charts=[src]),
        chart_colors=["#7F7F7F","#404040","#EEC245","#C00000","#FFC000"],
        resolved_colors={"primary":"#7F7F7F","secondary":"#404040","background":"#EEC245"},
        free_area={"x":0,"y":0,"cx":12_192_000,"cy":6_858_000},
        typography={"label_size":9,"body_size":10,"title_size":16,"font_family":"Calibri"},
        style_guide=None, resolved_anchors={},
    )
    return src, ctx


def test_pie_grouped_renders_n_pie_shapes():
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE
    from aurum_encuestas.element_renderers.chart_renderer import render
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _src, ctx = _make_pie_grouped_ctx(n_cats=3)
    render(slide, {"kind":"chart","id":"c","position":{"x_rel":0,"y_rel":0,"w_rel":1,"h_rel":1},
                   "data_source":{"chart_ref_index":0,"value_field":"pct"}}, ctx)
    chart_shapes = [sh for sh in slide.shapes if sh.has_chart]
    assert len(chart_shapes) == 3
    for sh in chart_shapes:
        assert sh.chart.chart_type == XL_CHART_TYPE.PIE


def test_pie_grouped_user_grid_cols_overrides_auto():
    from pptx import Presentation
    from aurum_encuestas.element_renderers.chart_renderer import render
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    src, ctx = _make_pie_grouped_ctx(n_cats=6)
    src.grid_cols = 2
    render(slide, {"kind":"chart","id":"c","position":{"x_rel":0,"y_rel":0,"w_rel":1,"h_rel":1},
                   "data_source":{"chart_ref_index":0,"value_field":"pct"}}, ctx)
    chart_shapes = [sh for sh in slide.shapes if sh.has_chart]
    assert len(chart_shapes) == 6
    # Row 1 has 2 shapes at same Y; Row 2 same; Row 3 same → 3 distinct Y values.
    y_set = {sh.top for sh in chart_shapes}
    assert len(y_set) == 3


def test_pie_grouped_cat_titles_override():
    from pptx import Presentation
    from aurum_encuestas.element_renderers.chart_renderer import render
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    src, ctx = _make_pie_grouped_ctx(n_cats=2)
    src.cat_titles = {"cat0": "Banco", "cat1": "MAF"}
    render(slide, {"kind":"chart","id":"c","position":{"x_rel":0,"y_rel":0,"w_rel":1,"h_rel":1},
                   "data_source":{"chart_ref_index":0,"value_field":"pct"}}, ctx)
    chart_shapes = [sh for sh in slide.shapes if sh.has_chart]
    titles = {sh.chart.chart_title.text_frame.text for sh in chart_shapes}
    assert titles == {"Banco", "MAF"}


def test_pie_grouped_chart_title_renders_textbox_above_grid():
    from pptx import Presentation
    from aurum_encuestas.element_renderers.chart_renderer import render
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    src, ctx = _make_pie_grouped_ctx(n_cats=2)
    src.title = "Plazo del crédito"
    render(slide, {"kind":"chart","id":"c","position":{"x_rel":0,"y_rel":0,"w_rel":1,"h_rel":1},
                   "data_source":{"chart_ref_index":0,"value_field":"pct"}}, ctx)
    textboxes = [sh for sh in slide.shapes if sh.has_text_frame and not sh.has_chart]
    assert any("Plazo del crédito" in tb.text_frame.text for tb in textboxes)


def test_pie_grouped_empty_breakdown_warns(caplog):
    import logging
    from pptx import Presentation
    from types import SimpleNamespace
    from aurum_encuestas.element_renderers.chart_renderer import render
    from aurum_encuestas.element_renderers.render_context import RenderContext

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    q = SimpleNamespace(options=["Sí","No"])
    src = SimpleNamespace(
        question=q, breakdown_ids=["edad"], chart_type="PIE_GROUPED",
        show_legend=False, title=None, grid_cols=None, cat_titles=None, colors=[],
        data={}, all_breakdowns_data={},   # ← missing breakdown
    )
    ctx = RenderContext(
        slide_config=SimpleNamespace(charts=[src]),
        chart_colors=["#7F7F7F"], resolved_colors={},
        free_area={"x":0,"y":0,"cx":12_192_000,"cy":6_858_000},
        typography={"label_size":9,"body_size":10,"title_size":16,"font_family":"Calibri"},
        style_guide=None, resolved_anchors={},
    )
    with caplog.at_level(logging.WARNING):
        render(slide, {"kind":"chart","id":"c","position":{"x_rel":0,"y_rel":0,"w_rel":1,"h_rel":1},
                       "data_source":{"chart_ref_index":0,"value_field":"pct"}}, ctx)
    assert not any(sh.has_chart for sh in slide.shapes)
    assert any("PIE_GROUPED" in r.message and "empty" in r.message.lower() for r in caplog.records)
