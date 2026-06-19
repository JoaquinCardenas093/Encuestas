import pytest
from pptx import Presentation
from types import SimpleNamespace

from aurum_encuestas.element_renderers.table_renderer import render
from aurum_encuestas.element_renderers.render_context import RenderContext


@pytest.fixture
def edad_chart():
    q = SimpleNamespace(options=["Sí", "No"])
    return SimpleNamespace(
        question=q,
        breakdown_ids=["edad"],
        chart_type="TABLE_WITH_MINIBARS",
        colors=[],
        data={},
        all_breakdowns_data={
            "edad": {
                "label": "Rango de edad",
                "categories": {
                    "De 18 a 39 años": {"Sí": {"pct": 0.92, "count": 230}, "No": {"pct": 0.08, "count": 20}},
                    "De 40 a 59 años": {"Sí": {"pct": 0.912, "count": 228}, "No": {"pct": 0.088, "count": 22}},
                },
            },
        },
    )


@pytest.fixture
def render_ctx(edad_chart):
    slide_config = SimpleNamespace(charts=[edad_chart], analyses=[], n_charts=1)
    return RenderContext(
        slide_config=slide_config,
        chart_colors=["#7F7F7F", "#404040", "#EEC245", "#C00000", "#FFC000"],
        resolved_colors={
            "primary": "#7F7F7F", "secondary": "#404040", "background": "#EEC245",
            "accent": "#C00000", "dark": "#FFC000", "light": "#7F7F7F",
        },
        free_area={"x": 487680, "y": 1097280, "cx": 11216640, "cy": 5212080},
        typography={"label_size": 9, "body_size": 10, "title_size": 16, "font_family": "Calibri"},
        style_guide=None,
        resolved_anchors={},
    )


def test_segmented_breakdowns_single_panel_image_style(render_ctx):
    """Single breakdown → 1 mini-table with 5 rows × 2 cols (Fase B: no label col):
    group_header / category_header / counts_row / Sí option_row / No option_row."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    element = {
        "kind": "table",
        "id": "edad_table",
        "position": {"x_rel": 0.17, "y_rel": 0.20, "w_rel": 0.65, "h_rel": 0.65},
        "structure": "segmented_breakdowns",
        "data_source": {"chart_ref_index": 0, "breakdown_groups": ["edad"]},
    }
    render(slide, element, render_ctx)

    tables = [sh for sh in slide.shapes if sh.has_table]
    assert len(tables) == 1, f"expected 1 table, got {len(tables)}"
    tbl = tables[0].table
    rows = list(tbl.rows)
    assert len(rows) == 5, f"expected 5 rows (group/cat/counts/Sí/No), got {len(rows)}"
    cols = list(tbl.columns)
    assert len(cols) == 2, f"expected 2 cols (cats only, no label), got {len(cols)}"

    # group_header row should contain "Rango de edad"
    group_texts = [tbl.cell(0, c).text_frame.text for c in range(2)]
    assert any("Rango de edad" in t for t in group_texts), f"group_header missing label: {group_texts}"

    # category_header row should contain both category labels
    cat_texts = [tbl.cell(1, c).text_frame.text for c in range(2)]
    assert "De 18 a 39 años" in cat_texts
    assert "De 40 a 59 años" in cat_texts

    # counts_row: cells [0] and [1] should be 250 (230+20) and 250 (228+22)
    assert tbl.cell(2, 0).text_frame.text.strip() == "250", f"counts col0: {tbl.cell(2,0).text_frame.text!r}"
    assert tbl.cell(2, 1).text_frame.text.strip() == "250", f"counts col1: {tbl.cell(2,1).text_frame.text!r}"

    # option_row "Sí": no label col; values start at col 0
    assert "92.0%" in tbl.cell(3, 0).text_frame.text
    assert "91.2%" in tbl.cell(3, 1).text_frame.text

    # option_row "No"
    assert "8.0%" in tbl.cell(4, 0).text_frame.text
    assert "8.8%" in tbl.cell(4, 1).text_frame.text

    # Style assertions: verify hex-extension + single-panel override fire
    from pptx.dml.color import RGBColor

    # group_header text color should be background role → palette[2] = #EEC245 yellow
    grp_run = tbl.cell(0, 0).text_frame.paragraphs[0].runs[0]
    assert str(grp_run.font.color.rgb) == "EEC245", \
        f"group_header text should be #EEC245 yellow, got {grp_run.font.color.rgb}"
    # group_header text size should be 11 (brief mandate, not the multi-panel 9 cap)
    assert grp_run.font.size is not None and grp_run.font.size.pt == 11, \
        f"group_header font_size should be 11pt, got {grp_run.font.size.pt if grp_run.font.size else None}"

    # option_row Sí cell (row 3, col 0): text_color should be raw hex #FFFFFF (new hex path)
    opt_run = tbl.cell(3, 0).text_frame.paragraphs[0].runs[0]
    assert str(opt_run.font.color.rgb) == "FFFFFF", \
        f"option_row text should be white via hex literal, got {opt_run.font.color.rgb}"


def test_segmented_breakdowns_single_panel_uses_fase_b_cells(edad_chart, render_ctx):
    """Fase B: single-panel uses _SEGMENTED_CELLS_FASE_B (no label col,
    fill=primary, white text on option rows)."""
    from pptx import Presentation
    from aurum_encuestas.element_renderers.table_renderer import render
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render(slide, {
        "kind":"table","id":"t","position":{"x_rel":0.1,"y_rel":0.1,"w_rel":0.8,"h_rel":0.8},
        "structure":"segmented_breakdowns",
        "data_source":{"chart_ref_index":0,"breakdown_groups":["edad"]},
    }, render_ctx)
    tbl = next(sh.table for sh in slide.shapes if sh.has_table)
    cols = list(tbl.columns)
    # No label col: 1 label + 2 cats = 3 cols becomes 2 cols (just cats)
    assert len(cols) == 2, f"expected 2 cols (cats only, no label), got {len(cols)}"
    rows = list(tbl.rows)
    assert len(rows) == 5, f"expected 5 rows, got {len(rows)}"


def test_segmented_breakdowns_multi_panel_uses_fase_b_cells():
    """Multi-panel must apply the same Fase B image-faithful styles
    (post-Fase A multi-panel used pattern's defaults)."""
    from pptx import Presentation
    from types import SimpleNamespace
    from aurum_encuestas.element_renderers.table_renderer import render
    from aurum_encuestas.element_renderers.render_context import RenderContext

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    q = SimpleNamespace(options=["Sí","No"])
    src = SimpleNamespace(
        question=q, breakdown_ids=["edad","sexo"], chart_type="TABLE_WITH_MINIBARS",
        show_legend=False, title=None, grid_cols=None, cat_titles=None, colors=[],
        data={},
        all_breakdowns_data={
            "edad": {"label":"Edad","categories":{
                "18-39":{"Sí":{"pct":0.92,"count":230},"No":{"pct":0.08,"count":20}},
                "40-59":{"Sí":{"pct":0.91,"count":228},"No":{"pct":0.09,"count":22}},
            }},
            "sexo": {"label":"Sexo","categories":{
                "F":{"Sí":{"pct":0.91,"count":250},"No":{"pct":0.09,"count":25}},
                "M":{"Sí":{"pct":0.92,"count":230},"No":{"pct":0.08,"count":20}},
            }},
        },
    )
    ctx = RenderContext(
        slide_config=SimpleNamespace(charts=[src]),
        chart_colors=["#7F7F7F","#404040","#EEC245","#C00000","#FFC000"],
        resolved_colors={"primary":"#7F7F7F","secondary":"#404040","background":"#EEC245"},
        free_area={"x":487680,"y":1097280,"cx":11216640,"cy":5212080},
        typography={"label_size":9,"body_size":10,"title_size":16,"font_family":"Calibri"},
        style_guide=None, resolved_anchors={},
    )
    render(slide, {
        "kind":"table","id":"t","position":{"x_rel":0.02,"y_rel":0.05,"w_rel":0.96,"h_rel":0.90},
        "structure":"segmented_breakdowns",
        "data_source":{"chart_ref_index":0,"breakdown_groups":["edad","sexo"]},
    }, ctx)
    tables = [sh.table for sh in slide.shapes if sh.has_table]
    assert len(tables) == 2, f"expected 2 panels (Edad + Sexo), got {len(tables)}"
    for tbl in tables:
        # 2 cats per panel; no label col → 2 cols
        assert len(list(tbl.columns)) == 2


def test_pack_panels_weight_uses_cats_only_when_label_col_zero():
    from aurum_encuestas.element_renderers.table_renderer import (
        _pack_panels_into_rows, _SEGMENTED_CELLS_FASE_B,
    )
    panels = [
        {"group_id":"edad","label":"Edad","cats":[("18-39",{}),("40-59",{})]},
        {"group_id":"sexo","label":"Sexo","cats":[("F",{}),("M",{})]},
        {"group_id":"nse","label":"NSE","cats":[("A",{}),("B",{}),("C",{}),("D",{}),("E",{})]},
    ]
    label_col_w = _SEGMENTED_CELLS_FASE_B["option_row"]["label_col_width_rel"]
    rows = _pack_panels_into_rows(panels, max_row_weight=12, label_col_width_rel=label_col_w)
    # Total weight without label col = 2+2+5 = 9 ≤ 12 → single row.
    assert len(rows) == 1
    assert len(rows[0]) == 3


def test_segmented_breakdowns_show_legend_renders_external_block(edad_chart, render_ctx):
    """show_legend=True → leftmost shape is the external legend block;
    panel tables shift right."""
    from pptx import Presentation
    from aurum_encuestas.element_renderers.table_renderer import render
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    edad_chart.show_legend = True
    render(slide, {
        "kind":"table","id":"t","position":{"x_rel":0.05,"y_rel":0.1,"w_rel":0.9,"h_rel":0.8},
        "structure":"segmented_breakdowns",
        "data_source":{"chart_ref_index":0,"breakdown_groups":["edad"]},
    }, render_ctx)
    tables = [sh for sh in slide.shapes if sh.has_table]
    assert len(tables) == 2, f"expected legend block + 1 panel = 2 tables, got {len(tables)}"
    # Leftmost table is the legend block
    leftmost = min(tables, key=lambda t: t.left)
    block_tbl = leftmost.table
    # 3 + len(options) rows = 3 + 2 = 5
    assert len(list(block_tbl.rows)) == 5
    # 1 col
    assert len(list(block_tbl.columns)) == 1
    # Row 2 = "Observaciones"
    assert block_tbl.cell(2, 0).text_frame.text.strip() == "Observaciones"
    # Rows 3,4 = options
    assert {block_tbl.cell(3,0).text_frame.text.strip(),
            block_tbl.cell(4,0).text_frame.text.strip()} == {"Sí","No"}


def test_segmented_breakdowns_show_legend_false_no_external_block(edad_chart, render_ctx):
    from pptx import Presentation
    from aurum_encuestas.element_renderers.table_renderer import render
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    edad_chart.show_legend = False
    render(slide, {
        "kind":"table","id":"t","position":{"x_rel":0.05,"y_rel":0.1,"w_rel":0.9,"h_rel":0.8},
        "structure":"segmented_breakdowns",
        "data_source":{"chart_ref_index":0,"breakdown_groups":["edad"]},
    }, render_ctx)
    tables = [sh for sh in slide.shapes if sh.has_table]
    assert len(tables) == 1, f"expected only the panel, got {len(tables)} tables"
