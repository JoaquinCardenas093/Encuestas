"""Test fixtures. Generates valid xlsx and pptx files on the fly per-test using tmp_path."""

import pytest
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt, Emu


@pytest.fixture
def valid_xlsx_path(tmp_path):
    """Synthesize an xlsx matching the BD Aurora schema (2 questions, 5 breakdowns, 3 blocks)."""
    wb = Workbook()
    ws = wb.active
    ws.title = "BD - Análisis"

    # Row 1: breakdown group headers (cols 4-5 edad, 6-7 sexo, 8-12 nse, 13-17 punto)
    ws.cell(1, 4, "Rango de edad")
    ws.cell(1, 6, "Sexo")
    ws.cell(1, 8, "NSE")
    ws.cell(1, 13, "Punto")
    # Repeat at cols 22, 24, 26, 31 (block 2 starts at col 20)
    ws.cell(1, 22, "Rango de edad")
    ws.cell(1, 24, "Sexo")
    ws.cell(1, 26, "NSE")
    ws.cell(1, 31, "Punto")

    # Row 2: subcategories
    ws.cell(2, 3, "General")
    ws.cell(2, 4, "De 18 a 39 años")
    ws.cell(2, 5, "de 40 a 59 años")
    ws.cell(2, 6, "Hombre")
    ws.cell(2, 7, "Mujer")
    ws.cell(2, 8, "Alto")
    ws.cell(2, 9, "Medio")
    ws.cell(2, 10, "Bajo superior")
    ws.cell(2, 11, "Bajo inferior")
    ws.cell(2, 12, "Marginal")
    ws.cell(2, 13, "Paradero")
    ws.cell(2, 14, "Mall")
    ws.cell(2, 15, "CC")
    ws.cell(2, 16, "Plaza")
    ws.cell(2, 17, "Open Plaza")
    ws.cell(2, 21, "General")  # block 2 start

    # Row 3: totals (Total = 500, distribution per breakdown)
    ws.cell(3, 2, "Total")
    ws.cell(3, 3, 500)
    for col, val in enumerate([250, 250, 250, 250, 38, 120, 276, 52, 14, 100, 100, 100, 100, 100], start=4):
        ws.cell(3, col, val)

    # Demographic distribution rows (Sexo)
    ws.cell(4, 1, "Sexo")
    ws.cell(4, 2, "Hombre")
    ws.cell(4, 3, 250)
    ws.cell(5, 2, "Mujer")
    ws.cell(5, 3, 250)

    # Question 1: $p1.label with 2 options Sí/No
    ws.cell(18, 1, "$p1.recordacion")
    ws.cell(18, 2, "Sí")
    ws.cell(18, 3, 458)
    for col, val in enumerate([230, 228, 229, 229, 35, 112, 245, 52, 14, 100, 86, 91, 90, 91], start=4):
        ws.cell(18, col, val)
    ws.cell(19, 2, "No")
    ws.cell(19, 3, 42)

    # Block 2 (percentages) — col 21 repeats General header in row 2
    ws.cell(18, 19, "$p1.recordacion")
    ws.cell(18, 20, "Sí")
    ws.cell(18, 21, 0.916)
    ws.cell(19, 20, "No")
    ws.cell(19, 21, 0.084)

    out = tmp_path / "valid.xlsx"
    wb.save(out)
    return out


@pytest.fixture
def valid_template_path(tmp_path):
    """Synthesize a 2-slide template: shell + separator, both with @Titulo placeholders."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    # Slide 1: SHELL with @Titulo top-left + @Notas bottom-left
    shell = prs.slides.add_slide(blank_layout)
    tb_title = shell.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(6), Inches(0.5))
    tb_title.text_frame.text = "@Titulo"
    tb_notes = shell.shapes.add_textbox(Inches(0.4), Inches(6.7), Inches(8), Inches(0.6))
    tb_notes.text_frame.text = "@Notas"

    # Slide 2: SEPARATOR with @Titulo middle
    sep = prs.slides.add_slide(blank_layout)
    tb_sep = sep.shapes.add_textbox(Inches(0.4), Inches(3.5), Inches(10), Inches(0.6))
    tb_sep.text_frame.text = "Análisis de resultados\n@Titulo"

    out = tmp_path / "valid_template.pptx"
    prs.save(out)
    return out


@pytest.fixture
def invalid_template_one_slide(tmp_path):
    """Template with only 1 slide → should fail validation."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank_layout)
    tb = s.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(6), Inches(0.5))
    tb.text_frame.text = "@Titulo"
    out = tmp_path / "invalid_one_slide.pptx"
    prs.save(out)
    return out


@pytest.fixture
def invalid_template_no_titulo(tmp_path):
    """Template with 2 slides but missing @Titulo placeholder."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for _ in range(2):
        s = prs.slides.add_slide(blank_layout)
        tb = s.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(6), Inches(0.5))
        tb.text_frame.text = "Sin marker"
    out = tmp_path / "invalid_no_titulo.pptx"
    prs.save(out)
    return out
