from types import SimpleNamespace

from pptx import Presentation
from pptx.util import Inches


def _make_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _ctx_with_source(source_chart):
    from aurum_encuestas.element_renderers.render_context import RenderContext
    slide_config = SimpleNamespace(charts=[source_chart], analyses=[], n_charts=1)
    return RenderContext(
        slide_config=slide_config,
        chart_colors=["#7F7F7F", "#404040", "#EEC245", "#C00000", "#FFC000"],
        resolved_colors={"primary": "#7F7F7F", "secondary": "#404040", "background": "#EEC245"},
        free_area={"x": 487680, "y": 1097280, "cx": 11216640, "cy": 5212080},
        typography={"label_size": 9, "body_size": 10, "title_size": 16, "font_family": "Calibri"},
        style_guide=None, resolved_anchors={},
    )


def _make_source():
    q = SimpleNamespace(options=["Sí", "No"])
    return SimpleNamespace(
        question=q,
        breakdown_ids=["edad"],
        all_breakdowns_data={
            "edad": {"label": "Edad", "categories": {
                "18-39": {"Sí": {"pct": 0.92, "count": 230}, "No": {"pct": 0.08, "count": 20}},
                "40-59": {"Sí": {"pct": 0.91, "count": 228}, "No": {"pct": 0.09, "count": 22}},
            }},
        },
    )


def test_render_creates_xlsx_part_image_part_and_graphicFrame():
    from aurum_encuestas.element_renderers.ole_table_renderer import render
    _prs, slide = _make_slide()
    src = _make_source()
    ctx = _ctx_with_source(src)
    element = {
        "kind": "ole_table",
        "id": "t1",
        "position": {"x_rel": 0.1, "y_rel": 0.1, "w_rel": 0.8, "h_rel": 0.6},
        "data_source": {"chart_ref_index": 0, "breakdown_groups": ["edad"]},
    }
    render(slide, element, ctx)
    package = slide.part.package
    partnames = [str(p.partname) for p in package.iter_parts()]
    assert any(p.startswith("/ppt/embeddings/oleObject") and p.endswith(".bin") for p in partnames)
    assert any(p.startswith("/ppt/media/image") and p.endswith(".png") for p in partnames)
    spTree = slide.shapes._spTree
    from lxml.etree import tostring
    xml = tostring(spTree, encoding="unicode")
    assert "graphicFrame" in xml
    assert 'progId="Excel.Sheet.12"' in xml


def test_render_skips_when_chart_ref_index_out_of_range(caplog):
    import logging
    from aurum_encuestas.element_renderers.ole_table_renderer import render
    _prs, slide = _make_slide()
    src = _make_source()
    ctx = _ctx_with_source(src)
    element = {
        "kind": "ole_table",
        "id": "t1",
        "position": {"x_rel": 0.1, "y_rel": 0.1, "w_rel": 0.8, "h_rel": 0.6},
        "data_source": {"chart_ref_index": 99, "breakdown_groups": ["edad"]},
    }
    with caplog.at_level(logging.WARNING):
        render(slide, element, ctx)
    # Returned silently — no shapes added
    assert "oleObj" not in _str_xml(slide)


def _str_xml(slide):
    from lxml.etree import tostring
    return tostring(slide.shapes._spTree, encoding="unicode")
