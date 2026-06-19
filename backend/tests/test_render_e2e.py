# backend/tests/test_render_e2e.py — new
# Task 7: end-to-end render + visual diff regression test.
#
# Adaptations from brief:
#   - `render_project` from `render_service` does not exist; real API is
#     `build_pptx(state, out_path)` from `aurum_encuestas.pptx_generator`.
#   - `build_pptx` requires a valid 2-slide template; we construct one on the
#     fly in tmp_path (same approach as conftest.valid_template_path fixture).
#   - `ProjectInputs.template_path` must point to the synthetic template.
import pytest
from pathlib import Path
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

MAF_XLSX = Path.home() / "Downloads" / "Aurum - Encuestas - Precancelaciones - MAF - Mayo 2026.xlsx"


@pytest.mark.skipif(not MAF_XLSX.exists(), reason="MAF reference xlsx not present")
def test_e2e_three_breakdown_demographics_slide(tmp_path):
    """Render a demographics slide with 3 separate single-series charts
    and assert geometry + chart_type fidelity."""
    from pptx.util import Inches
    from aurum_encuestas.xlsx_parser import parse_xlsx
    from aurum_encuestas.pptx_generator import build_pptx
    from aurum_encuestas.models import ProjectState, ProjectInputs, Slide, Chart

    # --- Build a minimal 2-slide template (shell + separator) in tmp_path ---
    tpl_prs = Presentation()
    tpl_prs.slide_width = Inches(13.33)
    tpl_prs.slide_height = Inches(7.5)
    blank_layout = tpl_prs.slide_layouts[6]

    # Shell slide: needs @Titulo + @Notas markers
    shell = tpl_prs.slides.add_slide(blank_layout)
    tb_title = shell.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(6), Inches(0.5))
    tb_title.text_frame.text = "@Titulo"
    tb_notes = shell.shapes.add_textbox(Inches(0.4), Inches(6.7), Inches(8), Inches(0.6))
    tb_notes.text_frame.text = "@Notas"

    # Separator slide: @Titulo only
    sep = tpl_prs.slides.add_slide(blank_layout)
    tb_sep = sep.shapes.add_textbox(Inches(0.4), Inches(3.5), Inches(10), Inches(0.6))
    tb_sep.text_frame.text = "Análisis de resultados\n@Titulo"

    template_path = str(tmp_path / "template.pptx")
    tpl_prs.save(template_path)

    # --- Parse the real xlsx ---
    parsed = parse_xlsx(str(MAF_XLSX))

    # --- Select first question + first 3 non-general breakdowns ---
    q_id = parsed.questions[0].id
    bds = [b.id for b in parsed.breakdowns if b.id != "general"][:3]
    if len(bds) < 3:
        pytest.skip("MAF xlsx has <3 non-general breakdowns")

    charts = [
        Chart(id=f"c{i}", question_id=q_id, breakdown_ids=[b], chart_type="BAR_HORIZONTAL", colors=[])
        for i, b in enumerate(bds)
    ]
    inputs = ProjectInputs(db_path=str(MAF_XLSX), template_path=template_path)
    state = ProjectState(
        project_name="e2e",
        inputs=inputs,
        parsed_db=parsed,
        slides=[Slide(id="s1", type="shell", title="Demo", charts=charts, analyses=[])],
    )

    # --- Render ---
    out = tmp_path / "out.pptx"
    build_pptx(state, str(out))

    # --- Assertions ---
    prs = Presentation(str(out))
    chart_slides = [s for s in prs.slides if any(sh.has_chart for sh in s.shapes)]
    assert chart_slides, "no chart slides rendered"

    chart_shapes = [sh for sh in chart_slides[0].shapes if sh.has_chart]
    assert len(chart_shapes) == 3, f"expected 3 chart shapes, got {len(chart_shapes)}"

    for sh in chart_shapes:
        assert sh.chart.chart_type == XL_CHART_TYPE.BAR_CLUSTERED, (
            f"expected BAR_CLUSTERED, got {sh.chart.chart_type}"
        )
        # Height must be at least 50% of slide height (Task 3: retuned layouts)
        assert sh.height >= 0.50 * prs.slide_height, (
            f"chart height {sh.height} < 50% of slide height {prs.slide_height}"
        )


def test_e2e_table_with_minibars_renders_single_panel_table(tmp_path, monkeypatch):
    """End-to-end: a Chart with chart_type=TABLE_WITH_MINIBARS and breakdown_ids=["edad"]
    flows through build_pptx and produces a table shape on the rendered slide."""
    from pptx import Presentation
    from pptx.util import Inches
    from aurum_encuestas.pptx_generator import build_pptx
    from aurum_encuestas.models import (
        ProjectState, ProjectInputs, Slide, Chart, ParsedDB, Question, Breakdown,
    )

    # Synthetic template with shell slide + separator slide (build_pptx requires both)
    tpl = tmp_path / "tpl.pptx"
    p = Presentation()
    p.slide_width = Inches(13.33)
    p.slide_height = Inches(7.5)
    blank = p.slide_layouts[6]
    # Shell slide: @Titulo + @Notas markers
    shell = p.slides.add_slide(blank)
    shell.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6)).text_frame.text = "@Titulo"
    shell.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5)).text_frame.text = "@Notas"
    # Separator slide: @Titulo only
    sep = p.slides.add_slide(blank)
    sep.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(0.6)).text_frame.text = "@Titulo"
    p.save(str(tpl))

    parsed = ParsedDB(
        questions=[Question(id="q1", code="Q1", text="¿Cliente actual?", options=["Sí", "No"], confidence=0.95)],
        breakdowns=[
            Breakdown(id="general", label="General", categories=["Total"]),
            Breakdown(id="edad", label="Rango de edad", categories=["De 18 a 39 años", "De 40 a 59 años"]),
        ],
        sample_size=500,
        data_blocks={"counts_cols": [], "pct_row_cols": [], "pct_col_cols": []},
    )

    def fake_extract(*args, **kwargs):
        return {"General": {"Sí": {"pct": 0.92, "count": 460}, "No": {"pct": 0.08, "count": 40}}}

    def fake_extract_all(*args, **kwargs):
        return {"edad": {"label": "Rango de edad", "categories": {
            "De 18 a 39 años": {"Sí": {"pct": 0.92, "count": 230}, "No": {"pct": 0.08, "count": 20}},
            "De 40 a 59 años": {"Sí": {"pct": 0.912, "count": 228}, "No": {"pct": 0.088, "count": 22}},
        }}}

    monkeypatch.setattr("aurum_encuestas.data_extractor.extract_chart_data", fake_extract)
    monkeypatch.setattr("aurum_encuestas.data_extractor.extract_all_breakdowns_data", fake_extract_all)

    state = ProjectState(
        project_name="e2e-table",
        inputs=ProjectInputs(db_path="/dev/null", template_path=str(tpl)),
        parsed_db=parsed,
        slides=[
            Slide(
                id="s1", type="shell", title="Demografía",
                charts=[Chart(id="c1", question_id="q1", breakdown_ids=["edad"],
                              chart_type="TABLE_WITH_MINIBARS", colors=[])],
                analyses=[],
            ),
        ],
    )

    out = tmp_path / "out.pptx"
    build_pptx(state, str(out))

    prs = Presentation(str(out))
    slides_with_tables = [s for s in prs.slides if any(sh.has_table for sh in s.shapes)]
    assert slides_with_tables, "expected at least one table shape rendered"
    slides_with_charts = [s for s in prs.slides if any(sh.has_chart for sh in s.shapes)]
    assert not slides_with_charts, "expected NO chart shape (TABLE_WITH_MINIBARS routes to table)"
